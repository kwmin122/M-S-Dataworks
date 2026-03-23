#!/usr/bin/env python3
"""
Generate Kira 사업계획서 PPTX — 학교 양식 + 패스트빌드 스타일.
Run: python scripts/generate_school_pptx.py
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Dimensions (16:9) ───────────────────────────────────────────────
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
_PX = SLIDE_W // 1920

def px(n): return int(n * _PX)

# ─── Colors ───────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
BLUE_D  = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_L  = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_50 = RGBColor(0xEF, 0xF6, 0xFF)
DARK    = RGBColor(0x1E, 0x29, 0x3B)
BODY    = RGBColor(0x47, 0x55, 0x69)
LIGHT   = RGBColor(0x94, 0xA3, 0xB8)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
BG      = RGBColor(0xF8, 0xFA, 0xFC)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
GREEN_L = RGBColor(0xDC, 0xFC, 0xE7)
RED     = RGBColor(0xEF, 0x44, 0x44)
YELLOW  = RGBColor(0xEA, 0xB3, 0x08)
ORANGE  = RGBColor(0xF9, 0x73, 0x16)
ORANGE_L= RGBColor(0xFF, 0xED, 0xD5)
G700    = RGBColor(0x33, 0x33, 0x33)
G500    = RGBColor(0x64, 0x74, 0x8B)
G400    = RGBColor(0x94, 0xA3, 0xB8)
G200    = RGBColor(0xE2, 0xE8, 0xF0)
G100    = RGBColor(0xF1, 0xF5, 0xF9)

FONT = '맑은 고딕'

# ─── Helpers ──────────────────────────────────────────────────────────

def _shape(slide, typ, x, y, w, h, fill=None, stroke=None, sw=Pt(1)):
    s = slide.shapes.add_shape(typ, px(x), px(y), px(w), px(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if stroke:
        s.line.color.rgb = stroke; s.line.width = sw
    else:
        s.line.fill.background()
    return s

def _rect(sl, x, y, w, h, **kw):
    return _shape(sl, MSO_SHAPE.RECTANGLE, x, y, w, h, **kw)

def _rrect(sl, x, y, w, h, r=0.06, **kw):
    s = _shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, **kw)
    try: s.adjustments[0] = r
    except Exception: pass
    return s

def _oval(sl, x, y, w, h, **kw):
    return _shape(sl, MSO_SHAPE.OVAL, x, y, w, h, **kw)

def _txt(sl, x, y, w, h, text, sz=14, bold=False, color=DARK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    try: tf.vertical_anchor = anchor
    except Exception: pass
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = FONT; p.alignment = align
    if spacing is not None:
        p.space_after = Pt(spacing)
    return tb

def _mtxt(sl, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    try: tf.vertical_anchor = anchor
    except Exception: pass
    for i, item in enumerate(lines):
        if len(item) == 4:
            text, sz, bold, color = item
            al = PP_ALIGN.LEFT
        else:
            text, sz, bold, color, al = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.bold = bold
        p.font.color.rgb = color; p.font.name = FONT; p.alignment = al
    return tb

def _header(sl, title):
    """패스트빌드 스타일 header bar."""
    _rect(sl, 0, 0, 1920, 185, fill=NAVY)
    _rect(sl, 0, 185, 1920, 4, fill=BLUE)
    _txt(sl, 80, 55, 800, 80, title, sz=28, bold=True, color=WHITE)

def _card(sl, x, y, w, h, title, body, accent=BLUE, title_sz=16):
    """Rounded card with title + body text."""
    _rrect(sl, x, y, w, h, fill=WHITE, stroke=BORDER, r=0.04)
    _txt(sl, x+20, y+16, w-40, 30, title, sz=title_sz, bold=True, color=DARK)
    _txt(sl, x+20, y+50, w-40, h-70, body, sz=12, color=BODY)

def _badge(sl, x, y, text, bg=BLUE, fg=WHITE, w=120):
    _rrect(sl, x, y, w, 28, fill=bg, r=0.5)
    _txt(sl, x, y+2, w, 24, text, sz=10, bold=True, color=fg, align=PP_ALIGN.CENTER)

def _add_table(sl, x, y, w, h, rows, col_widths=None):
    """Add a table. rows = [['h1','h2',...], ['r1c1','r1c2',...], ...]"""
    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 0
    tbl_shape = sl.shapes.add_table(n_rows, n_cols, px(x), px(y), px(w), px(h))
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = px(cw)
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                p.font.color.rgb = DARK if ri == 0 else BODY
                p.font.bold = (ri == 0)
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
    return tbl

# ═════════════════════════════════════════════════════════════════════
# SLIDES
# ═════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, 1920, 1080, fill=NAVY)
    # Decorative lines
    _rect(sl, 0, 0, 1920, 4, fill=BLUE)
    _rect(sl, 0, 1076, 1920, 4, fill=BLUE)
    # Title
    _txt(sl, 120, 300, 800, 100, "Kira", sz=64, bold=True, color=WHITE)
    _txt(sl, 120, 400, 800, 50, "공공조달 입찰 자동화 AI 플랫폼", sz=24, color=G400)
    _rect(sl, 120, 465, 80, 3, fill=BLUE)
    _txt(sl, 120, 490, 600, 30, "M&S Solutions  ·  2026년 3월", sz=16, color=G400)
    _txt(sl, 120, 530, 600, 30, "한성대학교 창업동아리 사업계획서", sz=18, bold=True, color=WHITE)
    # Tags
    _rrect(sl, 120, 590, 160, 32, fill=BLUE_D, r=0.5)
    _txt(sl, 120, 593, 160, 28, "AI · SaaS", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rrect(sl, 300, 590, 160, 32, fill=BLUE_D, r=0.5)
    _txt(sl, 300, 593, 160, 28, "공공조달", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rrect(sl, 480, 590, 200, 32, fill=BLUE_D, r=0.5)
    _txt(sl, 480, 593, 200, 28, "입찰 자동화", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Right side visual
    _rrect(sl, 1050, 180, 750, 720, fill=RGBColor(0x22, 0x33, 0x55), r=0.04)
    _txt(sl, 1100, 360, 650, 100, "AI Task\nAutomation", sz=48, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _txt(sl, 1100, 520, 650, 80, "분석 30초  ·  원클릭 생성\n154 테스트 통과  ·  E2E 검증", sz=16, color=G400, align=PP_ALIGN.CENTER)


def slide_02_motivation(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "창업의 계기")

    cards = [
        ("창업 아이디어 발상의 계기",
         "공공조달 시장 225조원, 등록 사업자 50만 개 중 96.7%가 중소기업.\n\n"
         "중소기업 BD팀은 1~3명. 공고 분석에 2~4시간, 제안서 작성에 2~4주.\n\n"
         "한 건에 올인하면 다른 기회를 놓치고, 동시 진행하면 품질이 떨어지는 구조적 문제를 발견."),
        ("창업 아이디어 구체화 과정 1",
         "기존 도구 분석:\n"
         "· 비드프로: 공고 검색만, AI 없음\n"
         "· 클라이원트: AI 분석만, 생성 없음\n"
         "· 디마툴즈: 낙찰가 예측만\n"
         "· 범용 AI (ChatGPT): 공공조달 전문성 없음, HWP 못 읽음\n\n"
         "→ '분석 + 생성 + 학습'을 하나의 플랫폼에서 제공하는 서비스가 없다!"),
        ("창업 아이디어 구체화 과정 2",
         "AI가 '분석만' 하는 것이 아니라 '문서까지 생성'하고,\n"
         "사용자 수정을 '학습'해서 점점 좋아지는 구조.\n\n"
         "→ 3계층 학습 엔진 설계:\n"
         "  Layer 1: 범용 지식 (전문가 노하우)\n"
         "  Layer 2: 회사 맞춤 학습\n"
         "  Layer 3: 승패 패턴 분석\n\n"
         "→ 'Kira' 개발 착수"),
    ]
    for i, (t, d) in enumerate(cards):
        x = 60 + i * 600
        _rrect(sl, x, 220, 570, 530, fill=WHITE, stroke=BORDER, r=0.04)
        _rrect(sl, x+20, 238, 530, 40, fill=NAVY, r=0.04)
        _txt(sl, x+30, 241, 510, 34, t, sz=14, bold=True, color=WHITE)
        _txt(sl, x+30, 300, 510, 430, d, sz=12, color=BODY)


def slide_03_social_venture(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "창업의 계기 (소셜 벤처)")

    cards = [
        ("사회적 문제 인식 계기",
         "중소기업은 전체 공공조달 사업자의 96.7%를 차지하지만,\n"
         "BD 인력 부족으로 양질의 공고를 놓치는 경우가 빈번.\n\n"
         "→ 대기업 대비 정보 비대칭·자원 격차로\n"
         "   중소기업의 공공조달 참여 기회가 구조적으로 제한됨"),
        ("사회적 문제 파악",
         "자격을 갖춘 중소기업도 제안서 품질 부족으로 탈락하는 사례 다수.\n\n"
         "· 제안서 외주 시 건당 300~1,000만원 (중소기업에겐 큰 부담)\n"
         "· BD 담당자 1~3명이 검색+분석+작성 전부 담당\n"
         "· 시간 부족 → 품질 저하 → 탈락 → 기회 상실의 악순환"),
        ("데이터를 통한 사회문제 파악",
         "공공조달 시장: 225조원/년 (2024, 조달청)\n"
         "등록 사업자: 50만 개\n"
         "중소기업 비율: 96.7%\n\n"
         "IT 서비스 분야 입찰 참여: 약 8만 개사\n"
         "→ AI 도구로 중소기업 경쟁력을 끌어올려\n"
         "   공공조달 참여 기회의 평등화를 실현"),
    ]
    for i, (t, d) in enumerate(cards):
        x = 60 + i * 600
        _rrect(sl, x, 220, 570, 530, fill=WHITE, stroke=BORDER, r=0.04)
        _rrect(sl, x+20, 238, 530, 40, fill=NAVY, r=0.04)
        _txt(sl, x+30, 241, 510, 34, t, sz=14, bold=True, color=WHITE)
        _txt(sl, x+30, 300, 510, 430, d, sz=12, color=BODY)


def slide_04_product(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "창업아이템 설명")

    # Before
    _txt(sl, 80, 210, 200, 30, "도입 전", sz=20, bold=True, color=RED)
    before = [
        ("공고 발견", "하루 수백 건에서\n수작업 검색\n→ 담당자 1명 × 2시간/일"),
        ("RFP 분석", "50~300쪽 PDF/HWP\n요건이 산재\n→ 2~4시간/건 소요"),
        ("제안서 작성", "50~200쪽 처음부터 집필\n→ 2~4주, 외주 300~1,000만원"),
        ("결과 학습", "낙찰/탈락 원인 분석 안 함\n→ 같은 실수 반복"),
    ]
    for i, (t, d) in enumerate(before):
        x = 80 + i * 220
        _rrect(sl, x, 250, 205, 210, fill=WHITE, stroke=RGBColor(0xFC, 0xBF, 0xBF), r=0.04)
        _txt(sl, x+14, 262, 177, 28, t, sz=13, bold=True, color=RED)
        _txt(sl, x+14, 298, 177, 150, d, sz=11, color=BODY)

    # Arrow
    _txt(sl, 900, 340, 80, 40, "→", sz=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # After
    _txt(sl, 1000, 210, 200, 30, "도입 후", sz=20, bold=True, color=GREEN)
    after = [
        ("자동 알림", "맞춤 키워드·금액·지역\n→ 관련 공고 자동 발견"),
        ("AI 분석 30초", "자격요건 자동 추출\nRFP 3섹션 요약\nGO/NO-GO 즉시 판단"),
        ("원클릭 생성", "제안서·PPT·WBS·\n실적기술서 5종 패키지\n→ 검토만 하면 끝"),
        ("수정 학습", "AI 초안 수정 → 패턴 감지\n3회 이상 → 자동 반영\n→ 쓸수록 좋아지는 AI"),
    ]
    for i, (t, d) in enumerate(after):
        x = 1000 + i * 220
        _rrect(sl, x, 250, 205, 210, fill=WHITE, stroke=RGBColor(0xBB, 0xF7, 0xD0), r=0.04)
        _txt(sl, x+14, 262, 177, 28, t, sz=13, bold=True, color=GREEN)
        _txt(sl, x+14, 298, 177, 150, d, sz=11, color=BODY)

    # Bottom summary
    _rrect(sl, 80, 490, 1760, 120, fill=BLUE_50, stroke=BLUE, sw=Pt(2), r=0.04)
    _mtxt(sl, 120, 498, 1680, 100, [
        ("핵심 가치: 시간과 품질의 트레이드오프 해결", 18, True, DARK),
        ("분석 4시간 → 30초  |  제안서 2~4주 → 원클릭 생성  |  외주 300만원+ → 월 9.9만원  |  학습 없음 → 자동 학습", 13, False, BODY),
    ])
    _txt(sl, 1500, 502, 300, 30, "ROI 20x", sz=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def slide_05_market(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "시장 분석")

    # TAM circle (largest)
    _oval(sl, 200, 240, 600, 480, fill=BLUE_50, stroke=BLUE, sw=Pt(2))
    _txt(sl, 350, 260, 300, 24, "TAM", sz=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _txt(sl, 280, 290, 440, 30, "전체 시장", sz=12, color=BODY, align=PP_ALIGN.CENTER)

    # SAM circle (medium)
    _oval(sl, 300, 360, 400, 320, fill=BLUE_L, stroke=BLUE_D, sw=Pt(2))
    _txt(sl, 400, 375, 200, 24, "SAM", sz=14, bold=True, color=BLUE_D, align=PP_ALIGN.CENTER)
    _txt(sl, 370, 400, 260, 24, "접근 가능 시장", sz=11, color=BODY, align=PP_ALIGN.CENTER)

    # SOM circle (smallest)
    _oval(sl, 380, 460, 240, 200, fill=BLUE, stroke=BLUE_D, sw=Pt(2))
    _txt(sl, 410, 510, 180, 24, "SOM", sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, 400, 540, 200, 24, "초기 핵심 시장", sz=11, color=WHITE, align=PP_ALIGN.CENTER)

    # Details (right side)
    details = [
        ("TAM  ·  전체 시장", "225조원 /년", [
            "대한민국 공공조달 시장 (2024, 조달청)",
            "등록 사업자 50만 개",
            "IT 서비스 분야 약 8만 개",
        ], BLUE),
        ("SAM  ·  접근 가능 시장", "720억원 /년", [
            "IT/SW 분야 활발한 입찰 기업: 약 2만 개",
            "평균 구독료 월 30만원 가정",
            "2만 × 360만원 = 720억원/년",
        ], BLUE_D),
        ("SOM  ·  획득 가능 시장 (3년)", "36억원 /년", [
            "1차 타겟: IT 서비스 중소기업 (10~100명)",
            "3년 내 목표 점유율 5%",
            "1,000개사 × 360만원 = 36억원/년",
        ], NAVY),
    ]
    for i, (label, amount, items, accent) in enumerate(details):
        y = 220 + i * 175
        _rrect(sl, 880, y, 960, 155, fill=WHITE, stroke=BORDER, r=0.04)
        _txt(sl, 900, y+12, 400, 24, label, sz=13, bold=True, color=accent)
        _txt(sl, 1540, y+8, 280, 34, amount, sz=26, bold=True, color=accent, align=PP_ALIGN.RIGHT)
        for j, item in enumerate(items):
            _txt(sl, 920, y+46 + j*28, 880, 24, "· " + item, sz=11, color=BODY)


def slide_06_competitor(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "경쟁사 분석 및 포지셔닝")

    _txt(sl, 80, 200, 300, 28, "경쟁사 분석", sz=18, bold=True, color=DARK)
    rows = [
        ["기능", "비드프로", "클라이원트", "디마툴즈", "Kira"],
        ["공고 검색", "●", "●", "—", "●"],
        ["AI 분석", "—", "●", "—", "●"],
        ["GO/NO-GO", "—", "—", "△", "●"],
        ["제안서 생성", "—", "—", "—", "●"],
        ["PPT/WBS", "—", "—", "—", "●"],
        ["학습 엔진", "—", "—", "—", "●"],
    ]
    _add_table(sl, 80, 235, 900, 350, rows)

    # Positioning map
    _txt(sl, 1020, 200, 300, 28, "포지셔닝 맵", sz=18, bold=True, color=DARK)
    _rrect(sl, 1020, 235, 830, 350, fill=G100, stroke=BORDER, r=0.04)
    # Axes
    _rect(sl, 1060, 420, 750, 2, fill=G400)
    _rect(sl, 1435, 260, 2, 300, fill=G400)
    _txt(sl, 1040, 260, 100, 20, "생성 능력 ↑", sz=9, color=G500)
    _txt(sl, 1710, 425, 100, 20, "범위 →", sz=9, color=G500)
    # Kira (top-right)
    _oval(sl, 1500, 280, 100, 50, fill=BLUE)
    _txt(sl, 1500, 287, 100, 36, "Kira", sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 클라이원트 (bottom-left)
    _oval(sl, 1200, 440, 100, 50, fill=G400)
    _txt(sl, 1200, 447, 100, 36, "클라이원트", sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 비드프로 (bottom-right)
    _oval(sl, 1600, 440, 100, 50, fill=G400)
    _txt(sl, 1600, 447, 100, 36, "비드프로", sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 디마툴즈 (bottom-center)
    _oval(sl, 1350, 460, 100, 50, fill=G400)
    _txt(sl, 1350, 467, 100, 36, "디마툴즈", sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Key differentiator
    _rrect(sl, 80, 610, 1760, 80, fill=BLUE_50, r=0.04)
    _txt(sl, 120, 622, 1680, 50,
         "핵심 차별화:  유일하게 '분석 → 생성 → 학습'을 한 플랫폼에서 제공  |  HWP 완벽 지원  |  3계층 학습 엔진",
         sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def slide_07_persona(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "고객 페르소나")

    # Persona card
    _rrect(sl, 80, 210, 800, 500, fill=WHITE, stroke=BORDER, r=0.04)
    # Avatar
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, px(120), px(250), px(100), px(100))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    _txt(sl, 120, 275, 100, 50, "김", sz=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _mtxt(sl, 250, 250, 600, 440, [
        ("김영수 (37세)", 22, True, DARK),
        ("IT 중소기업 영업/BD 팀장", 14, True, BLUE),
        ("", 6, False, DARK),
        ("회사 규모: 직원 25명, 연매출 40억원", 12, False, BODY),
        ("입찰 빈도: 연 8~12건 공공 입찰 참여", 12, False, BODY),
        ("BD팀: 본인 포함 2명", 12, False, BODY),
        ("", 6, False, DARK),
        ("Pain Points:", 13, True, RED),
        ("· 매일 2시간 공고 검색 — 관련 없는 것도 전부 봐야 함", 11, False, BODY),
        ("· 제안서 작성 시 주말·야근 필수 — 다른 영업 기회 놓침", 11, False, BODY),
        ("· 외주 비용 부담 — 건당 500만원 이상", 11, False, BODY),
        ("· 낙찰/탈락 원인을 몰라서 같은 실수 반복", 11, False, BODY),
        ("", 6, False, DARK),
        ("Needs:", 13, True, GREEN),
        ("· 관련 공고만 자동으로 알려주는 시스템", 11, False, BODY),
        ("· 30분 안에 GO/NO-GO 판단할 수 있는 분석 도구", 11, False, BODY),
        ("· 회사 특성을 반영한 자동 제안서 생성", 11, False, BODY),
    ])

    # Insight
    _rrect(sl, 920, 210, 920, 500, fill=G100, stroke=BORDER, r=0.04)
    _txt(sl, 950, 228, 860, 28, "고객 인사이트", sz=18, bold=True, color=DARK)

    insights = [
        ("핵심 고통", "시간과 품질의 트레이드오프\n한 건에 올인 → 다른 기회 놓침\n동시 진행 → 품질 하락", RED),
        ("의사결정 기준", "① 시간 절감 효과\n② 제안서 품질 (탈락 위험 감소)\n③ 비용 대비 효과 (ROI)", BLUE),
        ("도입 장벽", "① AI 생성 제안서 품질 불신\n② 기존 작업 방식 변경 저항\n③ 보안 우려 (회사 정보)", YELLOW),
        ("해결 전략", "① 무료 분석 5건 체험 제공\n② Before/After 비교 데모\n③ 점진적 도입 (분석→생성→학습)", GREEN),
    ]
    for i, (title, body, accent) in enumerate(insights):
        y = 268 + i * 106
        _rrect(sl, 950, y, 860, 96, fill=WHITE, stroke=BORDER, r=0.03)
        _txt(sl, 970, y+8, 200, 24, title, sz=13, bold=True, color=accent)
        _txt(sl, 970, y+34, 820, 56, body, sz=11, color=BODY)


def slide_08_customer_dist(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "고객 분포 현황")

    segments = [
        ("핵심 고객\n1차 타겟", "IT 서비스 중소기업", BLUE, 380, [
            "종업원 10~100명",
            "연 5건+ 공공 입찰 참여",
            "BD 담당자 1~3명",
            "AI 도구 수용도 높음",
            "약 2만 개사",
        ]),
        ("확장 고객\n2차 타겟", "건설·엔지니어링·컨설팅", GREEN, 280, [
            "입찰 빈도 높음",
            "제안서 부담 큼",
            "HWP 문서 의존도 높음",
            "6개월 후 진출",
        ]),
        ("미래 고객\n3차 타겟", "정부지원사업 신청자", ORANGE, 200, [
            "스타트업, 소상공인",
            "TIPS, AI 바우처",
            "창업패키지 계획서",
            "12개월 후 진출",
        ]),
    ]
    # Concentric layout
    for i, (label, sub, accent, size, items) in enumerate(segments):
        cx = 350
        cy = 500
        _oval(sl, cx - size//2, cy - size//2, size, size, fill=None, stroke=accent, sw=Pt(3))
        _txt(sl, cx - 60, cy - size//2 + 10, 120, 40, label, sz=11, bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Details (right side cards)
    for i, (label, sub, accent, _, items) in enumerate(segments):
        x = 700 + i * 400
        _rrect(sl, x, 240, 370, 380, fill=WHITE, stroke=accent, sw=Pt(2), r=0.04)
        _badge(sl, x+20, 258, label.split('\n')[0], bg=accent, w=140)
        _txt(sl, x+20, 300, 330, 28, sub, sz=16, bold=True, color=DARK)
        for j, item in enumerate(items):
            _txt(sl, x+20, 340 + j*30, 330, 26, "· " + item, sz=12, color=BODY)

    # Bottom
    _rrect(sl, 80, 650, 1760, 80, fill=G100, r=0.04)
    _txt(sl, 120, 660, 1680, 60,
         "진출 전략:  IT 서비스 (0~6개월) → 건설·엔지니어링 (6~12개월) → 정부지원사업 (12개월+)",
         sz=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def _bm_slide(prs, num, title, left_name, center_name, right_name,
              flows, revenue_title, revenue_body):
    """비즈니스 모델 슬라이드 공통 패턴 (3-element flow)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, f"비즈니스 모델 {num}")
    _txt(sl, 80, 200, 800, 30, title, sz=18, bold=True, color=DARK)

    # 3 Ovals
    _oval(sl, 150, 320, 220, 150, fill=G100, stroke=BORDER, sw=Pt(2))
    _txt(sl, 150, 365, 220, 40, left_name, sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    _oval(sl, 700, 280, 280, 190, fill=BLUE, stroke=BLUE_D, sw=Pt(2))
    _txt(sl, 700, 345, 280, 40, center_name, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _oval(sl, 1300, 320, 220, 150, fill=G100, stroke=BORDER, sw=Pt(2))
    _txt(sl, 1300, 365, 220, 40, right_name, sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Arrows
    _txt(sl, 400, 370, 100, 30, "→", sz=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _txt(sl, 550, 370, 100, 30, "←", sz=28, bold=True, color=BODY, align=PP_ALIGN.CENTER)
    _txt(sl, 1050, 370, 100, 30, "→", sz=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _txt(sl, 1180, 370, 100, 30, "←", sz=28, bold=True, color=BODY, align=PP_ALIGN.CENTER)

    # Flow labels
    for i, (direction, text, y_off) in enumerate(flows):
        if direction == "LtoC":
            _txt(sl, 380, 300 + y_off, 280, 24, text, sz=10, color=BODY, align=PP_ALIGN.CENTER)
        elif direction == "CtoL":
            _txt(sl, 380, 420 + y_off, 280, 24, text, sz=10, color=BODY, align=PP_ALIGN.CENTER)
        elif direction == "CtoR":
            _txt(sl, 1020, 300 + y_off, 280, 24, text, sz=10, color=BODY, align=PP_ALIGN.CENTER)
        elif direction == "RtoC":
            _txt(sl, 1020, 420 + y_off, 280, 24, text, sz=10, color=BODY, align=PP_ALIGN.CENTER)

    # Revenue box
    _rrect(sl, 80, 530, 1760, 170, fill=NAVY, r=0.04)
    _txt(sl, 120, 545, 600, 30, revenue_title, sz=16, bold=True, color=WHITE)
    _txt(sl, 120, 580, 1680, 100, revenue_body, sz=13, color=G400)


def slide_09_bm1(prs):
    _bm_slide(prs, 1, "B2B SaaS 구독 모델 (Freemium → PRO)",
        "중소기업\nBD 담당자", "Kira\n플랫폼", "공공조달\n공고",
        [
            ("LtoC", "구독료 ₩99,000/월", 0),
            ("CtoL", "AI 분석 + 문서 생성", 0),
            ("CtoR", "공고 크롤링", 0),
            ("RtoC", "공고 데이터", 0),
        ],
        "수익 모델:  FREE(₩0, 월5건) → PRO(₩99,000/월, 무제한+생성)",
        "핵심 지표: Free→PRO 전환율 15% 목표  |  ARPU ₩99,000  |  그로스 마진 ~82%  |  LTV/CAC ~6x")


def slide_10_bm2(prs):
    _bm_slide(prs, 2, "Enterprise 전담 모델",
        "중견·대기업\nBD팀", "Kira\nEnterprise", "전담 학습\n모델",
        [
            ("LtoC", "연간 계약 (별도 협의)", 0),
            ("CtoL", "전담 모델 + SLA + 온프레미스", 0),
            ("CtoR", "회사 데이터 학습", 0),
            ("RtoC", "맞춤 AI 모델", 0),
        ],
        "수익 모델:  연간 계약 ₩5,000만원~  |  전담 학습 모델 + SLA + API 연동",
        "Y2: 3건 (₩1.5억)  →  Y3: 10건 (₩5억)  |  높은 LTV, 낮은 이탈률")


def slide_11_bm3(prs):
    _bm_slide(prs, 3, "승패 분석 데이터 인사이트 (Phase 3)",
        "프리미엄\n구독자", "Kira\nInsights", "승패 데이터\n(익명화)",
        [
            ("LtoC", "프리미엄 구독", 0),
            ("CtoL", "승률 예측 + 인사이트 리포트", 0),
            ("CtoR", "낙찰/탈락 결과 축적", 0),
            ("RtoC", "패턴 분석 데이터", 0),
        ],
        "수익 모델:  승패 분석 리포트 프리미엄 구독 (PRO에 추가)",
        "데이터 플라이휠: 사용자↑ → 데이터↑ → 정확도↑ → 사용자↑  |  Layer 3 학습 엔진 기반")


def slide_12_bm4(prs):
    _bm_slide(prs, 4, "정부지원사업 사업계획서 확장",
        "스타트업\n소상공인", "Kira\nGov+", "정부지원사업\n(TIPS, AI바우처)",
        [
            ("LtoC", "구독료 또는 건별 결제", 0),
            ("CtoL", "사업계획서 자동 생성", 0),
            ("CtoR", "공고 정보 수집", 0),
            ("RtoC", "지원사업 데이터", 0),
        ],
        "수익 모델:  정부지원사업 계획서 생성 (TIPS, AI바우처, 창업패키지)",
        "2027 H1 런칭 예정  |  TAM 확장: 공공조달 + 정부지원사업 통합 플랫폼")


def slide_13_bm5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "비즈니스 모델 5")
    _txt(sl, 80, 200, 800, 30, "통합 생태계: 분석 + 생성 + 학습 + 인사이트", sz=18, bold=True, color=DARK)

    # Center oval
    _oval(sl, 760, 320, 400, 200, fill=BLUE, stroke=BLUE_D, sw=Pt(3))
    _txt(sl, 760, 385, 400, 50, "Kira\n플랫폼", sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 4 surrounding revenue streams
    revs = [
        (300, 260, "SaaS 구독", "FREE → PRO\n₩99,000/월"),
        (1350, 260, "Enterprise", "전담 모델\n연간 계약"),
        (300, 520, "데이터 인사이트", "승패 분석\n프리미엄 리포트"),
        (1350, 520, "정부지원사업", "사업계획서\n자동 생성"),
    ]
    for x, y, title, body in revs:
        _rrect(sl, x, y, 250, 130, fill=WHITE, stroke=BORDER, r=0.04)
        _txt(sl, x+15, y+12, 220, 24, title, sz=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        _txt(sl, x+15, y+42, 220, 70, body, sz=11, color=BODY, align=PP_ALIGN.CENTER)

    # Bottom summary
    _rrect(sl, 80, 690, 1760, 80, fill=NAVY, r=0.04)
    _txt(sl, 120, 702, 1680, 50,
         "3년 매출 목표:  Y1 ₩1.2억  →  Y2 ₩7.3억  →  Y3 ₩23억  |  그로스 마진 82%",
         sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_14_social_bm(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "소셜벤처 비즈니스 모델")

    values = [
        ("사회적 가치", BLUE, [
            "중소기업 공공조달 참여 기회 확대",
            "정보 비대칭 해소 — AI가 대기업급 분석 제공",
            "BD 인력 부족 해소 — 1인이 다건 동시 진행 가능",
            "중소기업 경쟁력 = 지역 경제 활성화",
        ]),
        ("경제적 가치", GREEN, [
            "제안서 외주 비용 절감 (300만원 → 9.9만원/월)",
            "시간 절감 (분석 4시간→30초, 제안서 2주→1시간)",
            "동시 다건 진행 → 매출 기회 확대",
            "Layer 2 학습 → 장기 품질 향상",
        ]),
        ("환경적 가치", ORANGE, [
            "종이 문서 → 디지털 전환 촉진",
            "불필요한 입찰 참여 감소 (GO/NO-GO 판단)",
            "업무 효율화 → 야근·과로 감소",
            "데이터 기반 의사결정 문화 확산",
        ]),
    ]
    for i, (title, accent, items) in enumerate(values):
        x = 60 + i * 620
        _rrect(sl, x, 220, 590, 350, fill=WHITE, stroke=accent, sw=Pt(2), r=0.04)
        _badge(sl, x+20, 238, title, bg=accent, w=140)
        for j, item in enumerate(items):
            _txt(sl, x+20, 285 + j*36, 550, 30, "· " + item, sz=12, color=BODY)

    # Impact metrics
    _txt(sl, 80, 595, 400, 28, "임팩트 지표", sz=18, bold=True, color=DARK)
    metrics = [
        ["지표", "현재", "Y1 목표", "Y3 목표"],
        ["지원 중소기업 수", "—", "100개사", "1,000개사"],
        ["절감 비용 (총합)", "—", "₩2억+", "₩50억+"],
        ["입찰 참여 건수 증가율", "—", "+30%", "+100%"],
    ]
    _add_table(sl, 80, 630, 1760, 140, metrics)


def slide_15_validation(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "시장 검증 및 제품 완성도")

    _txt(sl, 80, 200, 400, 28, "제품 완성도 (구현 현황)", sz=18, bold=True, color=DARK)

    stats = [
        ("154개", "테스트 통과", BLUE),
        ("495유닛", "Layer 1 지식 DB", GREEN),
        ("30초", "AI 분석 속도", ORANGE),
        ("5종", "문서 자동 생성", NAVY),
    ]
    for i, (num, label, accent) in enumerate(stats):
        x = 80 + i * 440
        _rrect(sl, x, 240, 410, 120, fill=WHITE, stroke=accent, sw=Pt(2), r=0.04)
        _txt(sl, x+20, 258, 200, 50, num, sz=36, bold=True, color=accent)
        _txt(sl, x+20, 315, 370, 24, label, sz=14, color=BODY)

    _txt(sl, 80, 385, 400, 28, "구현 완료 기능", sz=18, bold=True, color=DARK)
    features = [
        ["기능", "설명", "상태"],
        ["나라장터 실시간 검색", "키워드·업무구분·지역·금액 필터", "✅ 운영 중"],
        ["문서 자동 파싱", "PDF, DOCX, HWP, HWPX, Excel, PPT", "✅ 운영 중"],
        ["AI 자격요건 추출", "멀티패스 병렬 + 동의어 사전 17개", "✅ 운영 중"],
        ["제안서 DOCX 자동 생성", "Layer 1+2 지식 기반, 100+ 페이지", "✅ 운영 중"],
        ["PPT 발표자료 생성", "KRDS 디자인 + 예상질문 10개", "✅ 운영 중"],
        ["WBS/실적기술서 생성", "XLSX 간트차트 + DOCX 보고서", "✅ 운영 중"],
        ["수정 diff 학습", "패턴 감지 → 3회+ 자동 반영", "✅ 운영 중"],
        ["맞춤 공고 알림", "키워드·금액·지역 필터 + 이메일", "✅ 운영 중"],
    ]
    _add_table(sl, 80, 420, 1760, 300, features, col_widths=[400, 900, 460])


def slide_16_financials(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "소요자금 및 예상 매출")

    _txt(sl, 80, 200, 400, 28, "소요 자금: ₩2.0억", sz=18, bold=True, color=DARK)
    fund = [
        ["항목", "비율", "금액", "용도"],
        ["제품 개발", "50%", "₩1.0억", "개발자 채용 + Layer 3 + 고도화"],
        ["마케팅/GTM", "25%", "₩0.5억", "콘텐츠 + 커뮤니티 + 이벤트"],
        ["운영", "25%", "₩0.5억", "인프라 + 법무 + 운영비"],
    ]
    _add_table(sl, 80, 235, 860, 160, fund)

    _txt(sl, 980, 200, 400, 28, "3개년 매출 전망", sz=18, bold=True, color=DARK)
    rev = [
        ["항목", "Y1 (2026)", "Y2 (2027)", "Y3 (2028)"],
        ["유료 고객", "100", "400", "1,000"],
        ["MRR", "₩9.9M", "₩48M", "₩150M"],
        ["총 매출", "₩1.2억", "₩7.3억", "₩23억"],
        ["영업이익", "△₩2.0억", "₩0.4억", "₩11.4억"],
        ["영업이익률", "-167%", "5.5%", "49.6%"],
    ]
    _add_table(sl, 980, 235, 860, 210, rev)

    # Unit Economics
    _txt(sl, 80, 430, 400, 28, "Unit Economics (PRO 플랜)", sz=18, bold=True, color=DARK)
    ue = [
        ["항목", "금액"],
        ["ARPU (월)", "₩99,000"],
        ["LLM API 비용", "~₩15,000"],
        ["인프라 비용", "~₩3,000"],
        ["그로스 마진", "~82%"],
        ["CAC", "~₩200,000"],
        ["LTV (12개월)", "₩1,188,000"],
        ["LTV/CAC", "~6x"],
    ]
    _add_table(sl, 80, 465, 500, 280, ue)

    # BEP
    _rrect(sl, 640, 465, 1200, 80, fill=BLUE_50, r=0.04)
    _mtxt(sl, 670, 475, 1140, 60, [
        ("BEP (손익분기점): Y2 중반", 18, True, BLUE),
        ("유료 고객 ~250개사 시점에서 흑자 전환", 13, False, BODY),
    ])

    _rrect(sl, 640, 570, 1200, 170, fill=NAVY, r=0.04)
    _mtxt(sl, 670, 580, 1140, 150, [
        ("비용 구조 (Y1 → Y3)", 16, True, WHITE),
        ("", 6, False, DARK),
        ("인건비: ₩2.4억 → ₩7.2억  |  LLM API: ₩0.2억 → ₩1.8억", 12, False, G400),
        ("인프라: ₩0.1억 → ₩0.6억  |  마케팅: ₩0.3억 → ₩1.5억", 12, False, G400),
        ("총 비용: ₩3.2억 → ₩11.6억", 13, True, WHITE),
    ])


def slide_17_kpi(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "핵심 성과지표 및 재무 목표")

    kpis = [
        ("500", "MAU", "월간 활성 사용자\n(Y1 목표)", BLUE),
        ("15%", "Free→PRO 전환율", "무료 → 유료 전환", GREEN),
        ("₩10M", "MRR", "월간 반복 매출\n(Y1 목표)", ORANGE),
        ("< 5%", "Churn Rate", "월간 이탈률", RED),
    ]
    for i, (num, label, desc, accent) in enumerate(kpis):
        x = 80 + i * 450
        _rrect(sl, x, 220, 420, 150, fill=WHITE, stroke=accent, sw=Pt(2), r=0.04)
        _txt(sl, x+20, 235, 200, 50, num, sz=38, bold=True, color=accent)
        _txt(sl, x+20, 290, 380, 24, label, sz=14, bold=True, color=DARK)
        _txt(sl, x+20, 318, 380, 40, desc, sz=11, color=BODY)

    # Edit Rate (unique KPI)
    _rrect(sl, 80, 400, 1760, 200, fill=NAVY, r=0.04)
    _txt(sl, 120, 415, 1200, 30, "수정률 (Edit Rate) — Kira 고유 KPI", sz=18, bold=True, color=WHITE)
    _txt(sl, 120, 450, 1600, 30, "AI 초안 대비 사용자 수정 비율 → 낮아질수록 AI 학습 완성도 높음 → 이탈 방어 강화", sz=13, color=G400)

    stages = [("v1", "수정 45%", RED), ("v5", "수정 20%", YELLOW), ("v10", "수정 8%", GREEN)]
    for i, (ver, rate, clr) in enumerate(stages):
        x = 120 + i * 350
        _txt(sl, x, 500, 120, 30, ver, sz=16, bold=True, color=clr)
        _txt(sl, x+60, 500, 200, 30, rate, sz=16, bold=True, color=clr)
        if i < 2:
            _txt(sl, x+260, 500, 40, 30, "→", sz=20, bold=True, color=G400)
    _txt(sl, 1200, 500, 500, 30, '= "학습도 92%"', sz=18, bold=True, color=GREEN)

    # Financial goals
    _txt(sl, 80, 630, 400, 28, "재무 목표", sz=18, bold=True, color=DARK)
    goals = [
        ("Y1 목표", "유료 100개사 · MRR ₩10M", BLUE),
        ("Y2 목표", "유료 400개사 · BEP 달성 · Enterprise 3건", GREEN),
        ("Y3 목표", "유료 1,000개사 · 매출 ₩23억 · 영업이익률 49.6%", ORANGE),
    ]
    for i, (label, text, accent) in enumerate(goals):
        x = 80 + i * 600
        _rrect(sl, x, 665, 570, 60, fill=WHITE, stroke=accent, sw=Pt(2), r=0.03)
        _txt(sl, x+15, 672, 130, 36, label, sz=12, bold=True, color=accent)
        _txt(sl, x+150, 672, 400, 36, text, sz=12, color=BODY)


def slide_18_team(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "팀 소개")

    # Founder card
    _rrect(sl, 80, 220, 900, 400, fill=WHITE, stroke=BORDER, r=0.04)
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, px(120), px(270), px(120), px(120))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    _txt(sl, 120, 300, 120, 60, "민", sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _mtxt(sl, 270, 260, 680, 340, [
        ("민경욱", 28, True, DARK),
        ("대표 / 기술 총괄 (CEO & CTO)", 16, True, BLUE),
        ("bill.min122@gmail.com", 12, False, G500),
        ("", 8, False, DARK),
        ("Skills:", 13, True, DARK),
        ("· 풀스택 개발 (React, Python, FastAPI, Next.js)", 12, False, BODY),
        ("· AI/ML (OpenAI GPT-4, RAG, ChromaDB, 임베딩)", 12, False, BODY),
        ("· 공공조달 도메인 전문가", 12, False, BODY),
        ("· Kira 전체 아키텍처 설계·구현 (1인 개발)", 12, False, BODY),
        ("· 154+ 테스트, 30,000+ 코드 라인", 12, False, BODY),
    ])

    # Skills badges
    skills = ["React", "Python", "FastAPI", "AI/ML", "RAG", "ChromaDB"]
    for i, skill in enumerate(skills):
        _badge(sl, 270 + i * 100, 540, skill, bg=BLUE_D, w=90)

    # Tech proof
    _rrect(sl, 1020, 220, 820, 400, fill=NAVY, r=0.04)
    _txt(sl, 1050, 240, 760, 30, "구현 완료 증빙", sz=18, bold=True, color=WHITE)
    proofs = [
        ("총 테스트", "154개"),
        ("E2E 검증", "통과 (2026-02-28)"),
        ("Layer 1 지식", "495유닛 (ChromaDB)"),
        ("문서 파싱", "PDF, DOCX, HWP, HWPX, Excel, PPT"),
        ("문서 생성", "DOCX, PPTX, XLSX, PNG"),
        ("코드 라인", "30,000+"),
        ("보안 레이어", "HMAC, CSRF, SSRF, 입력 검증"),
    ]
    for i, (label, val) in enumerate(proofs):
        y = 285 + i * 36
        _txt(sl, 1050, y, 300, 30, label, sz=12, color=G400)
        _txt(sl, 1380, y, 460, 30, val, sz=12, bold=True, color=WHITE)

    # Hiring plan
    _txt(sl, 80, 650, 300, 28, "채용 계획", sz=16, bold=True, color=DARK)
    hires = [
        ("백엔드 개발자", "Y1 Q2"), ("프론트엔드 개발자", "Y1 Q3"),
        ("BD/마케팅", "Y1 Q2"), ("CS/제안서 전문가", "Y1 Q4"),
    ]
    for i, (role, when) in enumerate(hires):
        x = 80 + i * 460
        _rrect(sl, x, 685, 430, 50, fill=BLUE_50, r=0.03)
        _txt(sl, x+15, 692, 250, 30, role, sz=12, bold=True, color=DARK)
        _txt(sl, x+280, 692, 130, 30, when, sz=12, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)


def slide_19_progress(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "진행 과정")

    # Timeline phases
    phases = [
        ("문제/솔루션 진단", BLUE, True, [
            ("문제/대안 검증", True, "공공조달 시장 조사 + 경쟁사 분석 완료"),
            ("솔루션 검증", True, "3계층 학습 엔진 설계 + MVP 구현"),
        ]),
        ("제품·서비스 적합성", GREEN, True, [
            ("낮은 수준 MVP", True, "공고 검색 + AI 분석 + 제안서 생성 완성"),
            ("높은 수준 MVP", True, "PPT/WBS/실적기술서 + 알림 + 체크리스트"),
        ]),
        ("비즈니스 모델 검증", ORANGE, False, [
            ("핵심지표 검증", False, "베타 50개사 대상 수정률 추적"),
            ("수익모델 검증", False, "FREE→PRO 전환율 측정"),
            ("채널 검증", False, "콘텐츠 마케팅 + 커뮤니티"),
        ]),
    ]
    for i, (title, accent, done, items) in enumerate(phases):
        x = 60 + i * 620
        _rrect(sl, x, 220, 590, 70, fill=accent if done else G100,
               stroke=accent, sw=Pt(2), r=0.04)
        status = "✅ " if done else "⏳ "
        _txt(sl, x+20, 235, 550, 36, status + title,
             sz=16, bold=True, color=WHITE if done else DARK)

        for j, (item, item_done, desc) in enumerate(items):
            y = 310 + j * 110
            _rrect(sl, x, y, 590, 95, fill=WHITE, stroke=BORDER, r=0.03)
            mark = "✅ " if item_done else "○ "
            _txt(sl, x+16, y+10, 558, 24, mark + item, sz=13, bold=True,
                 color=GREEN if item_done else BODY)
            _txt(sl, x+16, y+40, 558, 50, desc, sz=11, color=BODY)

    # Current status
    _rrect(sl, 60, 680, 1800, 80, fill=BLUE_50, stroke=BLUE, sw=Pt(2), r=0.04)
    _txt(sl, 100, 692, 1720, 50,
         "현재 단계:  Phase 1-2 구현 완료 (154 테스트 통과, E2E 검증)  →  베타 런칭 준비 중",
         sz=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def slide_20_vision(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "창업회사의 비전")

    # Vision statement
    _oval(sl, 700, 240, 520, 200, fill=BLUE, stroke=BLUE_D, sw=Pt(3))
    _txt(sl, 700, 290, 520, 100,
         "VISION\n모든 회사 업무의\n효율성을 AI로 혁신",
         sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Values
    values = [
        (300, 330, "공공조달\n자동화", "입찰의 모든 순간을\nAI가 함께합니다"),
        (1300, 330, "회사 맞춤\n학습", "쓸수록 우리 회사\n색깔이 나오는 AI"),
        (300, 530, "중소기업\n경쟁력", "대기업급 분석·생성을\n모든 중소기업에게"),
        (1300, 530, "데이터\n플라이휠", "사용자가 많을수록\n더 정확해지는 AI"),
    ]
    for x, y, title, desc in values:
        _oval(sl, x, y, 240, 140, fill=WHITE, stroke=BLUE, sw=Pt(2))
        _txt(sl, x+20, y+15, 200, 50, title, sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        _txt(sl, x+20, y+75, 200, 50, desc, sz=10, color=BODY, align=PP_ALIGN.CENTER)

    # Bottom: Short/Mid/Long term
    _txt(sl, 80, 720, 200, 28, "목표 로드맵", sz=16, bold=True, color=DARK)
    terms = [
        ("단기 (2026)", "베타 런칭 + 100개사\nFREE→PRO 전환 최적화", BLUE),
        ("중기 (2027)", "500개사 + Enterprise 3건\nLayer 3 런칭", GREEN),
        ("장기 (2028+)", "1,000개사 + 데이터 플라이휠\n공공조달+민간 입찰 통합", ORANGE),
    ]
    for i, (label, goals, accent) in enumerate(terms):
        x = 320 + i * 540
        _rrect(sl, x, 720, 510, 70, fill=WHITE, stroke=accent, sw=Pt(2), r=0.03)
        _txt(sl, x+15, 724, 150, 24, label, sz=12, bold=True, color=accent)
        _txt(sl, x+15, 744, 480, 40, goals, sz=10, color=BODY)


def slide_21_marketing(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "마케팅 계획")

    channels = [
        ("마케팅 방안 1", "콘텐츠 마케팅",
         "\"공공입찰 GO/NO-GO 판단법\" 등 실용 콘텐츠\n블로그 주 2회, 유튜브 월 2회\n→ 오가닉 유입 + SEO 확보", BLUE),
        ("마케팅 방안 2", "무료 체험 (PLG)",
         "FREE 플랜 월 5건 분석 제공\n실제 공고로 체험 → 가치 확인 → PRO 전환\n→ 15% 전환율 목표", GREEN),
        ("마케팅 방안 3", "입찰 커뮤니티",
         "나라장터 관련 카페·포럼 참여\n무료 분석 데모 제공\n→ 초기 얼리어답터 확보", ORANGE),
        ("마케팅 방안 4", "파트너십",
         "제안서 컨설팅사와 협업 (AI 도구 제공)\n정부 디지털 바우처 활용\n→ 신뢰 확보 + 비용 부담 감소", NAVY),
        ("마케팅 방안 5", "직접 영업 (Enterprise)",
         "BD 담당자 대상 1:1 데모\n무료 분석 3건 + ROI 시뮬레이션\n→ 연간 계약 전환", RED),
    ]
    for i, (num, title, body, accent) in enumerate(channels):
        col = i % 3
        row = i // 3
        x = 80 + col * 610
        y = 210 + row * 280
        _rrect(sl, x, y, 580, 255, fill=WHITE, stroke=BORDER, r=0.04)
        _rrect(sl, x+16, y+16, 548, 36, fill=accent, r=0.04)
        _txt(sl, x+24, y+19, 532, 30, f"{num}: {title}", sz=13, bold=True, color=WHITE)
        _txt(sl, x+24, y+68, 532, 170, body, sz=12, color=BODY)

    # Metrics
    metrics = [("CAC", "~₩200,000"), ("LTV (12개월)", "₩1,188,000"), ("LTV/CAC", "~6x")]
    for i, (label, val) in enumerate(metrics):
        x = 80 + i * 300
        _rrect(sl, x, 770, 270, 50, fill=NAVY, r=0.03)
        _txt(sl, x+15, 778, 120, 30, label, sz=11, color=G400)
        _txt(sl, x+130, 778, 125, 30, val, sz=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def slide_22_prototype1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "시제품 개발 현황 (1)")
    _txt(sl, 80, 200, 800, 30, "Kira — 공고 검색 + AI 분석", sz=18, bold=True, color=DARK)

    # Screenshot placeholder
    _rrect(sl, 80, 250, 860, 480, fill=G100, stroke=BORDER, r=0.04)
    _txt(sl, 80, 440, 860, 60, "[채팅 UI 스크린샷 삽입]",
         sz=18, color=LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Feature cards
    features = [
        ("🔍 공고 검색", "나라장터 실시간 검색\n키워드·업무구분·지역·금액 필터\n첨부파일 자동 다운로드"),
        ("📊 AI 분석 (30초)", "자격요건 자동 추출\n멀티패스 병렬 처리\n동의어 사전 17개 카테고리"),
        ("📋 RFP 3섹션 요약", "사업개요·핵심요건·평가기준\n마크다운 렌더링\nreact-markdown + remark-gfm"),
        ("✅ GO/NO-GO 판단", "규칙 기반 우선 + LLM 보조\n자격요건별 개별 판정\n종합 참여 권고"),
    ]
    for i, (t, d) in enumerate(features):
        y = 250 + i * 120
        _rrect(sl, 980, y, 860, 110, fill=WHITE, stroke=BORDER, r=0.04)
        _txt(sl, 1000, y+10, 820, 28, t, sz=14, bold=True, color=DARK)
        _txt(sl, 1000, y+42, 820, 60, d, sz=11, color=BODY)


def slide_23_prototype2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "시제품 개발 현황 (2)")
    _txt(sl, 80, 200, 800, 30, "Kira — 제안서·PPT·WBS 자동 생성", sz=18, bold=True, color=DARK)

    _rrect(sl, 80, 250, 860, 480, fill=G100, stroke=BORDER, r=0.04)
    _txt(sl, 80, 440, 860, 60, "[문서 생성 결과 스크린샷 삽입]",
         sz=18, color=LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    features = [
        ("📝 제안서 DOCX", "Layer 1(495유닛) + Layer 2(회사맞춤)\nmistune 3.x AST → python-docx\n100+ 페이지 전문 구조"),
        ("📊 PPT 발표자료", "KRDS 디자인 가이드 적용\npython-pptx 자동 생성\n예상질문 10개 포함"),
        ("📅 WBS/수행계획서", "방법론 자동 감지 + LLM 태스크\nopenpyxl XLSX + matplotlib 간트\nDOCX 보고서 동시 생성"),
        ("📄 실적·경력 기술서", "CompanyDB 기반 자동 매칭\n유사 프로젝트 검색 + 서술 생성\nDOCX 표+서술 조립"),
    ]
    for i, (t, d) in enumerate(features):
        y = 250 + i * 120
        _rrect(sl, 980, y, 860, 110, fill=WHITE, stroke=BORDER, r=0.04)
        _txt(sl, 1000, y+10, 820, 28, t, sz=14, bold=True, color=DARK)
        _txt(sl, 1000, y+42, 820, 60, d, sz=11, color=BODY)


def slide_24_prototype3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "시제품 개발 현황 (3)")
    _txt(sl, 80, 200, 800, 30, "Kira — 3계층 학습 엔진 + 알림 시스템", sz=18, bold=True, color=DARK)

    # Layer diagram
    layers = [
        ("Layer 1: 범용 지식", "유튜브 40편 + 블로그 47편 + 공식문서 18편\n495개 구조화된 지식 유닛 (ChromaDB)\n\"전문 컨설턴트 50명분 노하우\"", BLUE, BLUE_50),
        ("Layer 2: 회사 맞춤", "과거 제안서 → 문체·구조·강점 분석\n수정 diff → 패턴 감지 (3회+ → 자동 반영)\n\"쓸수록 우리 회사 색깔이 나오는 AI\"", GREEN, GREEN_L),
        ("Layer 3: 승패 분석", "낙찰 vs 탈락 비교 → 패턴 추출 [Phase 3]\n발주처별 선호도 프로파일링\n\"데이터 플라이휠\" 경쟁사 진입 장벽", ORANGE, ORANGE_L),
    ]
    for i, (title, desc, accent, bg) in enumerate(layers):
        y = 250 + i * 145
        _rrect(sl, 80, y, 900, 130, fill=bg, stroke=accent, sw=Pt(2), r=0.04)
        _badge(sl, 100, y+12, title.split(":")[0], bg=accent, w=120)
        _txt(sl, 100, y+46, 860, 24, title.split(": ")[1], sz=15, bold=True, color=DARK)
        _txt(sl, 100, y+72, 860, 50, desc, sz=11, color=BODY)

    # Additional features
    extras = [
        ("🔔 맞춤 공고 알림", "키워드·금액·지역 필터\n이메일 자동 발송\n실시간 신규 공고 감지"),
        ("📋 제출 체크리스트", "RFP에서 필수 서류 자동 추출\n진행 상태 체크 기능\n누락 방지"),
        ("💬 대화형 Q&A", "RAG 하이브리드 검색\n(BM25 + 벡터 RRF)\n참조 페이지 표시"),
        ("🎯 프로덕트 투어", "신규 사용자 온보딩\ndriver.js 기반\n단계별 기능 안내"),
    ]
    for i, (t, d) in enumerate(extras):
        y = 250 + i * 120
        _rrect(sl, 1020, y, 830, 110, fill=WHITE, stroke=BORDER, r=0.04)
        _txt(sl, 1040, y+10, 790, 28, t, sz=14, bold=True, color=DARK)
        _txt(sl, 1040, y+42, 790, 60, d, sz=11, color=BODY)


def slide_25_monthly(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "창업동아리 월별 주요 활동 및 예산 사용 계획")

    activity = [
        ["월", "주요 활동", "마일스톤"],
        ["3월", "팀 빌딩 + 사업계획서 완성 + 시제품 데모 준비", "사업계획서 제출"],
        ["4월", "베타 테스트 (초대 10개사) + 피드백 수집 + 버그 수정", "베타 런칭"],
        ["5월", "베타 확대 (50개사) + 콘텐츠 마케팅 시작 + 고객 인터뷰", "50개사 돌파"],
        ["6월", "FREE→PRO 전환 퍼널 최적화 + 파트너십 협의", "첫 유료 전환"],
        ["7월", "Layer 2 회사 스타일 연동 + Enterprise 영업 시작", "Enterprise 파일럿"],
        ["8월", "성과 분석 + 투자 유치 준비 + 정식 런칭 준비", "데모데이 발표"],
    ]
    _add_table(sl, 80, 210, 1760, 330, activity, col_widths=[120, 900, 740])

    budget = [
        ["항목", "3월", "4월", "5월", "6월", "7월", "8월", "합계"],
        ["클라우드 인프라", "20", "30", "40", "50", "60", "70", "270만원"],
        ["LLM API 비용", "10", "20", "30", "40", "50", "60", "210만원"],
        ["마케팅/콘텐츠", "—", "10", "30", "50", "50", "60", "200만원"],
        ["도구/라이선스", "10", "10", "10", "10", "10", "10", "60만원"],
        ["기타 (교통/식비)", "10", "10", "10", "10", "10", "10", "60만원"],
        ["월 합계", "50", "80", "120", "160", "180", "210", "800만원"],
    ]
    _txt(sl, 80, 560, 400, 28, "예산 사용 계획 (단위: 만원)", sz=16, bold=True, color=DARK)
    _add_table(sl, 80, 595, 1760, 280, budget)


# ═════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_motivation(prs)
    slide_03_social_venture(prs)
    slide_04_product(prs)
    slide_05_market(prs)
    slide_06_competitor(prs)
    slide_07_persona(prs)
    slide_08_customer_dist(prs)
    slide_09_bm1(prs)
    slide_10_bm2(prs)
    slide_11_bm3(prs)
    slide_12_bm4(prs)
    slide_13_bm5(prs)
    slide_14_social_bm(prs)
    slide_15_validation(prs)
    slide_16_financials(prs)
    slide_17_kpi(prs)
    slide_18_team(prs)
    slide_19_progress(prs)
    slide_20_vision(prs)
    slide_21_marketing(prs)
    slide_22_prototype1(prs)
    slide_23_prototype2(prs)
    slide_24_prototype3(prs)
    slide_25_monthly(prs)

    out = "/Users/min-kyungwook/Downloads/Kira_사업계획서_학교제출.pptx"
    prs.save(out)
    print(f"✅ Saved: {out}")
    print(f"   Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
