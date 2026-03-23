"""
Generate Kira Bot business plan PPTX from the 한성대 창업동아리 template.
"""
import copy
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = Path("/Users/min-kyungwook/Downloads/(첨부2) 2026-1학기 창업동아리 사업계획서.pptx")
OUTPUT = Path("/Users/min-kyungwook/Downloads/Kira_Bot_창업동아리_사업계획서.pptx")


# ── helpers ────────────────────────────────────────────────────
def _clear_and_set(tf, lines, font_size=None, bold=None, alignment=None):
    """Clear a text frame and set new lines, preserving the first paragraph's formatting."""
    # Save format from existing first paragraph/run
    ref_para = tf.paragraphs[0] if tf.paragraphs else None
    ref_run = ref_para.runs[0] if (ref_para and ref_para.runs) else None

    # Clear all paragraphs except first
    for _ in range(len(tf.paragraphs) - 1):
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
            # clear existing runs
            for r in list(para.runs):
                r._r.getparent().remove(r._r)
        else:
            para = tf.add_paragraph()

        run = para.add_run()
        run.text = line

        # Copy formatting
        if ref_run:
            if ref_run.font.name:
                run.font.name = ref_run.font.name
            if ref_run.font.size and not font_size:
                run.font.size = ref_run.font.size
            if ref_run.font.bold is not None and bold is None:
                run.font.bold = ref_run.font.bold
            try:
                if ref_run.font.color and ref_run.font.color.rgb:
                    run.font.color.rgb = ref_run.font.color.rgb
            except AttributeError:
                pass

        if font_size:
            run.font.size = Pt(font_size)
        if bold is not None:
            run.font.bold = bold
        if alignment:
            para.alignment = alignment


def _set_shape_text(shape, lines, **kwargs):
    """Set text on a shape that has a text_frame."""
    if shape.has_text_frame:
        _clear_and_set(shape.text_frame, lines, **kwargs)


def _find_shape_by_text(slide, needle):
    """Find first shape containing needle text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = shape.text_frame.text
            if needle in full:
                return shape
    return None


def _set_cell(table, row, col, text, bold=None, font_size=None):
    """Set text in a table cell."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    if tf.paragraphs:
        para = tf.paragraphs[0]
        # clear runs
        for r in list(para.runs):
            r._r.getparent().remove(r._r)
        run = para.add_run()
        run.text = text
        if bold is not None:
            run.font.bold = bold
        if font_size:
            run.font.size = Pt(font_size)


def _find_table(slide):
    """Find first table shape on slide."""
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def _replace_text_in_shape(shape, old, new):
    """Replace text in shape runs."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def _replace_all_text(slide, old, new):
    """Replace text across all shapes in a slide."""
    for shape in slide.shapes:
        _replace_text_in_shape(shape, old, new)


# ── main ───────────────────────────────────────────────────────
def main():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)

    # ================================================================
    # Slide 1: Title
    # ================================================================
    s = slides[0]
    # Replace "창업동아리명" with our name
    _replace_all_text(s, "창업동아리명", "Kira Bot")
    # Replace guide notice
    guide_shape = _find_shape_by_text(s, "본 가이드")
    if guide_shape:
        _set_shape_text(guide_shape, [
            "공공조달 입찰 풀 라이프사이클 자동화 AI 플랫폼",
            "MS Solutions | 2026년 3월",
        ], font_size=14)

    # ================================================================
    # Slide 2: 창업의 계기
    # ================================================================
    s = slides[1]
    # Find and replace subtitle boxes
    shapes_with_text = []
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            shapes_with_text.append((t[:40], shape))

    # Replace "창업 아이디어 발상의 계기" section
    sh = _find_shape_by_text(s, "아이디어 발상의 계기")
    if sh:
        _replace_text_in_shape(sh, "창업 아이디어 발상의 계기", "문제 발견")

    sh = _find_shape_by_text(s, "구체화 과정 1")
    if sh:
        _replace_text_in_shape(sh, "창업 아이디어 구체화 과정 1", "아이디어 구체화")

    sh = _find_shape_by_text(s, "구체화 과정 2")
    if sh:
        _replace_text_in_shape(sh, "창업 아이디어 구체화 과정 2", "제품 개발")

    # Replace "텍스트 작성.." placeholders
    for shape in s.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == "텍스트 작성..":
                        break
                else:
                    continue
                break
            else:
                continue
            # Found a text box with placeholder
            txt = shape.text_frame.text
            if txt.count("텍스트 작성..") >= 1:
                # Check position to determine which box
                if shape.top < Emu(3000000):  # top section
                    _set_shape_text(shape, [
                        "공공조달 시장 225조원. 50만 중소기업이 매일",
                        "수백 건의 공고를 수작업으로 분석하고,",
                        "제안서 1건에 300~1,000만원을 외주하는 현실.",
                        "\"이 비효율을 AI로 해결할 수 있겠다\"는 확신.",
                    ], font_size=11)
                elif shape.left < Emu(5000000):  # middle-left
                    _set_shape_text(shape, [
                        "나라장터 API + LLM 결합으로",
                        "공고 검색 → RFP 분석 → GO/NO-GO를",
                        "30초 만에 자동 판단하는 MVP 개발.",
                        "HWP/PDF 파싱 + 동의어 사전 17개 카테고리 구축.",
                    ], font_size=11)
                else:
                    _set_shape_text(shape, [
                        "분석을 넘어 제안서·PPT·WBS까지",
                        "원클릭 자동 생성으로 확장.",
                        "3계층 학습 엔진(범용+회사맞춤+승패분석)",
                        "구축으로 \"쓸수록 좋아지는 AI\" 완성.",
                    ], font_size=11)

    # ================================================================
    # Slide 3: 소셜벤처 계기 → 문제 정의로 변환
    # ================================================================
    s = slides[2]
    sh = _find_shape_by_text(s, "창업의 계기(소셜 벤처)")
    if sh:
        _replace_text_in_shape(sh, "창업의 계기(소셜 벤처)", "풀고 있는 문제")

    sh = _find_shape_by_text(s, "사회적 문제 인식 계기")
    if sh:
        _replace_text_in_shape(sh, "사회적 문제 인식 계기", "중소기업 BD팀의 고통")

    sh = _find_shape_by_text(s, "사회적 문제 파악")
    if sh:
        _replace_text_in_shape(sh, "사회적 문제 파악", "비효율의 구조")

    sh = _find_shape_by_text(s, "데이터를 통한 사회문제 파악")
    if sh:
        _replace_text_in_shape(sh, "데이터를 통한 사회문제 파악", "데이터로 본 시장 규모")

    # Replace content boxes
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t == "텍스트 작성..":
                if shape.top < Emu(3000000):
                    _set_shape_text(shape, [
                        "BD 담당자 1~3명이 공고 발견, RFP 분석,",
                        "제안서 작성을 모두 수작업으로 수행.",
                        "1건에 올인하면 다른 기회를 놓치고,",
                        "동시 진행하면 품질이 떨어지는 딜레마.",
                    ], font_size=11)
                elif shape.left < Emu(5000000):
                    _set_shape_text(shape, [
                        "공고 분석: 2~4시간/건",
                        "제안서 작성: 2~4주/건",
                        "외주 비용: 300~1,000만원/건",
                        "서류 누락 시 감점/실격 → 취소 불가",
                    ], font_size=11)
                else:
                    _set_shape_text(shape, [
                        "공공조달 시장: 225조원/년 (2024, 조달청)",
                        "등록 사업자: 50만 개 (96.7% 중소기업)",
                        "나라장터 거래: 145조원 (전년 대비 11.1%↑)",
                        "2026년 발주계획: 85.6조원 (역대 최대)",
                    ], font_size=11)

            if "(장기적, 구조적, 다수의 피해 사실)" in t:
                _set_shape_text(shape, [""], font_size=10)
            if "(숫자)" in t:
                _set_shape_text(shape, [""], font_size=10)

    # ================================================================
    # Slide 4: 창업아이템 설명 (도입 전/후)
    # ================================================================
    s = slides[3]
    # Replace before/after content
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if "기존 솔루션의 한계" in t:
                if shape.left < Emu(5000000):  # left = 도입 전
                    _set_shape_text(shape, [
                        "공고 검색만 가능 (비드프로)",
                        "AI 분석만 가능, 생성 없음 (클라이원트)",
                        "HWP 파일 읽기 불가 (범용 AI)",
                        "학습 기능 없음 → 매번 처음부터",
                    ], font_size=10)
                else:
                    _set_shape_text(shape, [
                        "검색→분석→생성→학습 End-to-End 자동화",
                        "제안서·PPT·WBS·실적기술서 원클릭 생성",
                        "HWP/PDF/DOCX 완벽 파싱",
                        "3계층 학습으로 쓸수록 품질 향상",
                    ], font_size=10)
            elif "기존 솔루션 보다" in t:
                if shape.left < Emu(5000000):
                    _set_shape_text(shape, [
                        "제안서 외주: 300~1,000만원/건",
                        "분석 시간: 2~4시간/건",
                        "작성 기간: 2~4주/건",
                    ], font_size=10)
                else:
                    _set_shape_text(shape, [
                        "월 9.9만원으로 무제한 분석+5건 생성",
                        "분석 시간: 30초/건 (99.7% 절감)",
                        "제안서 검토: 1시간 (95% 절감)",
                    ], font_size=10)

    # ================================================================
    # Slide 5: 시장 분석
    # ================================================================
    s = slides[4]
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            # TAM
            if "전체 시장" in t and ":" in t:
                _set_shape_text(shape, [
                    "전체 시장(TAM): 공공조달 225조원, 등록 사업자 50만 개",
                ], font_size=14)
            # SAM
            elif "유효 시장" in t and ":" in t:
                _set_shape_text(shape, [
                    "유효 시장(SAM): IT/SW 활발 입찰 2만 기업 × 360만원 = 720억원/년",
                ], font_size=14)
            # SOM
            elif "핵심 시장" in t and ":" in t:
                _set_shape_text(shape, [
                    "핵심 시장(SOM): 3년 내 1,000개사 × 360만원 = 36억원/년",
                ], font_size=14)
            elif "시장규모" in t:
                _set_shape_text(shape, [
                    "시장규모",
                ], font_size=12, bold=True)
            elif "*최대한 숫자로 표현" in t:
                _set_shape_text(shape, [
                    "2024년 공공조달 225조원 (역대 최대, 전년 대비 7.9%↑)",
                    "나라장터 비중 64.5% (145조원), 중소기업 비중 63.1% (142조원)",
                    "출처: 조달청 2024년 공공조달 규모 보도자료",
                ], font_size=10)

    # ================================================================
    # Slide 6: 경쟁사 분석 + 포지셔닝 맵
    # ================================================================
    s = slides[5]
    # Fill competition table
    table = _find_table(s)
    if table:
        data = [
            ["분류", "비드프로", "클라이원트", "디마툴즈", "Kira Bot"],
            ["공고 검색", "O", "O", "X", "O"],
            ["AI 분석", "X", "O", "X", "O"],
            ["문서 생성", "X", "X", "X", "O"],
            ["학습 기능", "X", "X", "X", "O"],
        ]
        for r, row_data in enumerate(data):
            for c, val in enumerate(row_data):
                _set_cell(table, r, c, val, bold=(r == 0 or c == 0), font_size=10)

    # Update positioning map labels
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t == "핵심가치1":
                _replace_text_in_shape(shape, "핵심가치1", "문서 생성 능력")
            elif t == "핵심가치2":
                _replace_text_in_shape(shape, "핵심가치2", "AI 분석 깊이")
            elif "창업" in t and "회사" in t and len(t) < 10:
                _set_shape_text(shape, ["Kira Bot"], font_size=14, bold=True)

    # ================================================================
    # Slide 7: 고객 페르소나
    # ================================================================
    s = slides[6]
    table = _find_table(s)
    if table:
        persona_data = [
            ["인적사항", "니즈"],
            [
                "이   름: 김 영수\n나   이: 만 42세\n직   업: IT 서비스 기업 BD팀장\n회사규모: 30명",
                "고통: 매주 3~5건 공고 분석에 하루 4시간 소요.\n제안서 외주 시 건당 500만원 부담.\n\n필요한 것: 빠른 GO/NO-GO 판단 + 제안서 초안 자동 생성"
            ],
            ["활동과 경험", "그/그녀가 할 것 같은 말"],
            [
                "경력: SI 업계 15년 근무\n연 입찰: 10~20건 참여\n도구: 비드프로(검색), 엑셀(관리)\n고민: 동시 다건 진행 시 품질 저하",
                "\"공고 분석하느라 하루가 다 가고,\n정작 제안서 쓸 시간이 없어.\nAI가 초안이라도 잡아주면 좋겠다.\""
            ],
        ]
        for r, row_data in enumerate(persona_data):
            for c, val in enumerate(row_data):
                _set_cell(table, r, c, val, bold=(r == 0 or r == 2), font_size=9)

    # ================================================================
    # Slide 8: 고객 분포 현황
    # ================================================================
    s = slides[7]
    segments = [
        ("핵심고객\n특징1", "1차: IT 서비스 중소기업\n(10~100명, 연 5건+ 입찰)"),
        ("핵심고객\n특징2", "2차: 건설·엔지니어링·컨설팅\n(입찰 빈도 높고 제안서 부담 큼)"),
        ("핵심고객\n특징3", "3차: 정부지원사업 신청자\n(TIPS, AI바우처, 창업패키지)"),
    ]
    seg_idx = 0
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if "핵심고객" in t and "특징" in t and seg_idx < len(segments):
                _set_shape_text(shape, [segments[seg_idx][1]], font_size=12)
                seg_idx += 1

    # ================================================================
    # Slide 9: 비즈니스 모델 (replace example with ours)
    # ================================================================
    s = slides[8]
    _replace_all_text(s, "비즈니스 모델 1", "비즈니스 모델")
    _replace_all_text(s, "오픈갤러리", "Kira Bot")

    # Clear all content and replace with our BM
    # Find all text shapes and update
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t == "비즈니스 모델":
                continue  # keep title
            if t == "Kira Bot":
                continue  # keep subtitle
            if "작가" in t:
                _set_shape_text(shape, ["중소기업\n(고객)"], font_size=14, bold=True)
            elif "작품 대여" == t:
                _set_shape_text(shape, ["월 구독료 9.9만원"], font_size=12, bold=True)
            elif "대여료" == t and shape.left > Emu(5000000):
                _set_shape_text(shape, ["AI 분석+생성 서비스"], font_size=12, bold=True)
            elif "판매 수수료" in t:
                _set_shape_text(shape, ["엔터프라이즈 계약"], font_size=12, bold=True)
            elif "작품 판매 중개" in t:
                _set_shape_text(shape, ["전담 학습 모델 + SLA"], font_size=12, bold=True)
            elif "작품 대여/ 배송/ 교체" in t:
                _set_shape_text(shape, ["Free(5건)→PRO 전환"], font_size=12, bold=True)
            elif t == "대여료":
                _set_shape_text(shape, ["나라장터\nAPI"], font_size=12, bold=True)
            elif "마케팅 서비스 기획" in t:
                _set_shape_text(shape, ["공고 자동 수집"], font_size=12, bold=True)
            elif "기획비" in t:
                _set_shape_text(shape, ["맞춤 알림 발송"], font_size=12, bold=True)
            elif "작품 배송" in t:
                _set_shape_text(shape, ["문서 DOCX/PPTX/XLSX"], font_size=12, bold=True)
            elif "작품" in t and "대여" in t:
                pass
            elif "작품" in t and "구매" in t:
                _set_shape_text(shape, [""], font_size=10)
            elif "의뢰" in t:
                _set_shape_text(shape, ["OpenAI\nAPI"], font_size=14, bold=True)
            elif "전시" in t and "기획" in t:
                _set_shape_text(shape, ["LLM 분석+생성 비용"], font_size=12, bold=True)

    # ================================================================
    # Slides 10-14: Delete BM examples (reverse order)
    # ================================================================
    def _delete_slide(prs, slide_index):
        """Delete a slide by index using XML manipulation."""
        rId = prs.slides._sldIdLst[slide_index].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[slide_index]
        prs.slides._sldIdLst.remove(sldId)

    for idx in [13, 12, 11, 10, 9]:
        _delete_slide(prs, idx)

    # Refresh slides list after deletion
    slides = list(prs.slides)
    # Now indices shifted: what was 15 is now 10, etc.

    # ================================================================
    # Slide 10 (was 15): 고객 설문조사
    # ================================================================
    s = slides[9]
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if "조사기간:" in t:
                _set_shape_text(shape, [
                    "조사기간: 2026년 3~4월 (베타 테스트 기간)",
                    "조사대상자(명): 50개 IT 서비스 중소기업 BD 담당자",
                    "",
                    "주요 조사 항목:",
                    "1. 현재 공고 분석에 소요되는 시간",
                    "2. 제안서 작성 외주 비용 및 빈도",
                    "3. AI 제안서 생성 도구 사용 의향",
                    "4. 적정 구독 가격 수용도",
                    "5. 핵심 요구 기능 우선순위",
                ], font_size=11)

    # ================================================================
    # Slide 11 (was 16): 소요자금 및 예상 매출
    # ================================================================
    s = slides[10]
    table = _find_table(s)
    if table:
        data = [
            ["소요 자금", "예상 매출"],
            ["제품 개발: 1.0억", "Y1 (2026)"],
            ["마케팅/GTM: 0.5억", "총 비용: 3.2억"],
            ["운영 자금: 0.5억", "매출: 1.2억"],
            ["총 소요: 2.0억", "Y3 매출: 23억 (BEP: Y2)"],
            ["LLM API: 0.2억/년", "영업이익률: Y3 49.6%"],
        ]
        for r, row_data in enumerate(data):
            for c, val in enumerate(row_data):
                _set_cell(table, r, c, val, bold=(r == 0), font_size=10)

    # ================================================================
    # Slide 12 (was 17): 핵심 성과지표 및 재무 목표
    # ================================================================
    s = slides[11]
    sh = _find_shape_by_text(s, "작성이 어렵다면")
    if sh:
        _set_shape_text(sh, [
            "핵심 KPI:",
            "MAU: 500명 (Y1 목표)",
            "Free→PRO 전환율: 15%",
            "MRR: 1,000만원 (Y1 목표)",
            "Churn Rate: < 5%/월",
            "수정률(Edit Rate): 45% → 15% (6개월)",
            "제안서 생성 건수: 200건/월",
            "NPS: > 40",
            "",
            "고유 KPI — 수정률(Edit Rate):",
            "AI 초안 대비 사용자 수정 비율. 낮아질수록 학습 완성도 높음.",
            "v1: 45% → v5: 20% → v10: 8% = \"학습도 92%\"",
        ], font_size=11)

    # ================================================================
    # Slide 13 (was 18): 팀 소개
    # ================================================================
    s = slides[12]
    team_replacements = [
        {
            "name_old": "이름",
            "position_old": "Team Leader",
            "desc_old": "짧은 설명 작성",
            "name_new": "민경욱",
            "position_new": "대표 / 기술 총괄",
            "desc_new": "풀스택 개발, AI/ML, 공공조달 도메인",
        },
    ]
    name_count = 0
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t == "이름" and name_count == 0:
                _set_shape_text(shape, ["민경욱"], bold=True)
                name_count += 1
            elif t == "이름" and name_count == 1:
                _set_shape_text(shape, ["채용 예정"], bold=True)
                name_count += 1
            elif t == "이름" and name_count == 2:
                _set_shape_text(shape, ["채용 예정"], bold=True)
                name_count += 1
            elif t == "이름" and name_count == 3:
                _set_shape_text(shape, ["채용 예정"], bold=True)
                name_count += 1
            elif t == "Team Leader":
                _set_shape_text(shape, ["대표 / 기술 총괄"])
            elif "짧은 설명 작성" in t:
                if shape.left < Emu(3500000):
                    _set_shape_text(shape, ["풀스택 개발, AI/ML\n공공조달 도메인 전문"], font_size=11)
                elif shape.left < Emu(6000000):
                    _set_shape_text(shape, ["백엔드 개발자\nRAG 엔진·SaaS 인프라"], font_size=11)
                elif shape.left < Emu(8000000):
                    _set_shape_text(shape, ["프론트엔드 개발자\nChat UI·대시보드"], font_size=11)
                else:
                    _set_shape_text(shape, ["BD / 마케팅\n고객 확보·파트너십"], font_size=11)
            elif t == "Position":
                pass  # Keep position labels

    pos_count = 0
    for shape in s.shapes:
        if shape.has_text_frame:
            if shape.text_frame.text.strip() == "Position":
                if pos_count == 0:
                    _set_shape_text(shape, ["백엔드 (Y1 Q2)"])
                elif pos_count == 1:
                    _set_shape_text(shape, ["프론트엔드 (Y1 Q3)"])
                elif pos_count == 2:
                    _set_shape_text(shape, ["BD/마케팅 (Y1 Q2)"])
                pos_count += 1

    # ================================================================
    # Slide 14 (was 19): 진행 과정
    # ================================================================
    s = slides[13]
    # This slide has a stage diagram. Mark our current stage.
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if "*어느 단계" in t:
                _set_shape_text(shape, [
                    "현재 단계: 제품·서비스 검증 (높은 수준 MVP)",
                    "154개 테스트 통과, E2E 검증 완료 (2026-02-28)",
                    "Layer 1 지식 495유닛 탑재, 5종 문서 자동 생성 동작 중",
                ], font_size=11)

    # ================================================================
    # Slide 15 (was 20): 창업회사의 비전
    # ================================================================
    s = slides[14]
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t == "VISION":
                _set_shape_text(shape, ["VISION"], font_size=28, bold=True)
            elif "회사가 나아갈 방향" in t:
                _set_shape_text(shape, [
                    "입찰 공고가 나오면, AI가 알아서 분석하고, 제안서를 쓰고,",
                    "PPT를 만들고, 부족한 서류를 알려주고, 결과를 추적하여",
                    "매번 더 똑똑해지는 — 기업의 AI BD 팀이 되는 것.",
                ], font_size=11)
            elif "현재 지향하는 가치" in t:
                _set_shape_text(shape, [
                    "\"먼저 배우고, 회사를 알고, 그다음 생성한다\"",
                ], font_size=12)
            elif t == "Value 1":
                _set_shape_text(shape, ["End-to-End\n자동화"], font_size=11)
            elif t == "Value 2":
                _set_shape_text(shape, ["3계층\n학습 엔진"], font_size=11)
            elif t == "Value 3":
                _set_shape_text(shape, ["데이터\n플라이휠"], font_size=11)
            elif t == "Value 4":
                _set_shape_text(shape, ["쓸수록\n좋아지는 AI"], font_size=11)

    # ================================================================
    # Slide 16 (was 21): 마케팅 계획
    # ================================================================
    s = slides[15]
    marketing = [
        ("마케팅 방안 1", "콘텐츠 마케팅", "\"공공입찰 GO/NO-GO 판단법\" 블로그(주2)·유튜브(월2). SEO 오가닉 유입."),
        ("마케팅 방안 2", "프리미엄 모델", "FREE(월5건) → PRO 전환. 무료 체험에서 가치 증명 후 유료 전환 15% 목표."),
        ("마케팅 방안 3", "입찰 커뮤니티", "나라장터 관련 카페·포럼에서 초기 얼리어답터 확보. 무료 분석 이벤트."),
        ("마케팅 방안 4", "파트너십", "제안서 컨설팅사와 협업 (도구 제공). 중소기업 디지털 전환 바우처 활용."),
        ("마케팅 방안 5", "직접 영업", "BD 담당자 대상 데모 + 무료 분석. 엔터프라이즈 계약 전환."),
    ]
    for old_label, new_label, desc in marketing:
        sh = _find_shape_by_text(s, old_label)
        if sh:
            _replace_text_in_shape(sh, old_label, new_label)
        # Find the description shape (": 텍스트 작성...")
    desc_idx = 0
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if ": 텍스트 작성..." in t and desc_idx < len(marketing):
                _set_shape_text(shape, [": " + marketing[desc_idx][2]], font_size=10)
                desc_idx += 1

    # ================================================================
    # Slides 17-19 (was 22-24): 시제품 개발 현황
    # ================================================================
    for slide_idx, content in [
        (16, [
            "시제품 개발 현황 — 핵심 기능 (운영 중)",
            "",
            "1. 나라장터 실시간 검색: 키워드·업무구분·지역·금액 필터",
            "2. 문서 자동 파싱: PDF, DOCX, HWP, HWPX, Excel, PPT",
            "3. AI 자격요건 추출: 멀티패스 병렬 + 동의어 사전 17개 카테고리",
            "4. GO/NO-GO 자동 판단: 규칙 기반 우선 + LLM 보조 (하이브리드)",
            "5. 제안서 DOCX 자동 생성: Layer 1+2 지식 기반, 100+ 페이지",
        ]),
        (17, [
            "시제품 개발 현황 — 문서 생성 패키지",
            "",
            "1. 기술제안서 DOCX: 평가항목별 섹션 자동, 블라인드 위반 검증",
            "2. PPT 발표자료: KRDS 디자인 가이드, 예상질문 10개+답변",
            "3. 수행계획서/WBS: XLSX 간트차트 + DOCX 보고서",
            "4. 실적·경력 기술서: 회사 DB 기반 자동 매칭",
            "5. 제출 체크리스트: RFP에서 필수 서류 자동 추출",
        ]),
        (18, [
            "시제품 개발 현황 — 기술 스택 & 검증",
            "",
            "프론트엔드: React 19 + TypeScript + Vite + Tailwind",
            "백엔드: Python FastAPI + ChromaDB RAG + OpenAI GPT-4",
            "문서 생성: python-docx, python-pptx, openpyxl, mistune 3.x",
            "보안: HMAC 서명, CSRF/SSRF 방어, 입력 검증",
            "",
            "검증 현황: 154개 테스트 통과, E2E 검증 완료",
            "Layer 1 지식: 495유닛 (7카테고리), 코드: 30,000+ 라인",
        ]),
    ]:
        s = slides[slide_idx]
        # Replace title
        title_shape = _find_shape_by_text(s, "시제품 개발 현황")
        if title_shape:
            _replace_text_in_shape(title_shape, "시제품 개발 현황", content[0])

        # Try to find a content area or add text box
        # Since these slides might be mostly empty (for screenshots), add content
        content_added = False
        for shape in s.shapes:
            if shape.has_text_frame and "시제품" not in shape.text_frame.text:
                if shape.width > Emu(3000000):  # large enough to be content area
                    _set_shape_text(shape, content[2:], font_size=12)
                    content_added = True
                    break

        if not content_added:
            from pptx.util import Inches
            txBox = s.shapes.add_textbox(
                Emu(540000), Emu(1200000),
                Emu(9600000), Emu(5500000)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content[2:]):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                run = para.add_run()
                run.text = line
                run.font.size = Pt(14)
                run.font.name = "맑은 고딕"

    # ================================================================
    # Slide 20 (was 25): 월별 주요 활동 및 예산 사용 계획
    # ================================================================
    s = slides[19]

    # Find activity table (7 rows x 9 cols)
    for shape in s.shapes:
        if shape.has_table and len(shape.table.columns) == 9:
            table = shape.table
            activities = [
                ["추진 내용", "4월", "5월", "6월", "7월", "8월", "9월", "10월", "11월"],
                ["베타 테스트 (50개사)", "●", "●", "", "", "", "", "", ""],
                ["Layer 2 회사학습 연동", "", "●", "●", "", "", "", "", ""],
                ["정식 런칭 + 결제 연동", "", "", "●", "●", "", "", "", ""],
                ["콘텐츠 마케팅 (블로그/유튜브)", "", "●", "●", "●", "●", "●", "●", "●"],
                ["엔터프라이즈 영업 시작", "", "", "", "●", "●", "●", "●", "●"],
                ["Layer 3 승패분석 개발", "", "", "", "", "●", "●", "●", "●"],
            ]
            for r, row_data in enumerate(activities):
                for c, val in enumerate(row_data):
                    _set_cell(table, r, c, val, bold=(r == 0), font_size=8)
            break

    # Find budget table (5 rows x 3 cols)
    for shape in s.shapes:
        if shape.has_table and len(shape.table.columns) == 3:
            table = shape.table
            budget = [
                ["구분", "사용 내역", "예산(만원)"],
                ["클라우드 인프라", "Railway/AWS 서버 운영", "50"],
                ["LLM API", "OpenAI API 사용료", "100"],
                ["마케팅", "콘텐츠 제작 + 커뮤니티", "100"],
                ["합계", "", "250"],
            ]
            for r, row_data in enumerate(budget):
                for c, val in enumerate(row_data):
                    _set_cell(table, r, c, val, bold=(r == 0 or r == 4), font_size=9)
            break

    # ================================================================
    # Save
    # ================================================================
    prs.save(str(OUTPUT))
    print(f"Saved to {OUTPUT}")


if __name__ == "__main__":
    main()
