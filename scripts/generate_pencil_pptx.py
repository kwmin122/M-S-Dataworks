#!/usr/bin/env python3
"""
Generate Kira 사업계획서 PPTX — Swiss Clean design from Pencil MCP.
Run: python scripts/generate_pencil_pptx.py
Output: ~/Downloads/Kira_사업계획서_v2.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Dimensions ────────────────────────────────────────────────────────
SLIDE_W = Emu(12192000)  # 16:9 widescreen
SLIDE_H = Emu(6858000)
_PX = SLIDE_W // 1920     # 1 Pencil pixel in EMU (≈6350)

def px(n):
    """Convert Pencil pixel to EMU."""
    return int(n * _PX)

# ─── Colors (Swiss Clean) ─────────────────────────────────────────────
BLUE     = RGBColor(0x25, 0x63, 0xEB)
DARK     = RGBColor(0x18, 0x18, 0x1B)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
G700     = RGBColor(0x71, 0x71, 0x7A)
G500     = RGBColor(0x52, 0x52, 0x5B)
G400     = RGBColor(0xA1, 0xA1, 0xAA)
G200     = RGBColor(0xE4, 0xE4, 0xE7)
G100     = RGBColor(0xF4, 0xF4, 0xF5)
G50      = RGBColor(0xFA, 0xFA, 0xFA)
BLUE50   = RGBColor(0xEF, 0xF6, 0xFF)
BLUE100  = RGBColor(0xDB, 0xEA, 0xFE)
BLUE400  = RGBColor(0x60, 0xA5, 0xFA)
DSURF    = RGBColor(0x27, 0x27, 0x2A)
DBLUE    = RGBColor(0x1E, 0x3A, 0x5F)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)
GREEN50  = RGBColor(0xF0, 0xFD, 0xF4)
GREEN100 = RGBColor(0xDC, 0xFC, 0xE7)
ORANGE   = RGBColor(0xEA, 0x58, 0x0C)
ORANGE50 = RGBColor(0xFF, 0xF7, 0xED)
ORANGE100= RGBColor(0xFF, 0xED, 0xD5)
RED      = RGBColor(0xEF, 0x44, 0x44)
YELLOW   = RGBColor(0xEA, 0xB3, 0x08)

FONT = '맑은 고딕'

# ─── Helpers ───────────────────────────────────────────────────────────

def _bg(slide, color):
    """Fill entire slide with color using a full-size rectangle."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _rect(slide, x, y, w, h, fill=None, stroke=None, sw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if stroke: s.line.color.rgb = stroke; s.line.width = sw
    else: s.line.fill.background()
    return s

def _rrect(slide, x, y, w, h, fill=None, stroke=None, sw=Pt(1), r=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if stroke: s.line.color.rgb = stroke; s.line.width = sw
    else: s.line.fill.background()
    try: s.adjustments[0] = r
    except: pass
    return s

def _txt(slide, x, y, w, h, text, sz=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    try: tf.vertical_anchor = anchor
    except: pass
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = FONT; p.alignment = align
    return tb

def _mtxt(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """Multi-line text. lines: [(text, sz, bold, color), ...]"""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    try: tf.vertical_anchor = anchor
    except: pass
    for i, (text, sz, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.bold = bold
        p.font.color.rgb = color; p.font.name = FONT
    return tb

def _header(slide, num, title, subtitle=None, dark=False):
    """Standard section header: number + title + subtitle + divider."""
    tc = BLUE; ttc = WHITE if dark else DARK; stc = G400 if dark else G700
    dc = DSURF if dark else G200
    _txt(slide, 120, 50, 100, 40, num, sz=16, bold=True, color=tc)
    _txt(slide, 120, 85, 700, 60, title, sz=34, bold=True, color=ttc)
    if subtitle:
        _txt(slide, 120, 145, 1000, 40, subtitle, sz=18, color=stc)
    _rect(slide, 120, 195, 1680, 1, fill=dc)


# ═══════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Cover slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    # Accent bar
    _rect(sl, 140, 300, 5, 280, fill=BLUE)
    # Title
    _txt(sl, 180, 300, 700, 80, "Kira", sz=52, bold=True, color=DARK)
    _txt(sl, 180, 380, 700, 50, "업무 자동화 AI 플랫폼", sz=24, color=G700)
    # Divider
    _rect(sl, 180, 440, 60, 2, fill=G200)
    # Meta
    _txt(sl, 180, 460, 500, 30, "M&S Solutions  ·  2026년 3월", sz=16, color=G400)
    _txt(sl, 180, 495, 500, 30, "한성대학교 창업동아리 사업계획서", sz=16, bold=True, color=G700)
    # Badges
    _rrect(sl, 180, 550, 200, 36, fill=BLUE50, r=0.4)
    _txt(sl, 190, 553, 180, 30, "✓ 154 테스트 통과", sz=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _rrect(sl, 400, 550, 200, 36, fill=BLUE50, r=0.4)
    _txt(sl, 410, 553, 180, 30, "✓ E2E 검증 완료", sz=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    # Hero area
    _rrect(sl, 1050, 160, 700, 700, fill=G100, stroke=G200, r=0.06)
    _txt(sl, 1200, 420, 400, 120, "AI\nTask\nAutomation", sz=44, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    """문제 정의."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "01", "문제 정의", "공공조달 시장 225조원, 등록 사업자 50만 개 — 96.7%가 중소기업")

    cards = [
        ("🔍 공고 발견", "하루 수백 건의 공고에서\n관련 공고 찾기에\n담당자 1명 × 2시간/일"),
        ("📄 RFP 분석", "50~300페이지 PDF/HWP\n요건이 산재되어 있어\n2~4시간/건 소요"),
        ("✍️ 제안서 작성", "50~200페이지 처음부터\n2~4주 소요\n외주 시 300~1,000만원"),
        ("⚠️ 결과 학습 부재", "낙찰/탈락 원인 분석 안 함\n같은 실수 반복\n시간+품질 트레이드오프"),
    ]
    for i, (t, d) in enumerate(cards):
        x = 120 + i * 410
        _rrect(sl, x, 230, 380, 280, fill=WHITE, stroke=G200, r=0.05)
        _txt(sl, x+24, 250, 332, 36, t, sz=18, bold=True, color=DARK)
        _txt(sl, x+24, 300, 332, 180, d, sz=13, color=G700)

    # Insight box
    _rrect(sl, 120, 540, 1680, 100, fill=BLUE50, r=0.04)
    _txt(sl, 152, 548, 200, 24, "핵심 인사이트", sz=13, bold=True, color=BLUE)
    _txt(sl, 152, 575, 1600, 50, "중소기업 BD팀은 1~3명. 한 건에 올인하면 다른 기회를 놓치고, 동시 진행하면 품질이 떨어진다.", sz=15, color=DARK)

    # Market gap
    _rrect(sl, 120, 660, 1680, 70, fill=G50, stroke=G200, r=0.04)
    _txt(sl, 152, 670, 1600, 50, '시장 공백:  "분석 + 생성 + 학습"을 하나의 플랫폼에서 제공하는 서비스가 없다', sz=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def slide_03_solution(prs):
    """솔루션: Kira Bot (dark slide)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, DARK)
    _header(sl, "02", "솔루션: Kira", "공공조달 입찰 풀 라이프사이클 자동화 AI 플랫폼", dark=True)

    steps = [
        ("🔍", "공고 발견", "자동", DSURF),
        ("📊", "AI 분석", "30초", DSURF),
        ("✅", "GO/NO-GO", "즉시", DSURF),
        ("📝", "문서 생성", "원클릭 5종", BLUE),
        ("📋", "제출 관리", "체크리스트", DSURF),
        ("🧠", "결과 학습", "자동 반영", DSURF),
    ]
    for i, (icon, name, speed, bg) in enumerate(steps):
        x = 120 + i * 280
        _rrect(sl, x, 230, 240, 230, fill=bg, r=0.06)
        _txt(sl, x+20, 250, 200, 36, icon, sz=24)
        _txt(sl, x+20, 300, 200, 30, name, sz=16, bold=True, color=WHITE)
        _txt(sl, x+20, 335, 200, 24, speed, sz=13, bold=True, color=BLUE)
        if i < 5:
            _txt(sl, x+250, 320, 30, 30, "→", sz=20, bold=True, color=G700)

    # Before / After
    _rrect(sl, 120, 500, 780, 260, fill=DSURF, stroke=RGBColor(0x3F, 0x3F, 0x46), r=0.04)
    _mtxt(sl, 150, 515, 720, 230, [
        ("Before — 기존 방식", 16, True, RED),
        ("분석: 2~4시간/건", 14, False, G400),
        ("제안서: 2~4주, 외주 300~1,000만원", 14, False, G400),
        ("학습: 없음 — 같은 실수 반복", 14, False, G400),
    ])

    _rrect(sl, 960, 500, 780, 260, fill=DBLUE, stroke=BLUE, sw=Pt(2), r=0.04)
    _mtxt(sl, 990, 515, 720, 230, [
        ("After — Kira", 16, True, BLUE),
        ("분석: 30초 (AI 자동)", 14, False, WHITE),
        ("제안서: 원클릭 생성 (5종 문서 패키지)", 14, False, WHITE),
        ("학습: 수정할수록 AI가 자동 학습", 14, False, WHITE),
        ("→ ROI 20x (월 9.9만원으로 200만원+ 절감)", 14, True, BLUE400),
    ])


def slide_04_learning(prs):
    """3계층 학습 엔진."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "03", "3계층 학습 엔진", '"먼저 배우고, 회사를 알고, 그다음 생성한다" — 쓸수록 좋아지는 AI')

    layers = [
        ("Layer 1", "범용 지식", "모든 고객 공통", BLUE, BLUE50, [
            "· 유튜브 40편 + 블로그 47편 + 공식문서 18편",
            "· 495개 구조화된 지식 유닛 (ChromaDB)",
            "· 평가위원 심리, 감점 회피, 배점 전략",
            '"전문 컨설턴트 50명분 노하우"  구축 비용 $10 미만',
        ]),
        ("Layer 2", "회사 맞춤", "고객별 개별 학습", GREEN, GREEN50, [
            "· 과거 제안서 → 문체·구조·강점 자동 분석",
            "· 수정 diff → 패턴 감지 (3회+ → 자동 반영)",
            "· RAG 기반 — 재학습 불필요, 즉시 갱신",
            '"쓸수록 우리 회사 색깔이 나오는 AI"',
        ]),
        ("Layer 3", "승패 분석", "네트워크 효과 (Phase 3)", ORANGE, ORANGE50, [
            "· 낙찰 vs 탈락 비교 → 승패 패턴 추출",
            "· 발주처별 선호도 프로파일링",
            "· 사용자 증가 → 데이터 축적 → 정확도 향상",
            '"데이터 플라이휠"  경쟁사 진입 장벽',
        ]),
    ]
    for i, (tag, title, sub, accent, bg, details) in enumerate(layers):
        x = 120 + i * 560
        _rrect(sl, x, 220, 520, 460, fill=bg, r=0.05)
        # Badge
        _rrect(sl, x+24, 240, 100, 30, fill=accent, r=0.4)
        _txt(sl, x+24, 242, 100, 26, tag, sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        _txt(sl, x+24, 285, 472, 36, title, sz=22, bold=True, color=DARK)
        _txt(sl, x+24, 318, 472, 24, sub, sz=13, color=G700)
        # Details
        for j, line in enumerate(details):
            is_quote = line.startswith('"') or line.startswith('"')
            _txt(sl, x+24, 360 + j*32, 472, 28, line, sz=12,
                 bold=is_quote, color=accent if is_quote else G500)

    # Flywheel bar
    _rrect(sl, 120, 710, 1680, 60, fill=G50, stroke=G200, r=0.04)
    flow = "Layer 1 범용 지식  →  Layer 2 회사 맞춤  →  Layer 3 승패 분석  →  데이터 플라이휠"
    _txt(sl, 120, 720, 1680, 40, flow, sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def slide_05_market(prs):
    """시장 분석 — TAM/SAM/SOM."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "04", "시장 분석", "대한민국 공공조달 시장 — IT 서비스·소프트웨어 분야")

    # TAM
    _rrect(sl, 120, 220, 520, 420, fill=WHITE, stroke=G200, r=0.05)
    _txt(sl, 156, 240, 200, 24, "TAM", sz=14, bold=True, color=BLUE)
    _txt(sl, 156, 270, 400, 30, "전체 시장", sz=18, bold=True, color=DARK)
    _txt(sl, 156, 310, 400, 60, "225조원", sz=44, bold=True, color=DARK)
    _txt(sl, 156, 375, 200, 24, "/년", sz=16, color=G700)
    _txt(sl, 156, 410, 400, 24, "등록 사업자 50만 개", sz=14, color=G700)
    _txt(sl, 156, 438, 400, 24, "IT 서비스 분야 약 8만 개", sz=14, color=G700)

    # SAM
    _rrect(sl, 700, 220, 520, 420, fill=BLUE50, stroke=BLUE100, r=0.05)
    _txt(sl, 736, 240, 200, 24, "SAM", sz=14, bold=True, color=BLUE)
    _txt(sl, 736, 270, 400, 30, "접근 가능 시장", sz=18, bold=True, color=DARK)
    _txt(sl, 736, 310, 400, 60, "720억원", sz=44, bold=True, color=BLUE)
    _txt(sl, 736, 375, 200, 24, "/년", sz=16, color=G700)
    _txt(sl, 736, 410, 400, 24, "활발한 입찰 기업 약 2만 개", sz=14, color=G700)
    _txt(sl, 736, 438, 400, 24, "평균 월 30만원 구독 가정", sz=14, color=G700)

    # SOM
    _rrect(sl, 1280, 220, 520, 420, fill=DARK, r=0.05)
    _txt(sl, 1316, 240, 200, 24, "SOM (3년)", sz=14, bold=True, color=BLUE400)
    _txt(sl, 1316, 270, 400, 30, "획득 가능 시장", sz=18, bold=True, color=WHITE)
    _txt(sl, 1316, 310, 400, 60, "36억원", sz=44, bold=True, color=WHITE)
    _txt(sl, 1316, 375, 300, 24, "/년 (점유율 5%)", sz=16, color=G400)
    _txt(sl, 1316, 410, 400, 24, "타겟: IT 중소기업 (10~100명)", sz=14, color=G400)
    _txt(sl, 1316, 438, 400, 24, "목표 1,000개사 × 360만원/년", sz=14, color=G400)


def slide_06_competition(prs):
    """경쟁 분석."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "05", "경쟁 분석")

    # Table header
    cols = ["경쟁사", "공고 검색", "AI 분석", "GO/NO-GO", "제안서 생성", "PPT/WBS", "학습 엔진"]
    _rrect(sl, 120, 210, 1680, 44, fill=G100, r=0.03)
    for j, col in enumerate(cols):
        w = 240 if j == 0 else 240
        _txt(sl, 120 + j*240, 216, w, 32, col, sz=12, bold=True, color=G700,
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    # Rows
    rows = [
        ("비드프로",    ["●", "—", "—", "—", "—", "—"]),
        ("클라이원트",   ["●", "●", "—", "—", "—", "—"]),
        ("디마툴즈",    ["—", "—", "△", "—", "—", "—"]),
    ]
    for i, (name, vals) in enumerate(rows):
        y = 264 + i * 52
        _txt(sl, 120, y+6, 240, 28, name, sz=13, bold=True, color=DARK)
        for j, v in enumerate(vals):
            clr = GREEN if v == "●" else (YELLOW if v == "△" else G200)
            _txt(sl, 360 + j*240, y+6, 240, 28, v, sz=14, color=clr, align=PP_ALIGN.CENTER)
        _rect(sl, 120, y+44, 1680, 1, fill=G100)

    # Kira Bot row (highlighted)
    y_kira = 264 + 3 * 52
    _rrect(sl, 120, y_kira, 1680, 44, fill=BLUE50, stroke=BLUE, sw=Pt(2), r=0.03)
    _txt(sl, 120, y_kira+8, 240, 28, "Kira", sz=13, bold=True, color=BLUE)
    for j in range(6):
        _txt(sl, 360 + j*240, y_kira+8, 240, 28, "●", sz=14, color=GREEN, align=PP_ALIGN.CENTER)

    # Position map
    _rrect(sl, 120, 500, 800, 260, fill=G50, stroke=G200, r=0.05)
    _mtxt(sl, 150, 515, 740, 230, [
        ("경쟁 포지션 맵", 18, True, DARK),
        ("", 8, False, DARK),
        ("생성 능력 (높음) ↑", 12, True, G700),
        ("  ★ Kira — 분석 + 생성 + 학습", 15, True, BLUE),
        ("분석 깊이 ──────── 범위 (넓음)", 11, False, G400),
        ("  ● 클라이원트 (분석만)     ● 비드프로 (검색만)", 13, False, G700),
        ("생성 능력 (없음) ↓", 12, True, G700),
    ])

    # Trends
    _rrect(sl, 980, 500, 820, 260, fill=WHITE, stroke=G200, r=0.05)
    _txt(sl, 1010, 515, 700, 30, "시장 트렌드", sz=18, bold=True, color=DARK)
    trends = [
        "1. 공공조달 디지털 전환 가속 — 전자입찰 100%",
        "2. AI 업무 자동화 수용도 급증",
        "3. 중소기업 생산성 SaaS 도구 수요 증가",
        "4. 정부지원사업 증가 (TIPS, AI 바우처 등)",
    ]
    for i, t in enumerate(trends):
        _txt(sl, 1010, 555 + i*42, 760, 36, t, sz=13, color=G500)


def slide_07_business_model(prs):
    """비즈니스 모델 — 3-Tier Pricing."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "06", "비즈니스 모델", "B2B SaaS 구독 모델 — 3-Tier Pricing")

    # FREE
    _rrect(sl, 120, 220, 500, 440, fill=WHITE, stroke=G200, r=0.05)
    _mtxt(sl, 156, 240, 428, 400, [
        ("FREE", 13, True, G700),
        ("₩0", 40, True, DARK),
        ("리드 생성, 체험용", 13, False, G700),
        ("", 8, False, DARK),
        ("✓  월 5건 분석", 13, False, G500),
        ("✓  공고 검색", 13, False, G500),
        ("✓  RFP 요약", 13, False, G500),
        ("✗  제안서 생성", 13, False, G200),
        ("✗  학습 엔진", 13, False, G200),
    ])

    # PRO (highlighted)
    _rrect(sl, 680, 220, 540, 440, fill=DARK, stroke=BLUE, sw=Pt(2), r=0.05)
    _rrect(sl, 716, 238, 80, 24, fill=BLUE, r=0.4)
    _txt(sl, 716, 239, 80, 22, "BEST", sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _mtxt(sl, 716, 268, 468, 380, [
        ("PRO", 13, True, BLUE400),
        ("₩99,000", 40, True, WHITE),
        ("/월 — 소기업 (5~30명)", 13, False, G400),
        ("", 8, False, DARK),
        ("✓  무제한 분석", 13, False, WHITE),
        ("✓  제안서 5건/월 생성", 13, False, WHITE),
        ("✓  PPT·WBS·실적기술서", 13, False, WHITE),
        ("✓  맞춤 공고 알림", 13, False, WHITE),
        ("✓  수정 학습 (Layer 2)", 13, False, WHITE),
    ])

    # ENTERPRISE
    _rrect(sl, 1280, 220, 500, 440, fill=WHITE, stroke=G200, r=0.05)
    _mtxt(sl, 1316, 240, 428, 400, [
        ("ENTERPRISE", 13, True, G700),
        ("별도 협의", 40, True, DARK),
        ("중견·대기업 맞춤", 13, False, G700),
        ("", 8, False, DARK),
        ("✓  전담 학습 모델", 13, False, G500),
        ("✓  온프레미스 배포", 13, False, G500),
        ("✓  SLA 보장", 13, False, G500),
        ("✓  API 연동", 13, False, G500),
    ])

    # ROI bar
    _rrect(sl, 120, 690, 1680, 60, fill=BLUE50, r=0.04)
    _txt(sl, 156, 700, 120, 40, "ROI 20x", sz=18, bold=True, color=BLUE)
    _txt(sl, 300, 703, 800, 36, "제안서 외주 1건 300만원+ → PRO 월 9.9만원으로 5건 생성", sz=14, color=DARK)
    _txt(sl, 1500, 703, 260, 36, "그로스 마진 82%", sz=15, bold=True, color=GREEN)


def slide_08_financials(prs):
    """재무 계획 (dark slide)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, DARK)
    _header(sl, "07", "재무 계획", "3개년 매출 전망 (보수적 시나리오) — BEP: Y2 중반", dark=True)

    years = [
        ("Y1 — 2026", "₩1.2억", [
            "유료 고객 100개사", "MRR ₩9.9M",
        ], "영업이익 △₩2.0억 (-167%)", RED, DSURF, None),
        ("Y2 — 2027", "₩7.3억", [
            "유료 고객 400개사 + 엔터프라이즈 3건", "MRR ₩48M",
        ], "영업이익 ₩0.4억 (5.5%)", YELLOW, DSURF, None),
        ("Y3 — 2028", "₩23억", [
            "유료 고객 1,000개사 + 엔터프라이즈 10건", "MRR ₩150M",
        ], "영업이익 ₩11.4억 (49.6%)", GREEN, DBLUE, BLUE),
    ]
    for i, (tag, rev, details, op, opc, bg, stroke) in enumerate(years):
        x = 120 + i * 560
        kw = dict(fill=bg, r=0.05)
        if stroke: kw["stroke"] = stroke; kw["sw"] = Pt(2)
        _rrect(sl, x, 220, 520, 280, **kw)
        _txt(sl, x+28, 238, 300, 24, tag, sz=13, bold=True, color=BLUE400 if stroke else G700)
        _txt(sl, x+28, 268, 400, 60, rev, sz=38, bold=True, color=WHITE)
        for j, d in enumerate(details):
            _txt(sl, x+28, 340 + j*26, 460, 24, d, sz=13, color=BLUE100 if stroke else G400)
        _txt(sl, x+28, 405, 460, 24, op, sz=14, bold=True, color=opc)

    # Unit Economics row
    _rrect(sl, 120, 530, 1680, 230, fill=DSURF, r=0.05)
    ue_cols = [
        ("Unit Economics (PRO)", [
            ("ARPU (월)", "₩99,000"),
            ("LLM API 비용", "~₩15,000"),
            ("인프라 비용", "~₩3,000"),
            ("그로스 마진", "~82%"),
        ], None, GREEN),
        ("핵심 지표", [
            ("CAC", "~₩200,000"),
            ("LTV (12개월)", "₩1,188,000"),
            ("LTV/CAC", "~6x"),
        ], None, BLUE),
        ("자금 소요: ₩2.0억", [
            ("제품 개발 50%", "₩1.0억"),
            ("마케팅/GTM 25%", "₩0.5억"),
            ("운영 25%", "₩0.5억"),
        ], None, None),
    ]
    for i, (title, rows, _, highlight) in enumerate(ue_cols):
        x = 152 + i * 540
        _txt(sl, x, 548, 480, 28, title, sz=15, bold=True, color=WHITE)
        for j, (label, val) in enumerate(rows):
            line = f"{label}  {val}"
            is_last = j == len(rows) - 1 and highlight is not None
            _txt(sl, x, 585 + j*28, 480, 24, line, sz=12,
                 bold=is_last, color=highlight if is_last else G400)


def slide_09_gtm(prs):
    """GTM 전략."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "08", "Go-to-Market 전략")

    phases = [
        ("Phase 1", "0~6개월", "Product-Led Growth", BLUE, BLUE50, [
            "· FREE → PRO 전환율 최적화",
            "· 콘텐츠 마케팅 (블로그 주 2회)",
            "· 입찰 커뮤니티 + 무료 분석 데모",
            "· NPS + 수정률 추적",
        ], "목표: 100개사", BLUE100),
        ("Phase 2", "6~12개월", "확장", GREEN, GREEN50, [
            "· 건설·컨설팅 산업 진출",
            "· 엔터프라이즈 영업 시작",
            "· 파트너십 3개+ 확보",
            "· 정부 디지털 바우처 활용",
        ], "목표: 500개사", GREEN100),
        ("Phase 3", "12~24개월", "플랫폼화", ORANGE, ORANGE50, [
            "· 민간 입찰 + 정부지원사업 확장",
            "· Layer 3 승패 분석 런칭",
            "· 데이터 플라이휠 가동",
            "· 승률 대시보드 + 가격제안서",
        ], "목표: 1,000개사", ORANGE100),
    ]
    for i, (badge, period, subtitle, accent, bg, items, goal, goal_bg) in enumerate(phases):
        x = 120 + i * 560
        _rrect(sl, x, 210, 520, 430, fill=bg, r=0.05)
        _rrect(sl, x+28, 228, 120, 28, fill=accent, r=0.4)
        _txt(sl, x+28, 229, 120, 26, badge, sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(sl, x+28, 268, 460, 36, period, sz=22, bold=True, color=DARK)
        _txt(sl, x+28, 302, 460, 24, subtitle, sz=14, bold=True, color=accent)
        for j, item in enumerate(items):
            _txt(sl, x+28, 340 + j*30, 460, 26, item, sz=12, color=G500)
        _rrect(sl, x+28, 475, 460, 36, fill=goal_bg, r=0.04)
        _txt(sl, x+28, 478, 460, 30, goal, sz=13, bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Target segment
    _rrect(sl, 120, 660, 1680, 100, fill=G50, stroke=G200, r=0.04)
    _mtxt(sl, 152, 668, 1600, 85, [
        ("1차 타겟: IT 서비스 중소기업", 18, True, DARK),
        ("종업원 10~100명  |  연 5건+ 공공 입찰 참여  |  BD 담당자 1~3명", 14, False, G500),
        ("2차 → 건설·엔지니어링 (6개월 후)  |  3차 → 정부지원사업 (12개월 후)", 13, True, BLUE),
    ])


def slide_10_moat(prs):
    """경쟁 우위 및 해자."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "09", "경쟁 우위 및 해자")

    _txt(sl, 120, 200, 400, 28, "단기 해자 (0~12개월)", sz=18, bold=True, color=DARK)

    short_moats = [
        ("End-to-End\n자동화", "유일하게 분석→생성→학습까지 한 플랫폼"),
        ("HWP\n완벽 지원", "공공조달 핵심 포맷\n경쟁사 대부분 미지원"),
        ("3계층\n학습 엔진", "Layer 1(범용) + Layer 2(회사맞춤) 이미 동작"),
        ("공공조달\n특화 지식", "동의어 17개 카테고리\n495유닛 전문 지식"),
        ("제안서\n생성 품질", "KRDS 디자인 가이드\n전문 컨설턴트급"),
    ]
    for i, (t, d) in enumerate(short_moats):
        x = 120 + i * 332
        _rrect(sl, x, 240, 310, 180, fill=BLUE50, r=0.04)
        _txt(sl, x+20, 256, 270, 60, t, sz=15, bold=True, color=DARK)
        _txt(sl, x+20, 330, 270, 80, d, sz=11, color=G500)

    _txt(sl, 120, 445, 400, 28, "장기 해자 (12개월+)", sz=18, bold=True, color=DARK)

    long_moats = [
        ("수정 학습 Lock-in", '사용자가 수정할수록 AI가 학습 → "92% 학습한 AI를 버릴 건가요?"'),
        ("데이터 네트워크 효과", "고객 증가 → 승패 데이터 축적 → 인사이트 정확도 → 더 많은 고객"),
        ("회사 DB 축적", "실적·인력·과거 제안서가 쌓일수록 전환 비용 증가"),
        ("시간 선점 우위", "Layer 3 승패 분석은 데이터 축적 필수 → 후발주자 진입 장벽"),
    ]
    for i, (t, d) in enumerate(long_moats):
        col = i % 2
        row = i // 2
        x = 120 + col * 860
        y = 480 + row * 140
        _rrect(sl, x, y, 820, 120, fill=DARK, r=0.04)
        _txt(sl, x+24, y+16, 772, 28, t, sz=16, bold=True, color=WHITE)
        _txt(sl, x+24, y+52, 772, 56, d, sz=12, color=G400)


def slide_11_team(prs):
    """팀."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "10", "팀")

    # Founder card
    _rrect(sl, 120, 210, 800, 320, fill=G50, stroke=G200, r=0.05)
    # Avatar circle
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, px(156), px(245), px(90), px(90))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE; s.line.fill.background()
    _txt(sl, 156, 260, 90, 60, "민", sz=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _mtxt(sl, 270, 240, 620, 270, [
        ("민경욱", 26, True, DARK),
        ("대표 / 기술 총괄", 16, True, BLUE),
        ("", 6, False, DARK),
        ("· 풀스택 개발 (React, Python, AI/ML)", 13, False, G500),
        ("· 공공조달 도메인 전문가", 13, False, G500),
        ("· Kira 전체 아키텍처 설계·구현", 13, False, G500),
        ("· 154+ 테스트, 30,000+ 코드 라인 — 1인 개발", 13, False, G500),
    ])

    # Tech proof box
    _rrect(sl, 980, 210, 820, 320, fill=DARK, r=0.05)
    _mtxt(sl, 1012, 228, 756, 290, [
        ("구현 완료 증빙", 18, True, WHITE),
        ("", 8, False, DARK),
        ("총 테스트            154개", 13, False, G400),
        ("E2E 검증            통과 (2026-02-28)", 13, False, G400),
        ("Layer 1 지식 유닛     495개", 13, False, G400),
        ("지원 문서 포맷       PDF, DOCX, HWP, HWPX, Excel, PPT", 13, False, G400),
        ("생성 문서 타입       DOCX, PPTX, XLSX, PNG", 13, False, G400),
        ("코드 라인            30,000+", 13, False, G400),
    ])

    # Hiring
    _txt(sl, 120, 555, 300, 28, "채용 계획", sz=18, bold=True, color=DARK)
    hires = [
        ("백엔드 개발자", "Y1 Q2", "RAG 엔진·SaaS 인프라"),
        ("프론트엔드 개발자", "Y1 Q3", "Chat UI·대시보드 고도화"),
        ("BD / 마케팅", "Y1 Q2", "고객 확보·파트너십"),
        ("CS / 제안서 전문가", "Y1 Q4", "품질 관리·고객 성공"),
    ]
    for i, (title, when, desc) in enumerate(hires):
        x = 120 + i * 432
        _rrect(sl, x, 590, 400, 130, fill=BLUE50, r=0.04)
        _txt(sl, x+20, 600, 360, 28, title, sz=15, bold=True, color=DARK)
        _txt(sl, x+20, 630, 360, 24, when, sz=12, bold=True, color=BLUE)
        _txt(sl, x+20, 658, 360, 24, desc, sz=12, color=G700)


def slide_12_kpi(prs):
    """KPI (dark slide)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, DARK)
    _header(sl, "11", "핵심 지표 (KPI)", dark=True)

    kpis = [
        ("500", "MAU (Y1 목표)", "월간 활성 사용자", WHITE),
        ("15%", "Free→PRO 전환율", "무료 → 유료 전환", BLUE),
        ("₩10M", "MRR (Y1 목표)", "월간 반복 매출", GREEN),
        ("< 5%", "Churn Rate", "월간 이탈률", YELLOW),
    ]
    for i, (num, label, desc, clr) in enumerate(kpis):
        x = 120 + i * 420
        _rrect(sl, x, 210, 390, 160, fill=DSURF, r=0.04)
        _txt(sl, x+24, 224, 342, 60, num, sz=38, bold=True, color=clr)
        _txt(sl, x+24, 290, 342, 24, label, sz=13, color=G700)
        _txt(sl, x+24, 316, 342, 24, desc, sz=11, color=G500)

    # Edit Rate box
    _rrect(sl, 120, 400, 1680, 220, fill=DBLUE, stroke=BLUE, sw=Pt(2), r=0.05)
    _txt(sl, 156, 418, 1000, 30, "수정률 (Edit Rate) — Kira 고유 KPI", sz=18, bold=True, color=WHITE)
    _txt(sl, 156, 455, 1400, 50, "AI 초안 대비 사용자 수정 비율 → 낮아질수록 AI 학습 완성도 높음 → 이탈 방어 강화", sz=14, color=BLUE100)

    # Flow: v1 → v5 → v10
    _txt(sl, 156, 520, 200, 30, "v1: 수정 45%", sz=16, bold=True, color=RED)
    _txt(sl, 380, 520, 40, 30, "→", sz=18, bold=True, color=G700)
    _txt(sl, 440, 520, 200, 30, "v5: 수정 20%", sz=16, bold=True, color=YELLOW)
    _txt(sl, 660, 520, 40, 30, "→", sz=18, bold=True, color=G700)
    _txt(sl, 720, 520, 400, 30, 'v10: 수정 8% = "학습도 92%"', sz=16, bold=True, color=GREEN)


def slide_13_roadmap(prs):
    """제품 로드맵."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "12", "제품 로드맵")

    phases = [
        ("2026 상반기", "현재 → 런칭", BLUE, BLUE50, [
            ("✅ 코어 엔진 완성", GREEN),
            ("✅ Layer 1 지식 DB (495유닛)", GREEN),
            ("✅ 제안서·PPT·WBS·실적기술서", GREEN),
            ("✅ 알림 시스템", GREEN),
            ("✅ 프로덕트 투어", GREEN),
            ("✅ 154개 테스트 통과", GREEN),
            ("⏳ 베타 런칭 (초대 50개사)", YELLOW),
        ]),
        ("2026 하반기", "고도화 + 정식 런칭", GREEN, GREEN50, [
            ("○ 회사 DB 온보딩 UI", G500),
            ("○ Layer 2 회사 스타일 연동", G500),
            ("○ GO/NO-GO 2.0 (승률+기대수익)", G500),
            ("○ 정식 런칭 + 결제 연동", G500),
            ("○ 엔터프라이즈 1호 계약", G500),
        ]),
        ("2027", "플랫폼화 + 확장", ORANGE, ORANGE50, [
            ("○ Layer 3 승패 분석", G500),
            ("○ 가격제안서 자동 산출", G500),
            ("○ 정부지원사업 계획서", G500),
            ("○ 민간 입찰 확장 (LH, 한전 등)", G500),
            ("○ 승률 대시보드", G500),
            ("○ 데이터 플라이휠 가동", G500),
        ]),
    ]
    for i, (badge, title, accent, bg, items) in enumerate(phases):
        x = 120 + i * 560
        _rrect(sl, x, 210, 520, 560, fill=bg, r=0.05)
        _rrect(sl, x+28, 228, 160, 28, fill=accent, r=0.4)
        _txt(sl, x+28, 229, 160, 26, badge, sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(sl, x+28, 270, 460, 36, title, sz=20, bold=True, color=DARK)
        for j, (text, clr) in enumerate(items):
            _txt(sl, x+28, 318 + j*30, 460, 26, text, sz=12, bold=(clr == GREEN or clr == YELLOW), color=clr)


def slide_14_vision(prs):
    """비전 / 클로징 (dark slide)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, DARK)
    # Accent bar
    _rect(sl, 140, 280, 5, 240, fill=BLUE)
    # Title
    _txt(sl, 180, 280, 1200, 120, '"입찰의 모든 순간을\nAI가 함께합니다"', sz=44, bold=True, color=WHITE)
    # Subtitle
    _txt(sl, 180, 440, 1000, 80, "공고 발견 → 분석 → 판단 → 생성 → 학습\n— 입찰 라이프사이클의 완전한 자동화", sz=22, color=G400)
    # Brand
    _txt(sl, 180, 560, 300, 40, "Kira", sz=26, bold=True, color=BLUE)
    _txt(sl, 180, 600, 300, 30, "by M&S Solutions", sz=18, color=G700)
    _txt(sl, 180, 640, 400, 28, "bill.min122@gmail.com", sz=14, color=G500)
    # Badge
    _rrect(sl, 180, 690, 280, 36, fill=DBLUE, stroke=BLUE, r=0.4)
    _txt(sl, 180, 692, 280, 32, "한성대학교 창업동아리 2026", sz=12, bold=True, color=BLUE400, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_learning(prs)
    slide_05_market(prs)
    slide_06_competition(prs)
    slide_07_business_model(prs)
    slide_08_financials(prs)
    slide_09_gtm(prs)
    slide_10_moat(prs)
    slide_11_team(prs)
    slide_12_kpi(prs)
    slide_13_roadmap(prs)
    slide_14_vision(prs)

    out = "/Users/min-kyungwook/Downloads/Kira_사업계획서_v2.pptx"
    prs.save(out)
    print(f"✅ Saved: {out}")
    print(f"   Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
