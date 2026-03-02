"""Phase 2 실제 파일 생성 통합 테스트.

LLM 호출 없이, 각 assembler/generator를 직접 호출하여
실제 파일(XLSX, PNG, DOCX, PPTX)이 올바르게 생성되는지 검증.
"""
import os
import sys
import tempfile

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from phase2_models import (
    WbsTask, PersonnelAllocation, SlideType, SlideContent,
    QnaPair, TrackRecordEntry, PersonnelEntry, MethodologyType,
)


# ---------------------------------------------------------------------------
# Fixture data
# ---------------------------------------------------------------------------

def _sample_wbs_tasks() -> list[WbsTask]:
    return [
        WbsTask(phase="착수", task_name="사업 착수 보고", start_month=1, duration_months=1,
                deliverables=["착수보고서"], responsible_role="PM", man_months=0.5),
        WbsTask(phase="착수", task_name="현황 분석", start_month=1, duration_months=1,
                deliverables=["현황분석서"], responsible_role="PL", man_months=1.0),
        WbsTask(phase="분석", task_name="요구사항 정의", start_month=2, duration_months=2,
                deliverables=["요구사항정의서"], responsible_role="PL", man_months=2.0),
        WbsTask(phase="분석", task_name="현행시스템 분석", start_month=2, duration_months=1,
                deliverables=["현행시스템분석서"], responsible_role="개발자", man_months=1.5),
        WbsTask(phase="설계", task_name="아키텍처 설계", start_month=3, duration_months=2,
                deliverables=["아키텍처설계서"], responsible_role="아키텍트", man_months=2.0),
        WbsTask(phase="설계", task_name="UI/UX 설계", start_month=3, duration_months=2,
                deliverables=["UI/UX설계서", "화면정의서"], responsible_role="디자이너", man_months=2.0),
        WbsTask(phase="설계", task_name="DB 설계", start_month=4, duration_months=1,
                deliverables=["DB설계서"], responsible_role="DBA", man_months=1.0),
        WbsTask(phase="구현", task_name="프론트엔드 개발", start_month=4, duration_months=3,
                deliverables=["소스코드", "단위테스트결과서"], responsible_role="개발자", man_months=4.0),
        WbsTask(phase="구현", task_name="백엔드 개발", start_month=4, duration_months=3,
                deliverables=["소스코드", "API문서"], responsible_role="개발자", man_months=4.0),
        WbsTask(phase="구현", task_name="코드 리뷰", start_month=5, duration_months=2,
                deliverables=["코드리뷰결과서"], responsible_role="PL", man_months=1.0),
        WbsTask(phase="시험", task_name="통합 테스트", start_month=7, duration_months=1,
                deliverables=["통합테스트결과서"], responsible_role="QA", man_months=1.5),
        WbsTask(phase="시험", task_name="성능 테스트", start_month=7, duration_months=1,
                deliverables=["성능테스트결과서"], responsible_role="QA", man_months=1.0),
        WbsTask(phase="시험", task_name="보안 점검", start_month=7, duration_months=1,
                deliverables=["보안점검결과서"], responsible_role="QA", man_months=0.5),
        WbsTask(phase="이행/종료", task_name="데이터 이관", start_month=8, duration_months=1,
                deliverables=["이관결과서"], responsible_role="DBA", man_months=1.0),
        WbsTask(phase="이행/종료", task_name="운영자 교육", start_month=8, duration_months=1,
                deliverables=["교육자료"], responsible_role="PM", man_months=0.5),
        WbsTask(phase="이행/종료", task_name="완료 보고", start_month=8, duration_months=1,
                deliverables=["완료보고서"], responsible_role="PM", man_months=0.5),
    ]


def _sample_personnel() -> list[PersonnelAllocation]:
    return [
        PersonnelAllocation(role="PM", grade="특급", total_man_months=1.5,
                            monthly_allocation=[0.5, 0, 0, 0, 0, 0, 0, 1.0]),
        PersonnelAllocation(role="PL", grade="고급", total_man_months=4.0,
                            monthly_allocation=[1.0, 1.0, 0, 0, 0.5, 0.5, 0, 1.0]),
        PersonnelAllocation(role="개발자", grade="중급", total_man_months=9.5,
                            monthly_allocation=[0, 1.5, 0, 2.5, 2.5, 2.5, 0.5, 0]),
        PersonnelAllocation(role="디자이너", grade="중급", total_man_months=2.0,
                            monthly_allocation=[0, 0, 1.0, 1.0, 0, 0, 0, 0]),
        PersonnelAllocation(role="DBA", grade="중급", total_man_months=2.0,
                            monthly_allocation=[0, 0, 0, 1.0, 0, 0, 0, 1.0]),
        PersonnelAllocation(role="QA", grade="중급", total_man_months=3.0,
                            monthly_allocation=[0, 0, 0, 0, 0, 0, 3.0, 0]),
    ]


def _sample_slides() -> list[SlideContent]:
    return [
        SlideContent(
            slide_type=SlideType.COVER,
            title="스마트시티 통합플랫폼 구축 사업\n기술제안서 발표",
            speaker_notes="안녕하십니까. 스마트시티 통합플랫폼 구축 사업 기술제안서 발표를 시작하겠습니다.",
        ),
        SlideContent(
            slide_type=SlideType.CONTENT,
            title="목차",
            bullets=[
                "1. 사업 이해 및 분석",
                "2. 추진 전략",
                "3. 기술 방안",
                "4. 프로젝트 일정",
                "5. 투입 인력",
                "6. 유사 수행실적",
                "7. Q&A",
            ],
            speaker_notes="발표는 총 7개 장으로 구성되어 있습니다.",
        ),
        SlideContent(
            slide_type=SlideType.CONTENT,
            title="사업 이해 및 분석",
            body="스마트시티 통합플랫폼은 도시 인프라를 IoT 센서로 연결하고 빅데이터 분석으로 시민 서비스를 개선합니다.",
            bullets=[
                "IoT 센서 1,000개 이상 통합 관제",
                "실시간 대시보드 및 알림 시스템",
                "시민 참여 모바일 앱 연동",
                "GIS 기반 공간정보 통합",
            ],
            speaker_notes="본 사업의 핵심은 다양한 도시 인프라 데이터를 하나의 플랫폼에서 통합 관리하는 것입니다.",
        ),
        SlideContent(
            slide_type=SlideType.BULLET,
            title="추진 전략",
            body="Waterfall 방법론 기반 체계적 프로젝트 관리",
            bullets=[
                "1단계: 현황 분석 및 요구사항 정의 (1~2개월)",
                "2단계: 시스템 설계 및 아키텍처 수립 (3~4개월)",
                "3단계: 핵심 기능 개발 및 테스트 (4~7개월)",
                "4단계: 시스템 전환 및 안정화 (7~8개월)",
            ],
            speaker_notes="4단계에 걸쳐 체계적으로 프로젝트를 수행합니다.",
        ),
        SlideContent(
            slide_type=SlideType.TABLE,
            title="기술 방안 — 핵심 기술 스택",
            table_data=[
                ["구분", "기술", "적용 영역"],
                ["프론트엔드", "React + TypeScript", "관제 대시보드, 모바일 앱"],
                ["백엔드", "FastAPI + PostgreSQL", "API 서버, 데이터 처리"],
                ["IoT", "MQTT + InfluxDB", "센서 데이터 수집/저장"],
                ["AI/ML", "TensorFlow + scikit-learn", "이상 감지, 예측 분석"],
                ["인프라", "Kubernetes + AWS", "클라우드 인프라, 자동 스케일링"],
            ],
            speaker_notes="최신 오픈소스 기술을 적용하여 안정성과 확장성을 확보합니다.",
        ),
        SlideContent(
            slide_type=SlideType.TIMELINE,
            title="프로젝트 일정",
            body="총 8개월 (착수 1개월 → 분석 2개월 → 설계 2개월 → 구현 3개월 → 시험 1개월 → 이행 1개월)",
            speaker_notes="각 단계별 주요 마일스톤을 설명드리겠습니다.",
        ),
        SlideContent(
            slide_type=SlideType.TEAM,
            title="투입 인력",
            body="총 6명, 22 M/M 투입",
            bullets=[
                "PM (특급) - 1.5 M/M: 프로젝트 총괄",
                "PL (고급) - 4.0 M/M: 기술 리드",
                "개발자 (중급) - 9.5 M/M: 프론트/백엔드 개발",
                "디자이너 (중급) - 2.0 M/M: UI/UX 설계",
                "DBA (중급) - 2.0 M/M: DB 설계/이관",
                "QA (중급) - 3.0 M/M: 테스트/품질 보증",
            ],
            speaker_notes="핵심 인력은 유사 프로젝트 수행 경험이 풍부한 전문가로 구성됩니다.",
        ),
        SlideContent(
            slide_type=SlideType.CONTENT,
            title="유사 수행실적",
            body="최근 5년 내 유사 스마트시티/플랫폼 구축 프로젝트 10건 이상 수행",
            bullets=[
                "부산시 IoT 관제시스템 구축 (30억, 2024)",
                "인천시 스마트교통 플랫폼 (25억, 2023)",
                "세종시 통합데이터 플랫폼 (20억, 2023)",
            ],
            speaker_notes="풍부한 유사 수행실적을 보유하고 있습니다.",
        ),
        SlideContent(
            slide_type=SlideType.QNA,
            title="Q & A",
            speaker_notes="질문 받겠습니다.",
        ),
        SlideContent(
            slide_type=SlideType.CLOSING,
            title="감사합니다",
            speaker_notes="경청해 주셔서 감사합니다.",
        ),
    ]


def _sample_qna_pairs() -> list[QnaPair]:
    return [
        QnaPair(
            question="IoT 센서 1,000개의 동시 접속을 어떻게 처리하시겠습니까?",
            answer="MQTT 브로커의 클러스터 구성과 메시지 큐잉을 통해 초당 10,000건 이상의 메시지를 안정적으로 처리합니다.",
            category="기술",
        ),
        QnaPair(
            question="기존 시스템과의 연동 방안은 무엇입니까?",
            answer="표준 REST API와 메시지 큐 기반 비동기 연동을 통해 레거시 시스템과의 원활한 통합을 보장합니다.",
            category="기술",
        ),
        QnaPair(
            question="데이터 보안은 어떻게 보장합니까?",
            answer="전구간 TLS 암호화, 역할 기반 접근제어(RBAC), 개인정보 비식별화 처리를 적용합니다.",
            category="보안",
        ),
        QnaPair(
            question="프로젝트 지연 시 리스크 관리 방안은?",
            answer="주간 스프린트 리뷰로 조기 감지하고, 핵심 경로 태스크에 완충 기간을 배치하여 대응합니다.",
            category="관리",
        ),
        QnaPair(
            question="운영 전환 후 기술 지원 계획은?",
            answer="3개월 하자보수 기간 + 1년 유지보수 계약으로 안정적 운영을 지원합니다.",
            category="관리",
        ),
    ]


def _sample_track_records() -> list[TrackRecordEntry]:
    return [
        TrackRecordEntry(
            project_name="부산시 IoT 관제시스템 구축",
            client="부산광역시",
            period="2024.03~2024.11",
            amount=30.0,
            description="부산시 전역의 IoT 센서를 통합 관제하는 시스템을 구축하였습니다.",
            technologies=["React", "FastAPI", "MQTT", "PostgreSQL"],
            relevance_score=0.92,
            generated_text="부산광역시 IoT 관제시스템 구축 사업은 도시 전역에 설치된 500여 개의 IoT 센서를 "
                           "통합 모니터링하는 관제 플랫폼을 구축한 프로젝트입니다. MQTT 프로토콜 기반 실시간 "
                           "데이터 수집과 React 기반 대시보드를 개발하여 관제 효율을 40% 향상시켰습니다.",
        ),
        TrackRecordEntry(
            project_name="인천시 스마트교통 플랫폼",
            client="인천광역시",
            period="2023.06~2024.02",
            amount=25.0,
            description="스마트교통 데이터 분석 및 신호 최적화 플랫폼",
            technologies=["Python", "TensorFlow", "GIS", "Docker"],
            relevance_score=0.85,
            generated_text="인천광역시 스마트교통 플랫폼 구축 사업은 교통 데이터를 AI 기반으로 분석하여 "
                           "신호 체계를 최적화하는 플랫폼입니다. TensorFlow 기반 교통량 예측 모델과 "
                           "GIS 기반 시각화를 통해 교통 혼잡도를 20% 감소시켰습니다.",
        ),
        TrackRecordEntry(
            project_name="세종시 통합데이터 플랫폼",
            client="세종특별자치시",
            period="2023.01~2023.09",
            amount=20.0,
            description="도시 데이터 수집/분석/공유 통합 플랫폼",
            technologies=["Kubernetes", "PostgreSQL", "Apache Kafka", "React"],
            relevance_score=0.78,
            generated_text="세종특별자치시 통합데이터 플랫폼 구축 사업은 도시의 다양한 데이터(교통, 환경, 에너지)를 "
                           "하나의 플랫폼에서 수집·분석·공유하는 체계를 구축한 프로젝트입니다.",
        ),
    ]


def _sample_personnel_entries() -> list[PersonnelEntry]:
    return [
        PersonnelEntry(
            name="김철수", role="PM", grade="특급", experience_years=20,
            certifications=["PMP", "정보관리기술사"],
            key_projects=["서울시 스마트시티 PM", "부산시 IoT 관제 PM"],
            generated_text="김철수 PM은 20년의 IT 프로젝트 관리 경험을 보유한 정보관리기술사로, "
                           "서울시 스마트시티, 부산시 IoT 관제 등 대규모 공공 프로젝트를 성공적으로 이끌었습니다.",
        ),
        PersonnelEntry(
            name="이영희", role="PL", grade="고급", experience_years=15,
            certifications=["정보처리기사", "AWS SA"],
            key_projects=["인천시 스마트교통 PL", "세종시 데이터플랫폼 PL"],
            generated_text="이영희 PL은 15년의 시스템 개발 경험과 AWS SA 자격을 보유하며, "
                           "다수의 스마트시티 프로젝트에서 기술 리더로 활약했습니다.",
        ),
        PersonnelEntry(
            name="박민수", role="개발자", grade="중급", experience_years=8,
            certifications=["정보처리기사"],
            key_projects=["부산시 IoT 백엔드 개발"],
            generated_text="박민수 개발자는 8년간 Python/FastAPI 기반 백엔드 시스템 개발 전문가로, "
                           "IoT 데이터 처리 파이프라인 구축에 강점을 보유하고 있습니다.",
        ),
    ]


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

def test_wbs_xlsx_generation():
    """WBS XLSX 파일 생성 검증."""
    from wbs_generator import generate_wbs_xlsx

    tasks = _sample_wbs_tasks()
    personnel = _sample_personnel()

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test_wbs.xlsx")
        result = generate_wbs_xlsx(tasks, personnel, "스마트시티 구축 WBS", 8, path)

        assert os.path.isfile(result), f"XLSX 파일이 생성되지 않음: {result}"
        size = os.path.getsize(result)
        assert size > 5000, f"XLSX 파일이 너무 작음: {size} bytes"

        # openpyxl로 내용 확인
        from openpyxl import load_workbook
        wb = load_workbook(result)
        assert "WBS" in wb.sheetnames, "WBS 시트 없음"
        assert "인력배치표" in wb.sheetnames, "인력배치표 시트 없음"
        assert "산출물" in wb.sheetnames, "산출물 시트 없음"

        ws_wbs = wb["WBS"]
        assert ws_wbs.max_row > 1, "WBS 데이터 행 없음"
        assert ws_wbs.cell(1, 1).value == "No.", "헤더 불일치"

        ws_staff = wb["인력배치표"]
        assert ws_staff.max_row > 1, "인력배치표 데이터 행 없음"

        ws_del = wb["산출물"]
        assert ws_del.max_row > 1, "산출물 데이터 행 없음"

        print(f"  XLSX: {size:,} bytes, WBS {ws_wbs.max_row}행, 인력 {ws_staff.max_row}행, 산출물 {ws_del.max_row}행")


def test_gantt_chart_generation():
    """간트차트 PNG 생성 검증."""
    from wbs_generator import generate_gantt_chart

    tasks = _sample_wbs_tasks()

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test_gantt.png")
        result = generate_gantt_chart(tasks, 8, path)

        assert os.path.isfile(result), f"간트차트 PNG 파일이 생성되지 않음: {result}"
        size = os.path.getsize(result)
        assert size > 10000, f"간트차트 이미지가 너무 작음: {size} bytes"

        # PNG 매직 바이트 확인
        with open(result, "rb") as f:
            header = f.read(8)
        assert header[:4] == b"\x89PNG", "PNG 파일 형식이 아님"

        print(f"  Gantt PNG: {size:,} bytes")


def test_wbs_docx_generation():
    """수행계획서 DOCX 생성 검증."""
    from wbs_generator import generate_wbs_docx, generate_gantt_chart

    tasks = _sample_wbs_tasks()
    personnel = _sample_personnel()

    with tempfile.TemporaryDirectory() as tmpdir:
        gantt_path = os.path.join(tmpdir, "gantt.png")
        generate_gantt_chart(tasks, 8, gantt_path)

        docx_path = os.path.join(tmpdir, "test_wbs.docx")
        result = generate_wbs_docx(
            tasks, personnel,
            "스마트시티 통합플랫폼 구축 수행계획서", 8, "Waterfall",
            docx_path, gantt_path,
        )

        assert os.path.isfile(result), f"DOCX 파일이 생성되지 않음: {result}"
        size = os.path.getsize(result)
        assert size > 10000, f"DOCX 파일이 너무 작음: {size} bytes"

        # python-docx로 내용 확인
        from docx import Document
        doc = Document(result)
        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert "수행방법론" in full_text, "수행방법론 섹션 없음"
        assert "WBS" in full_text, "WBS 섹션 없음"
        assert "투입인력" in full_text, "투입인력 섹션 없음"

        table_count = len(doc.tables)
        assert table_count >= 2, f"표가 부족함: {table_count}개"

        print(f"  WBS DOCX: {size:,} bytes, {len(doc.paragraphs)} paragraphs, {table_count} tables")


def test_pptx_generation():
    """PPT 발표자료 PPTX 생성 검증."""
    from ppt_assembler import assemble_pptx

    slides = _sample_slides()
    qna_pairs = _sample_qna_pairs()

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test_presentation.pptx")
        result = assemble_pptx(
            title="스마트시티 통합플랫폼 구축",
            slides=slides,
            qna_pairs=qna_pairs,
            output_path=path,
            company_name="(주)MS솔루션즈",
        )

        assert os.path.isfile(result), f"PPTX 파일이 생성되지 않음: {result}"
        size = os.path.getsize(result)
        assert size > 20000, f"PPTX 파일이 너무 작음: {size} bytes"

        # python-pptx로 내용 확인
        from pptx import Presentation
        prs = Presentation(result)
        assert len(prs.slides) == len(slides), f"슬라이드 수 불일치: {len(prs.slides)} != {len(slides)}"

        # 표지 슬라이드 확인
        first_slide = prs.slides[0]
        texts = [shape.text for shape in first_slide.shapes if hasattr(shape, "text")]
        assert any("스마트시티" in t for t in texts), "표지에 제목 없음"
        assert any("MS솔루션즈" in t for t in texts), "표지에 회사명 없음"

        # 발표 노트 확인
        notes_count = sum(1 for s in prs.slides if s.notes_slide and s.notes_slide.notes_text_frame.text.strip())
        assert notes_count >= 5, f"발표 노트가 부족함: {notes_count}개"

        # QnA 슬라이드 노트에 예상질문 포함 확인
        qna_slide = prs.slides[-2]  # QnA is second-to-last
        qna_notes = qna_slide.notes_slide.notes_text_frame.text
        assert "IoT 센서" in qna_notes, "QnA 노트에 예상질문이 없음"

        print(f"  PPTX: {size:,} bytes, {len(prs.slides)} slides, {notes_count} slides with notes")


def test_track_record_docx_generation():
    """실적/경력 기술서 DOCX 생성 검증."""
    from track_record_assembler import assemble_track_record_docx

    records = _sample_track_records()
    personnel = _sample_personnel_entries()

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test_track_record.docx")
        result = assemble_track_record_docx(
            title="스마트시티 통합플랫폼 구축 - 유사수행실적 및 투입인력 기술서",
            records=records,
            personnel=personnel,
            output_path=path,
        )

        assert os.path.isfile(result), f"DOCX 파일이 생성되지 않음: {result}"
        size = os.path.getsize(result)
        assert size > 5000, f"DOCX 파일이 너무 작음: {size} bytes"

        # python-docx로 내용 확인
        from docx import Document
        doc = Document(result)
        full_text = "\n".join(p.text for p in doc.paragraphs)

        assert "유사수행실적" in full_text, "유사수행실적 섹션 없음"
        assert "투입인력" in full_text, "투입인력 섹션 없음"
        assert "부산시 IoT" in full_text, "실적 상세 내용 없음"
        assert "김철수" in full_text, "인력 상세 내용 없음"

        table_count = len(doc.tables)
        assert table_count >= 2, f"표가 부족함: {table_count}개 (요약표+인력표 필요)"

        # 표 내용 확인
        summary_table = doc.tables[0]
        header_text = [cell.text for cell in summary_table.rows[0].cells]
        assert "프로젝트명" in header_text, "요약표 헤더 불일치"

        print(f"  Track Record DOCX: {size:,} bytes, {len(doc.paragraphs)} paragraphs, {table_count} tables")


def test_wbs_planner_fallback():
    """WBS planner fallback (LLM 없이 템플릿 기반)."""
    from unittest.mock import patch
    from wbs_planner import _fallback_tasks, _extract_project_duration, _detect_methodology
    from wbs_planner import WATERFALL_TEMPLATE, AGILE_TEMPLATE

    # Duration extraction
    assert _extract_project_duration({"project_period": "8개월"}) == 8
    assert _extract_project_duration({"project_period": "1년"}) == 12
    assert _extract_project_duration({"project_period": "12 months"}) == 12
    assert _extract_project_duration({}) == 6

    # Methodology detection (mock LLM to test keyword fallback path)
    with patch("wbs_planner.call_with_retry", side_effect=Exception("mock")):
        assert _detect_methodology({"text": "스크럼 기반 개발"}).value == "agile"
        rfx_agile = {"requirements": [{"description": "애자일 스프린트"}, {"description": "칸반 보드"}]}
        result_agile = _detect_methodology(rfx_agile)
        assert result_agile.value == "agile"

    # Fallback tasks
    tasks = _fallback_tasks(WATERFALL_TEMPLATE, 8)
    assert len(tasks) > 10, f"Fallback 태스크 수 부족: {len(tasks)}"
    assert tasks[0].phase == "착수", f"첫 번째 단계 불일치: {tasks[0].phase}"
    assert all(t.start_month >= 1 for t in tasks), "start_month < 1인 태스크 존재"

    agile_tasks = _fallback_tasks(AGILE_TEMPLATE, 6)
    assert len(agile_tasks) > 5
    assert agile_tasks[0].phase == "착수/백로그"

    print(f"  WBS Fallback: Waterfall {len(tasks)} tasks, Agile {len(agile_tasks)} tasks")


def test_ppt_slide_planner_defaults():
    """PPT slide planner default slide structure (fallback 경로)."""
    from unittest.mock import patch
    from ppt_slide_planner import plan_slides

    rfx_result = {
        "title": "스마트시티 구축",
        "issuing_org": "서울시",
        "requirements": [
            {"category": "기술", "description": "IoT 센서 통합"},
            {"category": "관리", "description": "PM 방법론"},
        ],
        "evaluation_criteria": [
            {"category": "기술", "item": "기술제안", "score": 60},
            {"category": "가격", "item": "입찰가격", "score": 40},
        ],
    }

    with patch("ppt_slide_planner.call_with_retry", side_effect=Exception("mock LLM failure")):
        slides = plan_slides(rfx_result, target_slide_count=15, duration_min=20)
    assert len(slides) >= 10, f"슬라이드 수 부족: {len(slides)}"
    assert slides[0].slide_type == SlideType.COVER, "첫 슬라이드가 표지가 아님"
    assert slides[-1].slide_type == SlideType.CLOSING, "마지막 슬라이드가 마무리가 아님"
    assert any(s.slide_type == SlideType.QNA for s in slides), "Q&A 슬라이드 없음"
    # KRDS: TOC + DIVIDER 슬라이드 포함 확인
    assert any(s.slide_type == SlideType.TOC for s in slides), "목차(TOC) 슬라이드 없음"
    assert any(s.slide_type == SlideType.DIVIDER for s in slides), "간지(DIVIDER) 슬라이드 없음"

    total_sec = sum(s.duration_sec for s in slides)
    print(f"  Slide Planner: {len(slides)} slides, total {total_sec}s ({total_sec/60:.1f}min)")


def test_pptx_krds_full_pipeline():
    """KRDS 6종 슬라이드 타입 포함 PPTX 풀 파이프라인."""
    from ppt_assembler import assemble_pptx, KRDS_COLORS

    slides = [
        SlideContent(slide_type=SlideType.COVER, title="2026 스마트시티 구축 사업\n기술제안서 발표"),
        SlideContent(
            slide_type=SlideType.TOC, title="목차",
            bullets=["사업 이해", "추진 전략", "기술 방안", "수행 일정", "투입 인력", "Q&A"],
        ),
        SlideContent(
            slide_type=SlideType.DIVIDER, title="01 사업 이해",
            body="사업의 배경과 목적을 분석합니다.",
        ),
        SlideContent(
            slide_type=SlideType.CONTENT, title="사업 개요",
            body="IoT 기반 스마트시티 통합관제 플랫폼 구축",
            bullets=["IoT 센서 1,000개 통합", "실시간 대시보드", "AI 이상감지"],
        ),
        SlideContent(
            slide_type=SlideType.DIVIDER, title="02 추진 전략",
            body="Waterfall 기반 체계적 추진 방안을 제시합니다.",
        ),
        SlideContent(
            slide_type=SlideType.BULLET, title="추진 전략",
            bullets=["1단계: 현황 분석 (1~2M)", "2단계: 설계 (3~4M)", "3단계: 구현 (5~7M)", "4단계: 이행 (8M)"],
        ),
        SlideContent(
            slide_type=SlideType.TABLE, title="핵심 기술 스택",
            table_data=[
                ["구분", "기술", "적용 영역"],
                ["프론트", "React", "대시보드"],
                ["백엔드", "FastAPI", "API 서버"],
                ["AI/ML", "TensorFlow", "이상 감지"],
            ],
        ),
        SlideContent(slide_type=SlideType.TIMELINE, title="수행 일정", body="8개월 계획"),
        SlideContent(
            slide_type=SlideType.TEAM, title="투입 인력", body="총 6명",
            bullets=["PM 특급 1명", "PL 고급 1명", "개발자 중급 3명", "QA 중급 1명"],
        ),
        SlideContent(slide_type=SlideType.QNA, title="Q&A"),
        SlideContent(slide_type=SlideType.CLOSING, title="감사합니다"),
    ]
    qna_pairs = [
        QnaPair(question="IoT 보안은?", answer="TLS + RBAC 적용", category="기술"),
        QnaPair(question="일정 지연 시?", answer="버퍼 1개월 확보", category="관리"),
    ]

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "krds_full.pptx")
        result = assemble_pptx(
            title="KRDS 풀 파이프라인",
            slides=slides,
            qna_pairs=qna_pairs,
            output_path=path,
            company_name="(주)MS솔루션즈",
        )

        assert os.path.isfile(result)
        size = os.path.getsize(result)
        assert size > 20000, f"PPTX 파일이 너무 작음: {size} bytes"

        from pptx import Presentation
        prs = Presentation(result)
        assert len(prs.slides) == len(slides), f"슬라이드 수 불일치: {len(prs.slides)}"

        # KRDS 색상 적용 확인
        from pptx.dml.color import RGBColor
        assert KRDS_COLORS["primary"] == RGBColor(0x00, 0x37, 0x64)

        # 슬라이드 타입별 존재 확인
        type_counts = {}
        for s in slides:
            t = s.slide_type.value
            type_counts[t] = type_counts.get(t, 0) + 1
        assert type_counts.get("toc", 0) >= 1
        assert type_counts.get("divider", 0) >= 2
        assert type_counts.get("cover", 0) >= 1
        assert type_counts.get("closing", 0) >= 1

        print(f"  KRDS Full Pipeline: {size:,} bytes, {len(prs.slides)} slides")


if __name__ == "__main__":
    tests = [
        ("WBS XLSX 생성", test_wbs_xlsx_generation),
        ("간트차트 PNG 생성", test_gantt_chart_generation),
        ("수행계획서 DOCX 생성", test_wbs_docx_generation),
        ("PPT 발표자료 생성", test_pptx_generation),
        ("KRDS 풀 파이프라인", test_pptx_krds_full_pipeline),
        ("실적/경력 기술서 DOCX 생성", test_track_record_docx_generation),
        ("WBS Planner Fallback", test_wbs_planner_fallback),
        ("PPT Slide Planner", test_ppt_slide_planner_defaults),
    ]

    print("=" * 60)
    print("Phase 2 실제 파일 생성 통합 테스트")
    print("=" * 60)

    passed = 0
    failed = 0
    for name, test_fn in tests:
        try:
            print(f"\n[TEST] {name}")
            test_fn()
            print(f"  => PASS")
            passed += 1
        except Exception as e:
            print(f"  => FAIL: {e}")
            import traceback
            traceback.print_exc()
            failed += 1

    print("\n" + "=" * 60)
    print(f"결과: {passed} passed, {failed} failed / {len(tests)} total")
    print("=" * 60)
