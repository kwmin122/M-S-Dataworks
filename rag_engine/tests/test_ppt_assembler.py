"""Tests for ppt_assembler.py."""
import sys, os, tempfile
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx.dml.color import RGBColor

from phase2_models import SlideType, SlideContent, QnaPair
from ppt_assembler import assemble_pptx


def _sample_slides():
    return [
        SlideContent(
            slide_type=SlideType.COVER,
            title="스마트시티 통합플랫폼 구축",
            speaker_notes="표지 슬라이드입니다.",
        ),
        SlideContent(
            slide_type=SlideType.BULLET,
            title="목차",
            bullets=["사업 이해", "추진 전략", "기술 방안", "일정", "Q&A"],
        ),
        SlideContent(
            slide_type=SlideType.CONTENT,
            title="사업 이해",
            body="본 사업은 IoT 기반 스마트시티 구축",
            bullets=["IoT 센서 통합", "실시간 모니터링"],
            speaker_notes="사업의 핵심 내용을 설명합니다.",
        ),
        SlideContent(
            slide_type=SlideType.TABLE,
            title="유사 실적",
            table_data=[
                ["프로젝트", "발주처", "금액"],
                ["관제시스템", "부산시", "30억"],
                ["IoT 플랫폼", "과기부", "20억"],
            ],
        ),
        SlideContent(
            slide_type=SlideType.TIMELINE,
            title="수행 일정",
            body="8개월 수행 계획",
        ),
        SlideContent(
            slide_type=SlideType.QNA,
            title="Q&A",
        ),
        SlideContent(
            slide_type=SlideType.CLOSING,
            title="감사합니다",
        ),
    ]


def _sample_qna():
    return [
        QnaPair(question="IoT 보안 대책?", answer="암호화 + 인증 체계 구축", category="기술"),
        QnaPair(question="일정 지연 시 대응?", answer="버퍼 기간 1개월 확보", category="관리"),
    ]


def test_assemble_pptx_creates_file():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "test.pptx")
        result = assemble_pptx(
            title="스마트시티",
            slides=_sample_slides(),
            qna_pairs=_sample_qna(),
            output_path=path,
            company_name="테스트회사",
        )
        assert os.path.isfile(result)
        assert os.path.getsize(result) > 0


def test_assemble_pptx_empty_slides():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "empty.pptx")
        result = assemble_pptx(
            title="빈 발표자료",
            slides=[],
            qna_pairs=[],
            output_path=path,
        )
        assert os.path.isfile(result)


def test_assemble_pptx_without_company():
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "no_company.pptx")
        result = assemble_pptx(
            title="테스트",
            slides=_sample_slides(),
            qna_pairs=_sample_qna(),
            output_path=path,
        )
        assert os.path.isfile(result)


def test_assemble_pptx_krds_full_structure():
    """KRDS 6종 슬라이드 타입 모두 포함하여 생성."""
    slides = [
        SlideContent(slide_type=SlideType.COVER, title="2026년도 사업계획서"),
        SlideContent(
            slide_type=SlideType.TOC, title="목차",
            bullets=["사업 개요", "추진 현황", "성과 분석", "향후 계획"],
        ),
        SlideContent(
            slide_type=SlideType.DIVIDER, title="01 사업 개요",
            body="2026년 상반기 사업 추진 현황을 분석합니다.",
        ),
        SlideContent(
            slide_type=SlideType.CONTENT, title="사업 추진 현황",
            body="주요 추진 실적",
            bullets=["클라우드 전환 완료", "보안 체계 강화", "AI 도입 추진"],
        ),
        SlideContent(
            slide_type=SlideType.TABLE, title="성과 분석",
            table_data=[
                ["구분", "목표", "실적", "달성률"],
                ["매출", "500억", "471억", "94.2%"],
                ["고객수", "10만", "11.2만", "112%"],
            ],
        ),
        SlideContent(slide_type=SlideType.TIMELINE, title="수행 일정", body="12개월 계획"),
        SlideContent(
            slide_type=SlideType.DIVIDER, title="04 향후 계획",
            body="차년도 추진 방향을 제시합니다.",
        ),
        SlideContent(
            slide_type=SlideType.BULLET, title="기대 효과",
            bullets=["비용 20% 절감", "서비스 품질 향상"],
        ),
        SlideContent(slide_type=SlideType.QNA, title="Q&A"),
        SlideContent(slide_type=SlideType.CLOSING, title="감사합니다"),
    ]
    qna = [
        QnaPair(question="AI 도입 범위?", answer="자연어 처리 + 이미지 분석", category="기술"),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "krds_full.pptx")
        result = assemble_pptx(
            title="KRDS 풀 구조",
            slides=slides,
            qna_pairs=qna,
            output_path=path,
            company_name="한국oo공사",
        )
        assert os.path.isfile(result)
        assert os.path.getsize(result) > 0


def test_assemble_pptx_krds_colors_applied():
    """KRDS_COLORS가 기본으로 적용되는지 확인."""
    from ppt_assembler import KRDS_COLORS, DEFAULT_COLORS
    # KRDS_COLORS와 DEFAULT_COLORS가 동일 객체
    assert KRDS_COLORS is DEFAULT_COLORS
    # Blue 900 primary 색상 확인
    assert KRDS_COLORS["primary"] == RGBColor(0x00, 0x37, 0x64)


def test_assemble_pptx_overflow_part_labeling():
    """불릿 오버플로우 시 Part N/M 번호 및 speaker notes에 슬라이드 번호 삽입."""
    from ppt_assembler import _MAX_BULLETS_PER_SLIDE
    from pptx import Presentation as PptxPres

    # Create a content slide with more bullets than max
    many_bullets = [f"항목 {i+1}" for i in range(_MAX_BULLETS_PER_SLIDE + 3)]
    slides = [
        SlideContent(
            slide_type=SlideType.CONTENT,
            title="핵심 전략",
            body="주요 전략 항목",
            bullets=many_bullets,
            speaker_notes="전략을 설명합니다.",
        ),
    ]
    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, "overflow.pptx")
        assemble_pptx(
            title="오버플로우 테스트",
            slides=slides,
            qna_pairs=[],
            output_path=path,
        )
        assert os.path.isfile(path)

        # Verify the PPTX has 2 slides (overflow split)
        prs = PptxPres(path)
        assert len(prs.slides) == 2

        # First slide: original title (no Part suffix for first)
        first_texts = [sh.text for sh in prs.slides[0].shapes if hasattr(sh, 'text')]
        assert any("핵심 전략" in t for t in first_texts)
        # Should NOT contain "Part" in the first slide title
        assert not any("Part" in t and "핵심 전략" in t for t in first_texts)

        # Second slide: should have "Part 2/2" in the title
        second_texts = [sh.text for sh in prs.slides[1].shapes if hasattr(sh, 'text')]
        assert any("Part 2/2" in t for t in second_texts)

        # Speaker notes on second slide should contain "[슬라이드 2/2]"
        notes_text = prs.slides[1].notes_slide.notes_text_frame.text
        assert "[슬라이드 2/2]" in notes_text
