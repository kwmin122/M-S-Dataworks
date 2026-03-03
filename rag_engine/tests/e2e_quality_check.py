"""E2E 품질 검증: 실제 API 호출로 모든 문서 타입 생성 + 콘텐츠 품질 분석.

이 테스트는 실제 OpenAI API를 호출하므로 `pytest -m e2e`로 실행.
일반 테스트 스위트(pytest -q)에서는 자동 제외됨.

Usage:
    cd rag_engine && python tests/test_e2e_quality.py
"""
import json
import os
import sys
import tempfile
import time

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from dotenv import load_dotenv

load_dotenv()


# ---------------------------------------------------------------------------
# 테스트 RFP 데이터 (실제 스마트시티 프로젝트 시나리오)
# ---------------------------------------------------------------------------

RFX_RESULT = {
    "title": "스마트시티 통합플랫폼 구축",
    "issuing_org": "서울특별시 스마트도시정책관",
    "budget": "50억원",
    "project_period": "8개월",
    "rfp_text_summary": (
        "서울특별시 스마트도시정책관은 도시 인프라를 IoT 센서로 연결하고 "
        "빅데이터 분석 기반 시민 서비스 개선을 위한 통합플랫폼을 구축한다. "
        "본 사업은 IoT 센서 1,000개 이상의 데이터를 실시간 수집·분석하고 "
        "GIS 기반 관제 대시보드 및 시민 참여 모바일 앱을 개발하는 것을 목표로 한다. "
        "개발 완료 후 3개월 하자보수 기간이 포함되며, "
        "행정안전부 클라우드 보안 인증(CSAP) 환경에서 운영되어야 한다."
    ),
    "requirements": [
        {"category": "기술", "description": "IoT 센서 1,000개 이상 실시간 데이터 수집 및 통합 관제"},
        {"category": "기술", "description": "GIS 기반 공간정보 통합 관제 대시보드 개발"},
        {"category": "기술", "description": "빅데이터 분석 및 AI 이상 감지 엔진 구축"},
        {"category": "기술", "description": "시민 참여 모바일 앱 (iOS/Android) 개발"},
        {"category": "보안", "description": "CSAP 인증 클라우드 환경 구축 및 개인정보 비식별화"},
        {"category": "관리", "description": "PMO 방법론 기반 체계적 프로젝트 관리"},
        {"category": "관리", "description": "운영자/관리자 교육 및 매뉴얼 작성"},
    ],
    "evaluation_criteria": [
        {"category": "기술", "item": "시스템 아키텍처", "score": 25},
        {"category": "기술", "item": "핵심 기능 구현 방안", "score": 20},
        {"category": "기술", "item": "데이터 분석/AI", "score": 15},
        {"category": "관리", "item": "프로젝트 관리 체계", "score": 10},
        {"category": "관리", "item": "보안/품질 관리", "score": 10},
        {"category": "실적", "item": "유사수행실적", "score": 10},
        {"category": "가격", "item": "입찰가격", "score": 10},
    ],
}


def check_api_key():
    """OPENAI_API_KEY가 로드되었는지 확인."""
    key = os.environ.get("OPENAI_API_KEY", "")
    if not key or len(key) < 10:
        print("[FATAL] OPENAI_API_KEY not loaded from .env. Cannot proceed.")
        sys.exit(1)
    print(f"[OK] OPENAI_API_KEY loaded ({key[:8]}...)")


# ---------------------------------------------------------------------------
# 1. WBS 생성
# ---------------------------------------------------------------------------

def test_wbs_generation(output_dir: str) -> dict:
    """WBS 파이프라인 E2E 테스트."""
    from wbs_orchestrator import generate_wbs

    print("\n" + "=" * 60)
    print("[TEST] WBS 생성 (수행계획서 + XLSX + 간트차트)")
    print("=" * 60)

    result = generate_wbs(
        rfx_result=RFX_RESULT,
        output_dir=output_dir,
    )

    # File existence check
    assert os.path.isfile(result.xlsx_path), f"XLSX 미생성: {result.xlsx_path}"
    assert os.path.isfile(result.docx_path), f"DOCX 미생성: {result.docx_path}"

    xlsx_size = os.path.getsize(result.xlsx_path)
    docx_size = os.path.getsize(result.docx_path)

    print(f"  XLSX: {xlsx_size:,} bytes")
    print(f"  DOCX: {docx_size:,} bytes")
    if result.gantt_path:
        gantt_size = os.path.getsize(result.gantt_path)
        print(f"  Gantt PNG: {gantt_size:,} bytes")
    print(f"  Tasks: {len(result.tasks)}")
    print(f"  Personnel: {len(result.personnel)}")
    print(f"  Total months: {result.total_months}")
    print(f"  Time: {result.generation_time_sec}s")

    # Quality checks
    issues = []

    # Task count: 15~30 (RFP 맞춤이면 9건 모두 동일하지 않아야)
    if len(result.tasks) < 10:
        issues.append(f"[WARNING] 태스크 수 부족: {len(result.tasks)} (최소 10개 필요)")
    elif len(result.tasks) > 40:
        issues.append(f"[WARNING] 태스크 수 과다: {len(result.tasks)} (30개 이하 권장)")

    # Phase diversity
    phases = set(t.phase for t in result.tasks)
    if len(phases) < 3:
        issues.append(f"[WARNING] 단계 다양성 부족: {phases}")

    # RFP-specific content check
    task_names = " ".join(t.task_name for t in result.tasks)
    rfp_keywords = ["IoT", "센서", "GIS", "대시보드", "모바일", "빅데이터", "AI", "클라우드", "보안"]
    matched_keywords = [kw for kw in rfp_keywords if kw.lower() in task_names.lower()]
    if len(matched_keywords) < 3:
        issues.append(f"[WARNING] RFP 맞춤 부족: 태스크명에 '{matched_keywords}' 만 포함 (3개 이상 필요)")

    # Deliverables check
    all_deliverables = []
    for t in result.tasks:
        all_deliverables.extend(t.deliverables)
    if len(all_deliverables) < 10:
        issues.append(f"[WARNING] 산출물 수 부족: {len(all_deliverables)} (최소 10개 필요)")

    # Duration clamping check
    for t in result.tasks:
        end = t.start_month + t.duration_months - 1
        if end > result.total_months:
            issues.append(f"[ERROR] 기간 초과: {t.task_name} end_month={end} > total={result.total_months}")

    # Print tasks
    print("\n  [태스크 목록]")
    for i, t in enumerate(result.tasks, 1):
        print(f"    {i:2d}. [{t.phase}] {t.task_name} (M{t.start_month}~M{t.end_month()}, {t.man_months}MM)")

    if issues:
        print(f"\n  [품질 이슈] {len(issues)}건:")
        for issue in issues:
            print(f"    {issue}")
    else:
        print("\n  [품질] PASS - 모든 검증 통과")

    return {"issues": issues, "task_count": len(result.tasks), "time": result.generation_time_sec}


# ---------------------------------------------------------------------------
# 2. PPT 생성
# ---------------------------------------------------------------------------

def test_ppt_generation(output_dir: str) -> dict:
    """PPT 파이프라인 E2E 테스트."""
    from ppt_orchestrator import generate_ppt

    print("\n" + "=" * 60)
    print("[TEST] PPT 발표자료 생성 (PPTX + QnA)")
    print("=" * 60)

    proposal_sections = [
        {"name": "사업 이해", "text": "본 사업은 서울특별시의 스마트시티 통합플랫폼을 구축하는 것으로, IoT 센서 1,000개 이상의 데이터를 실시간 수집하고 GIS 기반 관제 대시보드를 통해 도시 인프라를 통합 모니터링합니다."},
        {"name": "추진 전략", "text": "PMO 방법론 기반의 체계적 프로젝트 관리를 적용하며, 착수→분석→설계→구현→시험→이행의 6단계로 진행합니다. 특히 IoT 데이터 수집 파이프라인과 AI 이상감지 엔진을 핵심 기술 요소로 개발합니다."},
        {"name": "기술 방안", "text": "마이크로서비스 아키텍처 기반으로 IoT 데이터 수집(MQTT), 실시간 처리(Apache Kafka), 빅데이터 저장(InfluxDB), AI 분석(TensorFlow), 관제 대시보드(React), 모바일 앱(Flutter)을 구축합니다."},
        {"name": "품질 관리", "text": "CSAP 인증 클라우드 환경에서 운영되며, 전구간 TLS 암호화, RBAC 기반 접근제어, 개인정보 비식별화를 적용합니다. ISO 27001 기반 정보보안 관리체계를 수립합니다."},
    ]

    result = generate_ppt(
        rfx_result=RFX_RESULT,
        output_dir=output_dir,
        proposal_sections=proposal_sections,
        duration_min=20,
        target_slide_count=15,
        qna_count=10,
        company_name="(주)스마트테크",
    )

    assert os.path.isfile(result.pptx_path), f"PPTX 미생성: {result.pptx_path}"
    pptx_size = os.path.getsize(result.pptx_path)

    print(f"  PPTX: {pptx_size:,} bytes")
    print(f"  Slides: {result.slide_count}")
    print(f"  QnA Pairs: {len(result.qna_pairs)}")
    print(f"  Duration: {result.total_duration_min} min")
    print(f"  Time: {result.generation_time_sec}s")

    # Quality checks
    issues = []

    if result.slide_count < 10:
        issues.append(f"[WARNING] 슬라이드 수 부족: {result.slide_count} (최소 10장 필요)")

    if len(result.qna_pairs) < 5:
        issues.append(f"[WARNING] QnA 수 부족: {len(result.qna_pairs)} (최소 5개 필요)")

    # Content check via python-pptx
    from pptx import Presentation

    prs = Presentation(result.pptx_path)
    all_text = []
    slides_with_notes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text)
        if slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            slides_with_notes += 1

    full_text = " ".join(all_text)

    # Placeholder check: "placeholder", "Lorem", "TODO" 등
    placeholder_words = ["placeholder", "lorem", "TODO", "여기에", "내용 입력"]
    for pw in placeholder_words:
        if pw.lower() in full_text.lower():
            issues.append(f"[ERROR] 플레이스홀더 텍스트 발견: '{pw}'")

    # RFP-specific content
    rfp_keywords = ["IoT", "센서", "GIS", "스마트시티", "CSAP", "대시보드"]
    matched = [kw for kw in rfp_keywords if kw.lower() in full_text.lower()]
    if len(matched) < 3:
        issues.append(f"[WARNING] PPT에 RFP 키워드 부족: {matched}")

    # Notes coverage
    if slides_with_notes < result.slide_count * 0.5:
        issues.append(f"[WARNING] 발표 노트 부족: {slides_with_notes}/{result.slide_count} 슬라이드")

    # QnA quality
    if result.qna_pairs:
        print("\n  [예상 질문/답변 샘플]")
        for i, qna in enumerate(result.qna_pairs[:3], 1):
            print(f"    Q{i}. [{qna.category}] {qna.question}")
            print(f"    A{i}. {qna.answer[:80]}...")

        # Check QnA category diversity
        categories = set(q.category for q in result.qna_pairs)
        if len(categories) < 2:
            issues.append(f"[WARNING] QnA 카테고리 다양성 부족: {categories}")

    if issues:
        print(f"\n  [품질 이슈] {len(issues)}건:")
        for issue in issues:
            print(f"    {issue}")
    else:
        print("\n  [품질] PASS - 모든 검증 통과")

    return {"issues": issues, "slide_count": result.slide_count, "qna_count": len(result.qna_pairs), "time": result.generation_time_sec}


# ---------------------------------------------------------------------------
# 3. 제안서 생성
# ---------------------------------------------------------------------------

def test_proposal_generation(output_dir: str) -> dict:
    """제안서 파이프라인 E2E 테스트."""
    from proposal_orchestrator import generate_proposal

    print("\n" + "=" * 60)
    print("[TEST] 제안서 DOCX 생성 (A-lite)")
    print("=" * 60)

    result = generate_proposal(
        rfx_result=RFX_RESULT,
        output_dir=output_dir,
        total_pages=30,
    )

    assert os.path.isfile(result.docx_path), f"DOCX 미생성: {result.docx_path}"
    docx_size = os.path.getsize(result.docx_path)

    print(f"  DOCX: {docx_size:,} bytes")
    print(f"  Sections: {len(result.sections)}")
    print(f"  Quality issues: {len(result.quality_issues)}")
    print(f"  Time: {result.generation_time_sec}s")

    issues = []

    # Section count
    if len(result.sections) < 3:
        issues.append(f"[WARNING] 섹션 수 부족: {len(result.sections)}")

    # Content analysis
    from docx import Document

    doc = Document(result.docx_path)
    full_text = "\n".join(p.text for p in doc.paragraphs)

    # Placeholder check
    placeholder_words = ["placeholder", "lorem", "TODO", "여기에"]
    for pw in placeholder_words:
        if pw.lower() in full_text.lower():
            issues.append(f"[ERROR] 플레이스홀더 텍스트 발견: '{pw}'")

    # Minimum text volume
    total_chars = len(full_text)
    if total_chars < 5000:
        issues.append(f"[WARNING] 제안서 분량 부족: {total_chars}자 (최소 5,000자 필요)")

    # RFP keywords
    rfp_keywords = ["IoT", "센서", "GIS", "스마트시티", "빅데이터", "CSAP"]
    matched = [kw for kw in rfp_keywords if kw.lower() in full_text.lower()]
    if len(matched) < 3:
        issues.append(f"[WARNING] RFP 키워드 부족: {matched}")

    # Quality checker results (blinding, vagueness)
    blind_issues = [q for q in result.quality_issues if q.category == "blind_violation"]
    vague_issues = [q for q in result.quality_issues if q.category == "vague_expression"]
    if blind_issues:
        print(f"\n  [블라인드 위반] {len(blind_issues)}건:")
        for bi in blind_issues[:3]:
            print(f"    - {bi.message}")
    if vague_issues:
        print(f"\n  [모호 표현] {len(vague_issues)}건:")
        for vi in vague_issues[:3]:
            print(f"    - {vi.message}")

    # Print section names
    print("\n  [섹션 목록]")
    for name, text in result.sections:
        print(f"    - {name}: {len(text)}자")

    if issues:
        print(f"\n  [품질 이슈] {len(issues)}건:")
        for issue in issues:
            print(f"    {issue}")
    else:
        print("\n  [품질] PASS - 모든 검증 통과")

    return {"issues": issues, "sections": len(result.sections), "chars": total_chars, "time": result.generation_time_sec}


# ---------------------------------------------------------------------------
# 4. auto_learner 학습 루프
# ---------------------------------------------------------------------------

def test_auto_learner_loop() -> dict:
    """auto_learner 학습 파이프라인 E2E 테스트."""
    from auto_learner import process_edit_feedback, get_learned_patterns, save_state, load_state

    print("\n" + "=" * 60)
    print("[TEST] auto_learner 학습 루프")
    print("=" * 60)

    issues = []

    # 1. Record edits (3회 반복으로 패턴 학습 트리거)
    original = "본 사업은 IoT 센서를 활용한 스마트시티 플랫폼을 구축합니다."
    edited = "본 사업은 IoT 센서 1,000개를 실시간 연동하여 도시 인프라 통합 관제 플랫폼을 구축합니다."

    for i in range(3):
        process_edit_feedback(
            company_id="test-company",
            section_name="사업 이해",
            original_text=original,
            edited_text=edited,
            doc_type="proposal",
        )

    # WBS doc_type
    wbs_original = '{"phase": "분석", "task_name": "요구사항 분석"}'
    wbs_edited = '{"phase": "분석", "task_name": "IoT 센서 요구사항 상세 분석"}'
    for i in range(3):
        process_edit_feedback(
            company_id="test-company",
            section_name="wbs_task",
            original_text=wbs_original,
            edited_text=wbs_edited,
            doc_type="wbs",
        )

    # 2. Check learned patterns
    proposal_patterns = get_learned_patterns(company_id="test-company", doc_type="proposal")
    wbs_patterns = get_learned_patterns(company_id="test-company", doc_type="wbs")

    print(f"  Proposal patterns: {len(proposal_patterns)}")
    print(f"  WBS patterns: {len(wbs_patterns)}")

    if not proposal_patterns:
        issues.append("[WARNING] proposal 학습 패턴 없음 (3회 반복 후)")
    if not wbs_patterns:
        issues.append("[WARNING] wbs 학습 패턴 없음 (3회 반복 후)")

    # 3. Save/Load persistence
    with tempfile.TemporaryDirectory() as tmpdir:
        save_state(tmpdir)
        state_file = os.path.join(tmpdir, "learning_state.json")
        if os.path.isfile(state_file):
            with open(state_file, "r") as f:
                state = json.load(f)
            print(f"  State file: {os.path.getsize(state_file):,} bytes")
            print(f"  Saved histories: {len(state.get('histories', {}))}")
            print(f"  Saved patterns: {len(state.get('learned_patterns', {}))}")
        else:
            issues.append("[ERROR] auto_learner 상태 파일 미생성")

        # Load and verify
        load_state(tmpdir)
        reloaded_patterns = get_learned_patterns(company_id="test-company", doc_type="proposal")
        if len(reloaded_patterns) != len(proposal_patterns):
            issues.append(f"[ERROR] 로드 후 패턴 불일치: {len(reloaded_patterns)} != {len(proposal_patterns)}")

    if issues:
        print(f"\n  [품질 이슈] {len(issues)}건:")
        for issue in issues:
            print(f"    {issue}")
    else:
        print("\n  [품질] PASS - 학습 루프 정상 동작")

    return {"issues": issues, "proposal_patterns": len(proposal_patterns), "wbs_patterns": len(wbs_patterns)}


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    check_api_key()

    total_start = time.time()
    results = {}

    with tempfile.TemporaryDirectory() as output_dir:
        print(f"\n출력 디렉토리: {output_dir}")

        # 1. WBS
        try:
            results["wbs"] = test_wbs_generation(output_dir)
        except Exception as e:
            print(f"\n  [FAIL] WBS 생성 실패: {e}")
            import traceback
            traceback.print_exc()
            results["wbs"] = {"error": str(e)}

        # 2. PPT
        try:
            results["ppt"] = test_ppt_generation(output_dir)
        except Exception as e:
            print(f"\n  [FAIL] PPT 생성 실패: {e}")
            import traceback
            traceback.print_exc()
            results["ppt"] = {"error": str(e)}

        # 3. 제안서
        try:
            results["proposal"] = test_proposal_generation(output_dir)
        except Exception as e:
            print(f"\n  [FAIL] 제안서 생성 실패: {e}")
            import traceback
            traceback.print_exc()
            results["proposal"] = {"error": str(e)}

        # 4. 학습 루프
        try:
            results["learner"] = test_auto_learner_loop()
        except Exception as e:
            print(f"\n  [FAIL] 학습 루프 실패: {e}")
            import traceback
            traceback.print_exc()
            results["learner"] = {"error": str(e)}

    # Summary
    total_time = round(time.time() - total_start, 1)
    print("\n" + "=" * 60)
    print("E2E 품질 검증 종합 결과")
    print("=" * 60)

    all_issues = []
    for name, result in results.items():
        if "error" in result:
            status = "FAIL"
            detail = result["error"]
            all_issues.append(f"{name}: {detail}")
        else:
            issue_count = len(result.get("issues", []))
            errors = [i for i in result.get("issues", []) if "[ERROR]" in i]
            warnings = [i for i in result.get("issues", []) if "[WARNING]" in i]
            if errors:
                status = "FAIL"
            elif warnings:
                status = "WARN"
            else:
                status = "PASS"
            detail = f"{issue_count} issues ({len(errors)} errors, {len(warnings)} warnings)"
            all_issues.extend(result.get("issues", []))

        print(f"  [{status}] {name}: {detail}")

    print(f"\n  총 소요 시간: {total_time}s")
    print(f"  총 이슈: {len(all_issues)}건")

    total_errors = sum(1 for i in all_issues if "[ERROR]" in i)
    if total_errors > 0:
        print(f"\n  [FAIL] {total_errors}개 ERROR 발견 — 배포 불가")
        return 1
    else:
        print(f"\n  [OK] ERROR 없음 — 배포 가능")
        return 0


if __name__ == "__main__":
    sys.exit(main())
