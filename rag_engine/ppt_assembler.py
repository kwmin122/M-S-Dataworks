"""PPT Assembler — KRDS 기반 공공기관 표준 PPTX 조립.

대한민국 정부 디자인 시스템(KRDS) 원칙에 기반한 프레젠테이션 생성.
디자인 가이드: docs/plans/ppt-template-guide.md
원본 PDF: data/templates/공공기관_PPT_디자인가이드_KRDS.pdf
"""
from __future__ import annotations

import logging
import os
from datetime import datetime, timezone, timedelta
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from phase2_models import SlideType, SlideContent, QnaPair

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# KRDS 디자인 토큰 — 공공기관 PPT 디자인 가이드 (2026.02)
# ---------------------------------------------------------------------------

# 색상 시스템 (최대 3색 원칙)
KRDS_COLORS = {
    # Primary — 정부 청색 (Government Blue)
    "primary":      RGBColor(0x00, 0x37, 0x64),  # Blue 900 — 제목 바, 표지 배경
    "secondary":    RGBColor(0x00, 0x4A, 0x8F),  # Blue 800 — 보조 강조
    "accent":       RGBColor(0x00, 0x70, 0xE0),  # Blue 600 — 링크, 강조 텍스트
    "blue_500":     RGBColor(0x1A, 0x85, 0xFF),  # Blue 500 — 아이콘, 불릿 번호
    "blue_400":     RGBColor(0x4D, 0xA3, 0xFF),  # Blue 400 — 차트 보조색
    "bg_card":      RGBColor(0xE6, 0xF0, 0xFF),  # Blue 100 — 카드/소제목 배경
    "bg_light":     RGBColor(0xF0, 0xF6, 0xFF),  # Blue 50  — 가장 연한 배경
    # Neutral — 무채색
    "text_dark":    RGBColor(0x1A, 0x1A, 0x1A),  # Gray 900 — 제목 텍스트
    "text_body":    RGBColor(0x44, 0x44, 0x44),  # Gray 700 — 본문 텍스트
    "text_caption":  RGBColor(0x88, 0x88, 0x88),  # Gray 500 — 캡션, 보조
    "text_light":   RGBColor(0xFF, 0xFF, 0xFF),  # White — 반전 텍스트
    "border":       RGBColor(0xE5, 0xE5, 0xE5),  # Gray 200 — 구분선, 테두리
    # Semantic — 의미 색상
    "success":      RGBColor(0x0D, 0x80, 0x50),  # Green — 성공/증가
    "danger":       RGBColor(0xE8, 0x40, 0x2D),  # Red — 경고/감소
    "warning":      RGBColor(0xE0, 0x7B, 0x54),  # Orange — 주의/진행중
    "info":         RGBColor(0x0D, 0x6E, 0x6E),  # Teal — 보조 강조
}

# 하위호환: 기존 colors dict 키 매핑
DEFAULT_COLORS = KRDS_COLORS

# 폰트 시스템 (Pretendard 1순위, 시스템에 없으면 fallback)
KRDS_FONT = "Pretendard"
KRDS_FONT_FALLBACK = "맑은 고딕"

# 타입 스케일 (pt 단위)
KRDS_TYPE = {
    "cover_title":  {"size": 40, "bold": True},     # ExtraBold
    "slide_title":  {"size": 26, "bold": True},      # Bold
    "subtitle":     {"size": 20, "bold": True},       # SemiBold
    "body":         {"size": 15, "bold": False},      # Regular
    "bullet":       {"size": 14, "bold": False},      # Regular
    "caption":      {"size": 11, "bold": False},      # Light
    "page_number":  {"size": 10, "bold": False},      # Light
}

# 레이아웃 수치 (16:9, 1920x1080 기준 → Inches 변환)
KRDS_LAYOUT = {
    "slide_width":      Inches(13.333),   # 16:9 표준
    "slide_height":     Inches(7.5),
    "title_bar_h":      Inches(0.67),     # 48px
    "footer_h":         Inches(0.39),     # 28px
    "margin":           Inches(0.44),     # 32px
    "margin_inner":     Inches(0.33),     # 24px
    "content_gap":      Inches(0.22),     # 16px
    "accent_bar_w":     Inches(0.56),     # 40px (구분선 바 두께)
    "divider_bar_w":    Inches(1.8),      # 간지 좌측 바 폭
}


# ---------------------------------------------------------------------------
# 유틸리티 함수
# ---------------------------------------------------------------------------

def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int = 15,
    bold: bool = False,
    color: RGBColor | None = None,
    alignment: int = PP_ALIGN.LEFT,
    font_name: str = KRDS_FONT,
) -> None:
    """슬라이드에 텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if color:
        p.font.color.rgb = color
    p.alignment = alignment


def _add_rect(slide, left, top, width, height, fill_color: RGBColor) -> None:
    """배경 사각형 추가 (테두리 없음)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()


def _add_title_bar(slide, title: str, colors: dict, page_num: str = "") -> None:
    """상단 제목 바 (KRDS: Blue 900 배경, 흰색 텍스트, 48px 높이)."""
    sw = KRDS_LAYOUT["slide_width"]
    h = KRDS_LAYOUT["title_bar_h"]
    m = KRDS_LAYOUT["margin"]

    _add_rect(slide, Inches(0), Inches(0), sw, h, colors["primary"])

    _add_textbox(
        slide, m, Inches(0.12),
        sw - m * 2, Inches(0.5),
        title,
        font_size=KRDS_TYPE["slide_title"]["size"],
        bold=True,
        color=colors["text_light"],
    )

    # 페이지 번호 (우측 상단)
    if page_num:
        _add_textbox(
            slide,
            sw - Inches(1.2), Inches(0.15),
            Inches(0.8), Inches(0.4),
            page_num,
            font_size=KRDS_TYPE["page_number"]["size"],
            color=colors["text_light"],
            alignment=PP_ALIGN.RIGHT,
        )


def _add_footer(slide, org_name: str, colors: dict) -> None:
    """하단 푸터 (기관 로고/명칭 + CONFIDENTIAL)."""
    sw = KRDS_LAYOUT["slide_width"]
    sh = KRDS_LAYOUT["slide_height"]
    fh = KRDS_LAYOUT["footer_h"]
    m = KRDS_LAYOUT["margin"]
    top = sh - fh

    # 푸터 배경선 (구분선)
    _add_rect(slide, Inches(0), top - Inches(0.02), sw, Inches(0.02), colors["border"])

    if org_name:
        _add_textbox(
            slide, m, top,
            Inches(4), fh,
            org_name,
            font_size=KRDS_TYPE["caption"]["size"],
            color=colors["text_caption"],
        )

    _add_textbox(
        slide,
        sw - Inches(2.5), top,
        Inches(2), fh,
        "CONFIDENTIAL",
        font_size=KRDS_TYPE["caption"]["size"],
        color=colors["text_caption"],
        alignment=PP_ALIGN.RIGHT,
    )


def _add_speaker_notes(slide, notes: str) -> None:
    """발표 노트 추가."""
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


def _add_accent_bar(slide, left, top, width, height, colors: dict) -> None:
    """소제목 앞 강조 바 (Blue 500 컬러 바)."""
    _add_rect(slide, left, top, width, height, colors["blue_500"])


# ---------------------------------------------------------------------------
# 슬라이드 타입별 렌더러 (KRDS 6종 + 기존 호환)
# ---------------------------------------------------------------------------

def _add_cover_slide(
    prs: Presentation,
    content: SlideContent,
    company_name: str,
    colors: dict,
) -> None:
    """A. 표지 슬라이드 — Blue 900 전면 배경, 중앙 제목."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    sw = KRDS_LAYOUT["slide_width"]
    sh = KRDS_LAYOUT["slide_height"]

    # 전면 배경 (Blue 900)
    _add_rect(slide, Inches(0), Inches(0), sw, sh, colors["primary"])

    # 상단 장식 (약간 더 밝은 파란 그라데이션 효과용 사각형)
    _add_rect(
        slide, Inches(0), Inches(0),
        sw, Inches(1.5),
        colors["secondary"],
    )

    # 중앙 제목
    _add_textbox(
        slide,
        Inches(1.5), Inches(2.5),
        sw - Inches(3), Inches(1.8),
        content.title,
        font_size=KRDS_TYPE["cover_title"]["size"],
        bold=True,
        color=colors["text_light"],
        alignment=PP_ALIGN.CENTER,
    )

    # 기관/회사명
    if company_name:
        _add_textbox(
            slide,
            Inches(1.5), Inches(4.5),
            sw - Inches(3), Inches(0.6),
            company_name,
            font_size=KRDS_TYPE["subtitle"]["size"],
            color=colors.get("bg_card", RGBColor(0xE6, 0xF0, 0xFF)),
            alignment=PP_ALIGN.CENTER,
        )

    # 구분선 (accent bar)
    bar_w = Inches(0.8)
    _add_rect(
        slide,
        (sw - bar_w) // 2, Inches(4.2),
        bar_w, Inches(0.04),
        colors["accent"],
    )

    # 날짜 (KST)
    kst = timezone(timedelta(hours=9))
    date_str = datetime.now(kst).strftime("%Y. %m")
    _add_textbox(
        slide,
        Inches(1.5), Inches(5.3),
        sw - Inches(3), Inches(0.5),
        date_str,
        font_size=KRDS_TYPE["caption"]["size"],
        color=colors.get("bg_card", RGBColor(0xE6, 0xF0, 0xFF)),
        alignment=PP_ALIGN.CENTER,
    )

    _add_speaker_notes(slide, content.speaker_notes)


def _add_toc_slide(
    prs: Presentation,
    content: SlideContent,
    colors: dict,
    company_name: str = "",
) -> None:
    """B. 목차 슬라이드 — 상단 제목 바 + 번호 목록."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = KRDS_LAYOUT["slide_width"]
    m = KRDS_LAYOUT["margin"]

    _add_title_bar(slide, "목차", colors)

    y = Inches(1.2)
    items = content.bullets if content.bullets else []

    for i, item in enumerate(items, 1):
        # 번호 (Blue 600 강조)
        _add_textbox(
            slide, m, y,
            Inches(0.6), Inches(0.5),
            f"{i:02d}",
            font_size=KRDS_TYPE["subtitle"]["size"],
            bold=True,
            color=colors["accent"],
        )
        # 섹션명
        _add_textbox(
            slide, m + Inches(0.8), y + Inches(0.03),
            Inches(8), Inches(0.5),
            item,
            font_size=KRDS_TYPE["body"]["size"],
            color=colors["text_dark"],
        )
        y += Inches(0.65)

    _add_footer(slide, company_name, colors)
    _add_speaker_notes(slide, content.speaker_notes)


_MAX_BULLETS_PER_SLIDE = 8  # 슬라이드당 최대 불릿 수 (overflow 방지)


def _add_content_slide(
    prs: Presentation,
    content: SlideContent,
    colors: dict,
    page_info: str = "",
    company_name: str = "",
) -> None:
    """C. 콘텐츠 슬라이드 — 상단 제목바 + 소제목 + 불릿.

    불릿이 _MAX_BULLETS_PER_SLIDE 초과 시 추가 슬라이드 자동 생성.
    """
    bullets = content.bullets or []

    # Split bullets into chunks for overflow handling
    chunks = [bullets[i:i + _MAX_BULLETS_PER_SLIDE]
              for i in range(0, max(len(bullets), 1), _MAX_BULLETS_PER_SLIDE)]
    if not chunks:
        chunks = [[]]

    total_parts = len(chunks)

    for chunk_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        sw = KRDS_LAYOUT["slide_width"]
        m = KRDS_LAYOUT["margin"]
        gap = KRDS_LAYOUT["content_gap"]

        # Smart title: "Part N/M" numbering for continuation slides
        title = content.title
        if total_parts > 1 and chunk_idx > 0:
            title = f"{content.title} — Part {chunk_idx + 1}/{total_parts}"

        _add_title_bar(slide, title, colors, page_num=page_info)

        y = KRDS_LAYOUT["title_bar_h"] + gap + Inches(0.1)
        footer_y = KRDS_LAYOUT["slide_height"] - KRDS_LAYOUT["footer_h"] - Inches(0.1)

        # 소제목 바 (첫 슬라이드만)
        if content.body and chunk_idx == 0:
            _add_accent_bar(slide, m, y, Inches(0.06), Inches(0.4), colors)

            _add_textbox(
                slide, m + Inches(0.15), y,
                sw - m * 2 - Inches(0.15), Inches(0.5),
                content.body,
                font_size=KRDS_TYPE["subtitle"]["size"],
                bold=True,
                color=colors["text_dark"],
            )
            y += Inches(0.65)

        # 불릿 목록
        for bullet in chunk:
            if y + Inches(0.5) > footer_y:
                break  # 푸터 영역 침범 방지

            _add_textbox(
                slide, m + Inches(0.1), y,
                Inches(0.3), Inches(0.4),
                "\u25AA",
                font_size=KRDS_TYPE["bullet"]["size"],
                color=colors["blue_500"],
            )
            _add_textbox(
                slide, m + Inches(0.4), y,
                sw - m * 2 - Inches(0.4), Inches(0.45),
                bullet,
                font_size=KRDS_TYPE["bullet"]["size"],
                color=colors["text_body"],
            )
            y += Inches(0.5)

        _add_footer(slide, company_name, colors)

        # Speaker notes: include slide count info for multi-part slides
        notes = content.speaker_notes or ""
        if total_parts > 1:
            part_note = f"[슬라이드 {chunk_idx + 1}/{total_parts}] "
            notes = part_note + notes if notes else part_note.strip()
        _add_speaker_notes(slide, notes)


def _add_bullet_slide(
    prs: Presentation,
    content: SlideContent,
    colors: dict,
    company_name: str = "",
) -> None:
    """불릿 목록 슬라이드 (콘텐츠와 동일 레이아웃)."""
    _add_content_slide(prs, content, colors, company_name=company_name)


def _add_table_slide(
    prs: Presentation,
    content: SlideContent,
    colors: dict,
    company_name: str = "",
) -> None:
    """표 슬라이드 — KRDS 표 규칙: 헤더 Blue 900, 교대행, 최소 테두리."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = KRDS_LAYOUT["slide_width"]
    m = KRDS_LAYOUT["margin"]

    _add_title_bar(slide, content.title, colors)

    if content.table_data and len(content.table_data) > 0:
        rows = len(content.table_data)
        cols = max(len(row) for row in content.table_data)
        table_w = sw - m * 2
        row_h = Inches(0.45)

        table_shape = slide.shapes.add_table(
            rows, cols,
            m, KRDS_LAYOUT["title_bar_h"] + Inches(0.3),
            table_w, row_h * rows,
        )
        table = table_shape.table

        for r_idx, row in enumerate(content.table_data):
            for c_idx, cell_text in enumerate(row):
                if c_idx < cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(cell_text)

                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(KRDS_TYPE["bullet"]["size"])
                        paragraph.font.name = KRDS_FONT

                    # 헤더 행: Blue 900 배경 + 흰색 텍스트
                    if r_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = colors["primary"]
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.color.rgb = colors["text_light"]
                            paragraph.font.bold = True
                    # 교대 행: Blue 50 배경
                    elif r_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = colors["bg_light"]
    else:
        _add_content_slide(prs, content, colors, company_name=company_name)
        return

    _add_footer(slide, company_name, colors)
    _add_speaker_notes(slide, content.speaker_notes)


def _add_timeline_slide(
    prs: Presentation,
    content: SlideContent,
    colors: dict,
    company_name: str = "",
) -> None:
    """일정 슬라이드 — 간트차트 이미지 삽입 또는 텍스트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = KRDS_LAYOUT["slide_width"]
    m = KRDS_LAYOUT["margin"]

    _add_title_bar(slide, content.title, colors)

    content_top = KRDS_LAYOUT["title_bar_h"] + Inches(0.3)

    if content.image_path and os.path.isfile(content.image_path):
        slide.shapes.add_picture(
            content.image_path,
            m, content_top,
            width=sw - m * 2,
        )
    elif content.body:
        _add_textbox(
            slide, m, content_top,
            sw - m * 2, Inches(5),
            content.body,
            font_size=KRDS_TYPE["body"]["size"],
            color=colors["text_body"],
        )

    _add_footer(slide, company_name, colors)
    _add_speaker_notes(slide, content.speaker_notes)


def _add_team_slide(
    prs: Presentation,
    content: SlideContent,
    colors: dict,
    company_name: str = "",
) -> None:
    """투입인력 슬라이드."""
    _add_content_slide(prs, content, colors, company_name=company_name)


def _add_qna_slide(
    prs: Presentation,
    content: SlideContent,
    qna_pairs: list[QnaPair],
    colors: dict,
    company_name: str = "",
) -> None:
    """Q&A 슬라이드 — Blue 800 전면 배경."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = KRDS_LAYOUT["slide_width"]
    sh = KRDS_LAYOUT["slide_height"]

    _add_rect(slide, Inches(0), Inches(0), sw, sh, colors["secondary"])

    _add_textbox(
        slide,
        Inches(1.5), Inches(2.5),
        sw - Inches(3), Inches(1.5),
        "Q & A",
        font_size=48,
        bold=True,
        color=colors["text_light"],
        alignment=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide,
        Inches(1.5), Inches(4.2),
        sw - Inches(3), Inches(0.8),
        "질의응답",
        font_size=KRDS_TYPE["subtitle"]["size"],
        color=colors.get("bg_card", RGBColor(0xE6, 0xF0, 0xFF)),
        alignment=PP_ALIGN.CENTER,
    )

    # 발표 노트에 전체 QnA 포함
    notes = "예상질문 및 모범답변:\n\n"
    for i, qna in enumerate(qna_pairs, 1):
        notes += f"Q{i}. {qna.question}\n"
        notes += f"A{i}. {qna.answer}\n\n"

    _add_speaker_notes(slide, notes)


def _add_closing_slide(
    prs: Presentation,
    content: SlideContent,
    company_name: str,
    colors: dict,
) -> None:
    """E. 마무리/감사 슬라이드 — Blue 900 전면 배경."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = KRDS_LAYOUT["slide_width"]
    sh = KRDS_LAYOUT["slide_height"]

    _add_rect(slide, Inches(0), Inches(0), sw, sh, colors["primary"])

    _add_textbox(
        slide,
        Inches(1.5), Inches(2.5),
        sw - Inches(3), Inches(1.5),
        "감사합니다",
        font_size=44,
        bold=True,
        color=colors["text_light"],
        alignment=PP_ALIGN.CENTER,
    )

    # 구분선
    bar_w = Inches(0.8)
    _add_rect(
        slide,
        (sw - bar_w) // 2, Inches(4.2),
        bar_w, Inches(0.04),
        colors["accent"],
    )

    if company_name:
        _add_textbox(
            slide,
            Inches(1.5), Inches(4.6),
            sw - Inches(3), Inches(0.8),
            company_name,
            font_size=KRDS_TYPE["body"]["size"],
            color=colors.get("bg_card", RGBColor(0xE6, 0xF0, 0xFF)),
            alignment=PP_ALIGN.CENTER,
        )

    _add_speaker_notes(slide, content.speaker_notes)


def _add_divider_slide(
    prs: Presentation,
    content: SlideContent,
    colors: dict,
    company_name: str = "",
) -> None:
    """F. 간지(구분) 슬라이드 — 좌측 Blue 900 바 + 섹션 번호 / 우측 제목+설명."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = KRDS_LAYOUT["slide_width"]
    sh = KRDS_LAYOUT["slide_height"]
    bar_w = KRDS_LAYOUT["divider_bar_w"]

    # 좌측 컬러 바 (Blue 900)
    _add_rect(slide, Inches(0), Inches(0), bar_w, sh, colors["primary"])

    # 섹션 번호 추출 (title에서 "01", "02" 등)
    section_num = ""
    title_text = content.title
    if title_text and len(title_text) >= 2 and title_text[:2].isdigit():
        section_num = title_text[:2]
        title_text = title_text[2:].strip(". ").strip()

    if section_num:
        _add_textbox(
            slide,
            Inches(0.3), Inches(2.5),
            bar_w - Inches(0.6), Inches(1.2),
            section_num,
            font_size=48,
            bold=True,
            color=colors["text_light"],
            alignment=PP_ALIGN.CENTER,
        )

    # 우측: 섹션 제목
    right_left = bar_w + Inches(0.8)
    _add_textbox(
        slide,
        right_left, Inches(2.5),
        sw - right_left - Inches(1), Inches(1.0),
        title_text or content.title,
        font_size=KRDS_TYPE["slide_title"]["size"],
        bold=True,
        color=colors["text_dark"],
    )

    # 우측: 설명 텍스트
    if content.body:
        _add_textbox(
            slide,
            right_left, Inches(3.6),
            sw - right_left - Inches(1), Inches(1.5),
            content.body,
            font_size=KRDS_TYPE["body"]["size"],
            color=colors["text_body"],
        )

    _add_footer(slide, company_name, colors)
    _add_speaker_notes(slide, content.speaker_notes)


# ---------------------------------------------------------------------------
# Main assembler
# ---------------------------------------------------------------------------

def assemble_pptx(
    title: str,
    slides: list[SlideContent],
    qna_pairs: list[QnaPair],
    output_path: str,
    company_name: str = "",
    colors: Optional[dict] = None,
) -> str:
    """KRDS 기반 PPTX 조립.

    Args:
        title: 전체 발표 제목
        slides: 슬라이드 콘텐츠 목록
        qna_pairs: 예상질문+답변 (Q&A 슬라이드 발표노트에 삽입)
        output_path: 출력 파일 경로
        company_name: 기관/회사명 (표지, 마무리, 푸터에 표시)
        colors: 커스텀 색상 (None이면 KRDS_COLORS 사용)
    """
    prs = Presentation()
    prs.slide_width = KRDS_LAYOUT["slide_width"]
    prs.slide_height = KRDS_LAYOUT["slide_height"]

    c = colors or KRDS_COLORS

    # 페이지 번호 계산 (표지/마무리 제외)
    total_content = sum(
        1 for s in slides
        if s.slide_type not in (SlideType.COVER, SlideType.CLOSING, SlideType.QNA)
    )
    content_idx = 0

    for slide_content in slides:
        st = slide_content.slide_type

        # 페이지 번호 문자열
        page_str = ""
        if st not in (SlideType.COVER, SlideType.CLOSING, SlideType.QNA, SlideType.DIVIDER):
            content_idx += 1
            page_str = f"{content_idx:02d} / {total_content:02d}"

        if st == SlideType.COVER:
            _add_cover_slide(prs, slide_content, company_name, c)
        elif st == SlideType.TOC:
            _add_toc_slide(prs, slide_content, c, company_name=company_name)
        elif st == SlideType.DIVIDER:
            _add_divider_slide(prs, slide_content, c, company_name=company_name)
        elif st == SlideType.QNA:
            _add_qna_slide(prs, slide_content, qna_pairs, c, company_name=company_name)
        elif st == SlideType.CLOSING:
            _add_closing_slide(prs, slide_content, company_name, c)
        elif st == SlideType.TIMELINE:
            _add_timeline_slide(prs, slide_content, c, company_name=company_name)
        elif st == SlideType.TABLE:
            _add_table_slide(prs, slide_content, c, company_name=company_name)
        elif st == SlideType.TEAM:
            _add_team_slide(prs, slide_content, c, company_name=company_name)
        else:
            _add_content_slide(prs, slide_content, c, page_info=page_str, company_name=company_name)

    prs.save(output_path)
    return output_path
