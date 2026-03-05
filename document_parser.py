"""
RFx AI Assistant - 문서 파서/청커

지원 포맷:
- PDF (.pdf)
- DOCX (.docx)
- TXT (.txt)
"""

from __future__ import annotations

import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional


@dataclass
class ParsedDocument:
    """파싱된 문서"""

    filename: str
    text: str
    pages: list[str] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)

    @property
    def char_count(self) -> int:
        return len(self.text or "")

    @property
    def page_count(self) -> int:
        return len(self.pages) if self.pages else 1


@dataclass
class TextChunk:
    """RAG 저장 단위 텍스트 청크"""

    text: str
    chunk_id: int
    source_file: str
    page_number: Optional[int] = None
    metadata: dict = field(default_factory=dict)


class TextChunker:
    """문서 텍스트를 겹침(overlap) 기반으로 분할"""

    def __init__(self, chunk_size: int = 900, chunk_overlap: int = 150):
        if chunk_size <= 0:
            raise ValueError("chunk_size must be > 0")
        if chunk_overlap < 0:
            raise ValueError("chunk_overlap must be >= 0")
        if chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap must be < chunk_size")
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    @staticmethod
    def _normalize_text(text: str) -> str:
        # 일부 PDF 추출기(pdfplumber)에서 NUL 문자가 섞여 들어오는 케이스를 정리.
        text = text.replace("\x00", "")
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def _slice_text(self, text: str) -> list[str]:
        normalized = self._normalize_text(text)
        if not normalized:
            return []
        if len(normalized) <= self.chunk_size:
            return [normalized]

        chunks: list[str] = []
        start = 0
        text_len = len(normalized)
        while start < text_len:
            end = min(start + self.chunk_size, text_len)
            candidate = normalized[start:end]

            # 문장 경계 쪽으로 컷 위치를 맞춰 가독성을 높인다.
            if end < text_len:
                boundary = max(
                    candidate.rfind(". "),
                    candidate.rfind("! "),
                    candidate.rfind("? "),
                    candidate.rfind("\n"),
                )
                if boundary > int(self.chunk_size * 0.5):
                    end = start + boundary + 1
                    candidate = normalized[start:end]

            candidate = candidate.strip()
            if candidate:
                chunks.append(candidate)

            if end >= text_len:
                break
            start = max(0, end - self.chunk_overlap)

        return chunks

    def chunk_document(self, doc: ParsedDocument) -> list[TextChunk]:
        chunks: list[TextChunk] = []
        chunk_id = 0

        source_name = Path(doc.filename).name or "unknown"

        # 페이지 정보가 있으면 페이지 단위로 먼저 분할
        if doc.pages:
            for idx, page_text in enumerate(doc.pages, start=1):
                for part in self._slice_text(page_text):
                    chunks.append(
                        TextChunk(
                            text=part,
                            chunk_id=chunk_id,
                            source_file=source_name,
                            page_number=idx,
                            metadata={"type": "text"},
                        )
                    )
                    chunk_id += 1
            return chunks

        # 페이지 정보가 없으면 문서 전체 기준으로 분할
        for part in self._slice_text(doc.text):
            chunks.append(
                TextChunk(
                    text=part,
                    chunk_id=chunk_id,
                    source_file=source_name,
                    page_number=None,
                    metadata={"type": "text"},
                )
            )
            chunk_id += 1

        return chunks


class DocumentParser:
    """파일 포맷별 파싱 + 청킹 엔트리포인트"""

    def __init__(self, chunk_size: int = 900, chunk_overlap: int = 150):
        self.chunker = TextChunker(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    def parse(self, file_path: str) -> ParsedDocument:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"문서 파일을 찾을 수 없습니다: {file_path}")

        ext = path.suffix.lower()
        if ext == ".pdf":
            return self._parse_pdf(path)
        if ext == ".docx":
            return self._parse_docx(path)
        if ext == ".txt":
            return self._parse_txt(path)
        if ext == ".hwp":
            return self._parse_hwp(path)
        if ext == ".hwpx":
            return self._parse_hwpx(path)
        if ext in (".xlsx", ".xls"):
            return self._parse_excel(path)
        if ext == ".csv":
            return self._parse_csv(path)
        if ext in (".pptx", ".ppt"):
            return self._parse_pptx(path)

        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

    def parse_and_chunk(self, file_path: str) -> list[TextChunk]:
        parsed = self.parse(file_path)
        return self.chunker.chunk_document(parsed)

    def _parse_hwp(self, path: Path) -> ParsedDocument:
        """HWP 5.x 파일 파싱 (olefile + zlib).

        HWP 5.x 구조:
        - OLE Compound Document (olefile로 열기)
        - BodyText/Section0 스트림: zlib raw deflate 압축
        - 레코드 타입 67 (HWPTAG_PARA_TEXT): UTF-16-LE 텍스트

        가비지 필터: 한글/영문/숫자/공통 기호 외 문자 제거
        """
        try:
            import olefile
        except ImportError as exc:
            raise ImportError("olefile 미설치: pip install olefile") from exc

        import zlib
        import struct

        try:
            ole = olefile.OleFileIO(str(path))
        except Exception as exc:
            raise ValueError(f"HWP 파일을 열 수 없습니다: {exc}") from exc

        if not ole.exists("BodyText/Section0"):
            ole.close()
            raise ValueError("HWP 5.x 형식이 아닙니다 (BodyText/Section0 없음)")

        raw_data = ole.openstream("BodyText/Section0").read()
        ole.close()

        try:
            decompressed = zlib.decompress(raw_data, -15)
        except zlib.error as exc:
            raise ValueError(f"HWP 압축 해제 실패: {exc}") from exc

        # 레코드 파싱: type 67 = HWPTAG_PARA_TEXT
        paragraphs: list[str] = []
        i = 0
        while i < len(decompressed) - 4:
            tag_header = struct.unpack_from("<I", decompressed, i)[0]
            record_type = tag_header & 0x3FF
            size = (tag_header >> 20) & 0xFFF
            if size == 0xFFF:
                if i + 8 > len(decompressed):
                    break
                size = struct.unpack_from("<I", decompressed, i + 4)[0]
                i += 8
            else:
                i += 4
            if i + size > len(decompressed):
                break
            if record_type == 67:  # HWPTAG_PARA_TEXT
                chunk = decompressed[i : i + size]
                try:
                    text = chunk.decode("utf-16-le", errors="ignore").strip()
                except Exception:
                    text = ""
                if text and len(text) > 1:
                    # 한글/영문/숫자/공통 기호만 유지, 가비지 제거
                    cleaned = re.sub(r"[^\uAC00-\uD7A3\u3130-\u318F\w\s\-.,()[\]{}:;/\\@#%&*+=<>?!\n]", "", text).strip()
                    if cleaned:
                        paragraphs.append(cleaned)
            i += size

        text = "\n".join(paragraphs)
        normalized = self.chunker._normalize_text(text)
        pages = [normalized] if normalized else []
        return ParsedDocument(
            filename=path.name,
            text=normalized,
            pages=pages,
            metadata={"file_type": "hwp"},
        )

    def _parse_hwpx(self, path: Path) -> ParsedDocument:
        """HWPX 파일 파싱 (OOXML 기반, zipfile + xml)."""
        import zipfile
        from xml.etree import ElementTree as ET

        if not zipfile.is_zipfile(str(path)):
            raise ValueError("유효한 HWPX 파일이 아닙니다.")

        paragraphs: list[str] = []
        with zipfile.ZipFile(str(path), "r") as zf:
            # Contents/section*.xml 에 본문이 있음
            section_files = sorted(
                [n for n in zf.namelist() if n.startswith("Contents/section") and n.endswith(".xml")]
            )
            if not section_files:
                # 대안: Contents/ 내 모든 xml
                section_files = sorted(
                    [n for n in zf.namelist() if n.startswith("Contents/") and n.endswith(".xml")]
                )

            for section_file in section_files:
                try:
                    xml_data = zf.read(section_file)
                    root = ET.fromstring(xml_data)
                    # 모든 텍스트 노드 추출
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            paragraphs.append(elem.text.strip())
                        if elem.tail and elem.tail.strip():
                            paragraphs.append(elem.tail.strip())
                except Exception:
                    continue

        text = "\n".join(paragraphs)
        normalized = self.chunker._normalize_text(text)
        pages = [normalized] if normalized else []
        return ParsedDocument(
            filename=path.name,
            text=normalized,
            pages=pages,
            metadata={"file_type": "hwpx"},
        )

    def _parse_excel(self, path: Path) -> ParsedDocument:
        """Excel (.xlsx, .xls) 파싱 (openpyxl)."""
        try:
            import openpyxl
        except ImportError as exc:
            raise ImportError("openpyxl 미설치: pip install openpyxl") from exc

        if path.suffix.lower() == ".xls":
            raise ValueError(
                "구형 Excel(.xls) 형식은 지원하지 않습니다. "
                ".xlsx로 변환 후 업로드해주세요."
            )

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets_text: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell or "").strip() for cell in row]
                line = "\t".join(cells).strip()
                if line.replace("\t", ""):
                    rows.append(line)
            if rows:
                sheets_text.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))

        wb.close()
        text = "\n\n".join(sheets_text)
        normalized = self.chunker._normalize_text(text)
        pages = [normalized] if normalized else []
        return ParsedDocument(
            filename=path.name,
            text=normalized,
            pages=pages,
            metadata={"file_type": "excel"},
        )

    def _parse_csv(self, path: Path) -> ParsedDocument:
        """CSV 파일 파싱."""
        import csv

        rows: list[str] = []
        for encoding in ("utf-8", "cp949", "euc-kr"):
            try:
                with open(str(path), "r", encoding=encoding, errors="strict") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        line = "\t".join(cell.strip() for cell in row).strip()
                        if line.replace("\t", ""):
                            rows.append(line)
                break
            except (UnicodeDecodeError, UnicodeError):
                rows = []
                continue

        text = "\n".join(rows)
        normalized = self.chunker._normalize_text(text)
        pages = [normalized] if normalized else []
        return ParsedDocument(
            filename=path.name,
            text=normalized,
            pages=pages,
            metadata={"file_type": "csv"},
        )

    def _parse_pptx(self, path: Path) -> ParsedDocument:
        """PowerPoint (.pptx) 파싱 (python-pptx)."""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError("python-pptx 미설치: pip install python-pptx") from exc

        if path.suffix.lower() == ".ppt":
            raise ValueError(
                "구형 PowerPoint(.ppt) 형식은 지원하지 않습니다. "
                ".pptx로 변환 후 업로드해주세요."
            )

        prs = Presentation(str(path))
        slides_text: list[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        line = "\t".join(cells).strip()
                        if line.replace("\t", ""):
                            texts.append(line)
            if texts:
                slides_text.append(f"[슬라이드 {slide_idx}]\n" + "\n".join(texts))

        text = "\n\n".join(slides_text)
        normalized = self.chunker._normalize_text(text)
        # 슬라이드별 페이지 분할
        pages = []
        for slide_text in slides_text:
            norm_slide = self.chunker._normalize_text(slide_text)
            if norm_slide:
                pages.append(norm_slide)
        if not pages and normalized:
            pages = [normalized]

        return ParsedDocument(
            filename=path.name,
            text=normalized,
            pages=pages,
            metadata={"file_type": "pptx"},
        )

    def _parse_txt(self, path: Path) -> ParsedDocument:
        text = path.read_text(encoding="utf-8", errors="ignore")
        normalized = self.chunker._normalize_text(text)
        pages = [normalized] if normalized else []
        return ParsedDocument(
            filename=path.name,
            text=normalized,
            pages=pages,
            metadata={"file_type": "txt"},
        )

    def _parse_docx(self, path: Path) -> ParsedDocument:
        try:
            import docx
        except ImportError as exc:
            raise ImportError("python-docx 미설치. pip install python-docx") from exc

        document = docx.Document(str(path))
        paragraphs = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]
        text = "\n".join(paragraphs)
        normalized = self.chunker._normalize_text(text)
        pages = [normalized] if normalized else []
        return ParsedDocument(
            filename=path.name,
            text=normalized,
            pages=pages,
            metadata={"file_type": "docx"},
        )

    @staticmethod
    def _table_to_markdown(table) -> str:
        """pdfplumber 표 데이터 → GFM 마크다운 테이블.

        Args:
            table: list[list] — pdfplumber extract_tables() 결과

        Returns:
            마크다운 테이블 문자열. 행 1개 이하면 빈 문자열.
        """
        if not table or len(table) < 2:
            return ""

        def _cell(v):
            s = str(v or "").strip()
            return s.replace("|", "\\|").replace("\n", " ")

        header = [_cell(c) for c in table[0]]
        col_count = len(header)
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * col_count) + " |",
        ]

        for row in table[1:]:
            cells = [_cell(c) for c in row]
            while len(cells) < col_count:
                cells.append("")
            lines.append("| " + " | ".join(cells[:col_count]) + " |")

        return "\n".join(lines)

    def _parse_pdf(self, path: Path) -> ParsedDocument:
        pages: list[str] = []

        # 1차: pdfplumber 시도
        try:
            import pdfplumber

            with pdfplumber.open(str(path)) as pdf:
                for page in pdf.pages:
                    parts: list[str] = []

                    # 텍스트 추출
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        parts.append(page_text)

                    # 표 추출 → 마크다운 변환
                    try:
                        tables = page.extract_tables() or []
                        for table in tables:
                            md = self._table_to_markdown(table)
                            if md:
                                parts.append(md)
                    except Exception:
                        pass  # 표 추출 실패 시 텍스트만 사용

                    combined = "\n\n".join(parts)
                    normalized = self.chunker._normalize_text(combined)
                    pages.append(normalized)
        except ImportError:
            pass
        except Exception:
            # 2차 파서로 폴백
            pages = []

        # 2차: PyMuPDF 폴백
        if not any(pages):
            try:
                import fitz
            except ImportError as exc:
                raise ImportError("PDF 파서 미설치. pip install pdfplumber pymupdf") from exc

            doc = fitz.open(str(path))
            try:
                for page in doc:
                    page_text = page.get_text("text") or ""
                    normalized = self.chunker._normalize_text(page_text)
                    pages.append(normalized)
            finally:
                doc.close()

        # 빈 페이지 제거
        pages = [p for p in pages if p]
        full_text = "\n\n".join(pages).strip()

        return ParsedDocument(
            filename=path.name,
            text=full_text,
            pages=pages,
            metadata={"file_type": "pdf"},
        )
